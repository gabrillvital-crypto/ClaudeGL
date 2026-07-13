#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

# Template path
template_path = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\Documentos\Plano de Ação Efcaz _ DockBrasil - modelo.pptx"

# Output path
output_dir = Path(r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\clientes\NPS")
output_dir.mkdir(parents=True, exist_ok=True)
output_path = output_dir / "NPS_Carteira_Efcaz_2026_06_19.pptx"

# Load template
print(f"Carregando template: {template_path}")
prs = Presentation(template_path)

print(f"\nTemplate carregado com {len(prs.slides)} slides\n")

# Inspecionar template
print("=== INSPEÇÃO DO TEMPLATE ===\n")
for i, slide in enumerate(prs.slides):
    print(f"--- Slide {i+1} ---")
    for shape in slide.shapes:
        if shape.has_text_frame:
            texto = "".join(run.text for para in shape.text_frame.paragraphs for run in para.runs).strip()
            if texto:
                print(f"  {shape.name}: {texto[:100]}")

# Dados para substituição
substituicoes_gerais = {
    "{{CLIENTE}}": "NPS Carteira Efcaz",
    "{{DATA}}": "Junho 2026",
    "{{MES_ANO}}": "Junho/2026",
    "{{CS}}": "Gabriel Vital",
}

# Função para substituir textos preservando formatação
def corrigir_runs_fragmentados(text_frame, substituicoes):
    for paragraph in text_frame.paragraphs:
        # Concatenar texto completo
        texto_completo = "".join(run.text for run in paragraph.runs)

        # Verificar se há algum placeholder
        alterado = False
        for placeholder, valor in substituicoes.items():
            if placeholder in texto_completo:
                texto_completo = texto_completo.replace(placeholder, valor)
                alterado = True

        # Reescrever no primeiro run
        if alterado and paragraph.runs:
            paragraph.runs[0].text = texto_completo
            for run in paragraph.runs[1:]:
                run.text = ""

def substituir_texto(prs, substituicoes):
    """Substitui textos no template preservando formatação"""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                corrigir_runs_fragmentados(shape.text_frame, substituicoes)

# Aplicar substituições gerais
print("\n=== APLICANDO SUBSTITUIÇÕES ===\n")
substituir_texto(prs, substituicoes_gerais)

# Dados específicos de NPS para inserir manualmente nos slides
nps_data = {
    "respondidos": "15 clientes (60%)",
    "sem_resposta": "8 clientes (32%)",
    "sem_acesso": "2 clientes (8%)",
    "nps_medio": "+73",
    "taxa_resposta": "22.6%",
    "total_acessos": "164",
    "total_respondentes": "37",
    "detratores": "0",
    "promoters": "100%",
}

# Adaptar slides com dados de NPS
# Se há slides vazios, preencher com conteúdo de NPS

if len(prs.slides) >= 2:
    # Slide 2: Resumo Executivo
    slide = prs.slides[1]
    print("Preenchendo Slide 2 - Resumo Executivo...")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                texto = "".join(run.text for run in para.runs).strip().lower()

                if "respondido" in texto or "resposta" in texto:
                    para.runs[0].text = f"Respondidos: {nps_data['respondidos']}"
                elif "sem resposta" in texto or "não responderam" in texto:
                    para.runs[0].text = f"Sem Resposta: {nps_data['sem_resposta']}"
                elif "sem acesso" in texto:
                    para.runs[0].text = f"Sem Acesso: {nps_data['sem_acesso']}"

if len(prs.slides) >= 3:
    # Slide 3: Distribuição
    slide = prs.slides[2]
    print("Preenchendo Slide 3 - Distribuição por Status...")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                texto = "".join(run.text for run in para.runs).strip().lower()

                if "problema" in texto or "insight" in texto or "taxa" in texto:
                    para.runs[0].text = f"Taxa Geral: {nps_data['taxa_resposta']} | {nps_data['total_respondentes']} de {nps_data['total_acessos']} usuários responderam"

if len(prs.slides) >= 4:
    # Slide 4: Qualidade
    slide = prs.slides[3]
    print("Preenchendo Slide 4 - Qualidade dos Respondentes...")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                texto = "".join(run.text for run in para.runs).strip().lower()

                if "nps" in texto:
                    para.runs[0].text = f"NPS Médio: {nps_data['nps_medio']}"
                elif "detrator" in texto:
                    para.runs[0].text = f"Detratores Genuínos: {nps_data['detratores']}"
                elif "promoter" in texto:
                    para.runs[0].text = f"Promoters entre Respondentes: {nps_data['promoters']}"

# Salvar arquivo
print(f"\n=== SALVANDO ===\n")
prs.save(str(output_path))
print(f"✅ Arquivo salvo em: {output_path}")
print(f"   Abra em PowerPoint para revisar")
