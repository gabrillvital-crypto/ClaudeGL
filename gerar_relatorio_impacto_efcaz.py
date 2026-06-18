# -*- coding: utf-8 -*-
"""Slide executivo — Atualizacao de Base e Governanca de Dados (Zurich Airport).
Gerado do zero (layout de 2 colunas nao existe no template Dock Brasil).

Sempre que o dash Zurich for atualizado com nova base de CSVs, reexecutar este
script com os parametros abaixo atualizados (datas, linhas da tabela e lista de
CNPJs descontinuados) para gerar uma nova copia datada, destinada ao envio para
a Zurich. Nunca sobrescrever copias anteriores — cada atualizacao gera um novo
arquivo com a data de referencia no nome e no titulo do slide.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── PARAMETROS DA ATUALIZACAO (editar a cada nova base) ─────────────────────
DATA_REF_ANTERIOR = "12/06/2026"
DATA_REF_NOVA     = "16/06/2026"
MES_ANO_RODAPE    = "Junho/2026"
TOTAL_FORN_ANTES  = 57
TOTAL_FORN_DEPOIS = 50

# ── CORES — PALETA OFICIAL EFCAZ (extraída do template real / demais decks) ─
AZUL_ESCURO   = RGBColor(0x0E, 0x8F, 0xA3)   # institucional primária — #0E8FA3
AZUL_ESCURO_2 = RGBColor(0x0A, 0x6A, 0x7A)   # institucional, tom mais escuro (rodapé)
TEAL_CLARO    = RGBColor(0x14, 0xB3, 0xCC)   # institucional, tom claro (linha de destaque)
CINZA_CLARO   = RGBColor(0xF2, 0xF4, 0xF4)   # fundo do slide
CINZA_LINHA   = RGBColor(0xE8, 0xEC, 0xEC)   # linha zebrada da tabela
CINZA_MEDIO   = RGBColor(0x85, 0x92, 0x9E)   # texto secundario
BRANCO        = RGBColor(0xFF, 0xFF, 0xFF)
VERDE_ESMER   = RGBColor(0x27, 0xAE, 0x60)   # indicador positivo
VERMELHO_COR  = RGBColor(0xE7, 0x4C, 0x3C)   # alerta / descontinuado
PRETO_TEXTO   = RGBColor(0x1A, 0x2A, 0x2A)

FONTE = "Calibri"

W, H = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

# ── HELPERS ──────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill, border=None, bw=1.0):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s

def tb(slide, text, x, y, w, h, size=10, bold=False, color=PRETO_TEXTO,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = FONTE
    return box

def delta_color(valor):
    if valor > 0:
        return VERMELHO_COR
    if valor < 0:
        return VERDE_ESMER
    return CINZA_MEDIO

# ── SLIDE ────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, CINZA_CLARO)

# Cabecalho institucional
rect(s, 0, 0, W, 1.15, AZUL_ESCURO)
rect(s, 0, 1.08, W, 0.07, TEAL_CLARO)
tb(s, "Atualização de Base e Governança de Dados", 0.5, 0.16, 10.6, 0.5,
   size=22, bold=True, color=BRANCO)
tb(s, f"Zurich Airport  |  Dashboard Efcaz SRM  •  Base de {DATA_REF_NOVA}",
   0.5, 0.70, 10.6, 0.34, size=12, color=RGBColor(0xC9, 0xD6, 0xE4))

logo_path = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\Geral\logo_efcaz_clean.png"
if not os.path.exists(logo_path):
    logo_path = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\Geral\logo_efcaz.png"
if os.path.exists(logo_path):
    s.shapes.add_picture(logo_path, Inches(11.35), Inches(0.30), width=Inches(1.6))

# ── COLUNA 1 (ESQUERDA) — TABELA COMPARATIVA ───────────────────────────────
LX, LW = 0.5, 6.55
tb(s, "Comparativo de Impacto — Antes vs. Depois", LX, 1.35, LW, 0.36,
   size=14, bold=True, color=AZUL_ESCURO)
tb(s, f"8 indicadores ativos do dashboard  •  base anterior ({DATA_REF_ANTERIOR}) vs. nova base ({DATA_REF_NOVA})",
   LX, 1.68, LW, 0.28, size=9.5, color=CINZA_MEDIO, italic=True)

ROWS = [
    ("Docs Esperados Fornecedor (R4)", 1149, 1149),
    ("Docs Esperados Terceiros (R3)",  5475, 5480),
    ("Documentos Aprovados",           1343, 1341),
    ("Documentos Não Aprovados",       2275, 2363),
    ("Documentos Não Enviados",        1526, 1518),
    ("Aguardando Submissão",           1279, 1209),
    ("Documentos Em Análise",            58,   50),
    ("Documentos Vencidos",             143,  148),
]

COLW = [3.05, 1.15, 1.15, 1.20]
COLX = [LX, LX+COLW[0], LX+COLW[0]+COLW[1], LX+COLW[0]+COLW[1]+COLW[2]]
HEAD_Y = 2.05
ROW_H = 0.46

rect(s, LX, HEAD_Y, LW, 0.4, AZUL_ESCURO)
for lbl, cx, cw, al in [
    ("Métrica", COLX[0], COLW[0], PP_ALIGN.LEFT),
    ("Base Anterior", COLX[1], COLW[1], PP_ALIGN.CENTER),
    ("Nova Base", COLX[2], COLW[2], PP_ALIGN.CENTER),
    ("Δ", COLX[3], COLW[3], PP_ALIGN.CENTER),
]:
    tb(s, lbl, cx+0.08, HEAD_Y+0.06, cw-0.1, 0.3, size=10, bold=True,
       color=BRANCO, align=al)

for i, (nome, antes, depois) in enumerate(ROWS):
    y = HEAD_Y + 0.4 + i*ROW_H
    bg = BRANCO if i % 2 == 0 else CINZA_LINHA
    rect(s, LX, y, LW, ROW_H, bg)
    delta = depois - antes
    dcolor = delta_color(delta)
    tb(s, nome, COLX[0]+0.1, y+0.08, COLW[0]-0.12, 0.3, size=9.5, color=PRETO_TEXTO)
    tb(s, f"{antes:,}".replace(",", "."), COLX[1], y+0.08, COLW[1]-0.05, 0.3,
       size=9.5, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)
    tb(s, f"{depois:,}".replace(",", "."), COLX[2], y+0.08, COLW[2]-0.05, 0.3,
       size=9.5, bold=True, color=AZUL_ESCURO, align=PP_ALIGN.CENTER)
    sinal = "+" if delta > 0 else ("" if delta == 0 else "-")
    tb(s, f"{sinal}{abs(delta):,}".replace(",", "."), COLX[3], y+0.08, COLW[3]-0.05, 0.3,
       size=9.5, bold=True, color=dcolor, align=PP_ALIGN.CENTER)

table_bottom = HEAD_Y + 0.4 + len(ROWS)*ROW_H
rect(s, LX, table_bottom, LW, 0.02, AZUL_ESCURO)

tb(s, f"Total de Fornecedores na base: {TOTAL_FORN_ANTES} → {TOTAL_FORN_DEPOIS}  (queda associada ao saneamento — ver coluna ao lado)",
   LX, table_bottom + 0.14, LW, 0.32, size=9, color=CINZA_MEDIO, italic=True)

# ── COLUNA 2 (DIREITA) — EXPURGOS E GOVERNANÇA ─────────────────────────────
RX, RW = 7.35, 5.45
BOX_Y = 1.35
BOX_H = 5.55

CNPJS = [
    ("CONSTRUTORA HERA LTDA",                          "10.251.085/0001-34"),
    ("DMW METALMECANICA E CONSTRUTORA LTDA",            "55.054.247/0001-82"),
    ("MHS - SERVIÇOS DE CONSTRUÇÃO E MANUTENÇÃO LTDA",  "44.970.916/0001-40"),
    ("SAMTAL LTDA",                                     "04.827.603/0001-12"),
    ("T & M PINTURAS E PLACAS LTDA",                    "22.412.410/0001-82"),
    ("TK ELEVADORES BRASIL LTDA",                       "90.347.840/0009-75"),
    ("UNIÃO SUL CONTROLE DE PRAGAS LTDA",                "07.817.370/0001-65"),
]

rect(s, RX, BOX_Y, RW, BOX_H, BRANCO, border=VERMELHO_COR, bw=1.4)
rect(s, RX, BOX_Y, RW, 0.62, VERMELHO_COR)
tb(s, "Expurgos e Descontinuações (Limpeza de Base)", RX+0.22, BOX_Y+0.13, RW-0.4, 0.4,
   size=13, bold=True, color=BRANCO)

tb(s, f"{len(CNPJS)} fornecedores removidos do Relatório de Códigos de Contrato — saneamento oficial, não erro de carga.",
   RX+0.25, BOX_Y+0.74, RW-0.5, 0.42, size=9.5, color=CINZA_MEDIO, italic=True)

ITEM_Y0 = BOX_Y + 1.20
ITEM_H = 0.50
for i, (nome, cnpj) in enumerate(CNPJS):
    y = ITEM_Y0 + i*ITEM_H
    rect(s, RX+0.25, y+0.05, 0.1, 0.1, VERMELHO_COR)
    tb(s, cnpj, RX+0.48, y+0.01, RW-0.75, 0.24, size=10.5, bold=True, color=AZUL_ESCURO)
    tb(s, nome, RX+0.48, y+0.25, RW-0.75, 0.24, size=8.8, color=CINZA_MEDIO)

footer_box_y = ITEM_Y0 + len(CNPJS)*ITEM_H + 0.10
rect(s, RX+0.2, footer_box_y, RW-0.4, 0.02, VERMELHO_COR)
tb(s, "Lista preparada para envio ao time de BPO para validação operacional.",
   RX+0.25, footer_box_y + 0.10, RW-0.5, 0.3, size=9, italic=True, color=CINZA_MEDIO)

# ── RODAPÉ ───────────────────────────────────────────────────────────────
rect(s, 0, H-0.42, W, 0.42, AZUL_ESCURO_2)
tb(s, f"Gabriel Vital  •  Customer Success Specialist  •  Efcaz SRM  •  {MES_ANO_RODAPE}",
   0.5, H-0.40, W-1.0, 0.3, size=9, color=RGBColor(0xC9, 0xD6, 0xE4), align=PP_ALIGN.CENTER)

# ── SALVAR — nova copia datada a cada execucao, nunca sobrescreve a anterior ─
# Unico destino: dashboard-react/ (Gabriel pediu para nao duplicar na raiz do projeto)
data_arquivo = DATA_REF_NOVA.replace("/", "-")
nome_arquivo = f"relatorio_impacto_efcaz_{data_arquivo}.pptx"
output = os.path.join(r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\dashboard-react", nome_arquivo)
prs.save(output)
print(f"PPTX salvo: {output}")
