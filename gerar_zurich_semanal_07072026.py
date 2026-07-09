from pptx import Presentation
from pptx.util import Pt
import copy

TEMPLATE = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\dashboard-react\relatorio_impacto_efcaz_30-06-2026_v2.pptx"
OUTPUT   = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\dashboard-react\relatorio_impacto_efcaz_07-07-2026.pptx"

prs = Presentation(TEMPLATE)

# ── Dados do comparativo 03/07 → 07/07 ────────────────────────────────────────
ant = "03/07/2026"
atu = "07/07/2026"

metricas = {
    # shape_name : novo_texto
    # Slide 1 — cabeçalho
    "TextBox 5":  f"Zurich Airport  |  Dashboard Efcaz SRM  •  Base de {atu}",
    "TextBox 8":  f"8 indicadores ativos do dashboard  •  base anterior ({ant}) vs. nova base ({atu})",
    # Docs Esperados R4
    "TextBox 16": "1.478",
    "TextBox 17": "1.491",
    "TextBox 18": "+13",
    # Docs Esperados R3
    "TextBox 21": "8.461",
    "TextBox 22": "8.602",
    "TextBox 23": "+141",
    # Documentos Aprovados
    "TextBox 26": "2.142",
    "TextBox 27": "2.145",
    "TextBox 28": "+3",
    # Documentos Não Aprovados
    "TextBox 31": "3.627",
    "TextBox 32": "4.460",
    "TextBox 33": "+833",
    # Documentos Não Enviados (R3 + R4 não anexados)
    "TextBox 36": "1.384",
    "TextBox 37": "1.501",
    "TextBox 38": "+117",
    # Aguardando Submissão
    "TextBox 41": "2.425",
    "TextBox 42": "1.545",
    "TextBox 43": "-880",
    # Documentos Em Análise (R3 + R4)
    "TextBox 46": "353",
    "TextBox 47": "441",
    "TextBox 48": "+88",
    # Documentos Vencidos (R4)
    "TextBox 51": "8",
    "TextBox 52": "1",
    "TextBox 53": "-7",
    # Total fornecedores
    "TextBox 55": "Total de Fornecedores na base: 49 → 49  (base estável entre as referências)",
    # Rodapé slide 1
    "TextBox 84": "Gabriel Vital  •  Customer Success Specialist  •  Efcaz SRM  •  Julho/2026",
}

# Slide 2 — cabeçalho e rodapé (conteúdo de contatos fica igual ao v2, atualizar manualmente se necessário)
slide2_updates = {
    "TextBox 5":  f"Zurich Airport  |  Dashboard Efcaz SRM  •  Julho/2026",
    "TextBox 84": "Gabriel Vital  •  Customer Success Specialist  •  Efcaz SRM  •  Julho/2026",
}

def substituir_shape(slide, updates):
    for shape in slide.shapes:
        if shape.name in updates and shape.has_text_frame:
            novo = updates[shape.name]
            tf = shape.text_frame
            # consolida todos os runs no primeiro parágrafo/run preservando formatação
            for para in tf.paragraphs:
                texto_atual = "".join(r.text for r in para.runs)
                if texto_atual.strip():
                    if para.runs:
                        para.runs[0].text = novo
                        for run in para.runs[1:]:
                            run.text = ""
                    break

substituir_shape(prs.slides[0], metricas)
substituir_shape(prs.slides[1], slide2_updates)

prs.save(OUTPUT)
print(f"Salvo: {OUTPUT}")
