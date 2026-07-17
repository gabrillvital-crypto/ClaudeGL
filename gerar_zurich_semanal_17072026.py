from pptx import Presentation
from pptx.util import Pt
import copy

TEMPLATE = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\dashboard-react\relatorio_impacto_efcaz_15-07-2026.pptx"
OUTPUT   = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\dashboard-react\relatorio_impacto_efcaz_17-07-2026.pptx"

prs = Presentation(TEMPLATE)

# ── Dados do comparativo 15/07 → 17/07 ────────────────────────────────────────
ant = "15/07/2026"
atu = "17/07/2026"

metricas = {
    # Slide 1 — cabeçalho
    "TextBox 5":  f"Zurich Airport  |  Dashboard Efcaz SRM  •  Base de {atu}",
    "TextBox 8":  f"8 indicadores ativos do dashboard  •  base anterior ({ant}) vs. nova base ({atu})",
    # Docs Esperados R4
    "TextBox 16": "1.643",
    "TextBox 17": "1.715",
    "TextBox 18": "+72",
    # Docs Esperados R3
    "TextBox 21": "9.188",
    "TextBox 22": "9.430",
    "TextBox 23": "+242",
    # Documentos Aprovados (R3 + R4)
    "TextBox 26": "2.227",
    "TextBox 27": "2.276",
    "TextBox 28": "+49",
    # Documentos Não Aprovados (Reprovados R3 + R4)
    "TextBox 31": "4.909",
    "TextBox 32": "4.929",
    "TextBox 33": "+20",
    # Documentos Não Enviados (R3 + R4)
    "TextBox 36": "1.540",
    "TextBox 37": "1.557",
    "TextBox 38": "+17",
    # Aguardando Submissão (R3)
    "TextBox 41": "1.528",
    "TextBox 42": "1.884",
    "TextBox 43": "+356",
    # Documentos Em Análise (R3 + R4)
    "TextBox 46": "624",
    "TextBox 47": "496",
    "TextBox 48": "-128",
    # Documentos Vencidos (R4)
    "TextBox 51": "3",
    "TextBox 52": "3",
    "TextBox 53": "= 0",
    # Total fornecedores
    "TextBox 55": "Total de Fornecedores na base: 49 → 49  (base estável entre as referências)",
    # Rodapé slide 1
    "TextBox 84": "Gabriel Vital  •  Customer Success Specialist  •  Efcaz SRM  •  Julho/2026",
}

# Slide 2 — cabeçalho e rodapé
slide2_updates = {
    "TextBox 5":  f"Zurich Airport  |  Dashboard Efcaz SRM  •  Julho/2026",
    "TextBox 84": "Gabriel Vital  •  Customer Success Specialist  •  Efcaz SRM  •  Julho/2026",
}

def substituir_shape(slide, updates):
    for shape in slide.shapes:
        if shape.name in updates and shape.has_text_frame:
            novo = updates[shape.name]
            tf = shape.text_frame
            for para in tf.paragraphs:
                texto_atual = "".join(r.text for r in para.runs)
                if texto_atual.strip():
                    if para.runs:
                        para.runs[0].text = novo
                        for run in para.runs[1:]:
                            run.text = ""
                    break

substituir_shape(prs.slides[0], metricas)
substituir_shape(prs.slides[1], slide2_updates)

prs.save(OUTPUT)
print(f"Salvo: {OUTPUT}")
