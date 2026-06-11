# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── CORES EFCAZ (extraídas do template real) ──────────────────────────────────
TEAL        = RGBColor(0x0E, 0x8F, 0xA3)
TEAL_CLARO  = RGBColor(0x0B, 0xA2, 0xB9)
TEAL_ESCURO = RGBColor(0x08, 0x71, 0x82)
VERDE       = RGBColor(0x7B, 0xB9, 0x30)
VERMELHO    = RGBColor(0xB3, 0x2B, 0x2B)
LARANJA     = RGBColor(0xF3, 0x9C, 0x12)
BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLARO = RGBColor(0xF2, 0xF4, 0xF4)
CINZA_MED   = RGBColor(0x66, 0x66, 0x66)
CINZA_ESC   = RGBColor(0x75, 0x75, 0x75)
PRETO       = RGBColor(0x43, 0x43, 0x43)
TEXT_LIGHT  = RGBColor(0xD6, 0xF5, 0xF8)
ECF         = RGBColor(0xEC, 0xF0, 0xF1)

F_TITLE = "Barlow Condensed"
F_BODY  = "Barlow"

W = 10.0
H = 5.62

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

# ── HELPERS ───────────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill, border=None, bw=1.2):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    return s

def tb(slide, text, x, y, w, h, size=10, bold=False, color=PRETO,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame; tf.word_wrap = wrap
    p   = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = font or F_BODY
    return box

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, W, 0.80, TEAL)
    rect(slide, 0, 0.73, W, 0.06, TEAL_CLARO)
    tb(slide, title, 0.25, 0.08, 7.8, 0.50, size=18, bold=False, color=BRANCO, font=F_TITLE)
    if subtitle:
        tb(slide, subtitle, 0.25, 0.52, 7.8, 0.24, size=9, color=TEXT_LIGHT, font=F_BODY)

def footer(slide):
    rect(slide, 0, 5.30, W, 0.32, TEAL_ESCURO)
    tb(slide, "Gabriel Vital  •  Customer Success Specialist  •  Efcaz SRM  •  Mai/2026",
       0.2, 5.32, 9.6, 0.26, size=8, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, font=F_TITLE)

def logo(slide):
    path = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\logo_efcaz.png"
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(8.55), Inches(0.10), width=Inches(1.25))

def bg(slide):
    rect(slide, 0, 0, W, H, CINZA_CLARO)

# ── SLIDE 1 — CAPA ────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
rect(s1, 0, 0, W, H, TEAL)
rect(s1, 0, 3.5, W, 2.12, TEAL_ESCURO)
rect(s1, 0, 3.42, W, 0.12, TEAL_CLARO)
tb(s1, "Critérios de Cálculo", 0.5, 1.2, 9.0, 0.9, size=30, bold=True,
   color=BRANCO, align=PP_ALIGN.CENTER, font=F_TITLE)
tb(s1, "Dashboard Zurich Airport", 0.5, 2.05, 9.0, 0.65, size=22,
   color=TEXT_LIGHT, align=PP_ALIGN.CENTER, font=F_TITLE)
tb(s1, "Lógica por trás dos indicadores  |  Efcaz × Zurich Airport",
   0.5, 2.65, 9.0, 0.4, size=11, color=BRANCO, align=PP_ALIGN.CENTER, font=F_BODY)
tb(s1, "Maio / 2026", 0.5, 3.9, 9.0, 0.4, size=11,
   color=BRANCO, align=PP_ALIGN.CENTER, font=F_TITLE)
logo(s1); footer(s1)

# ── SLIDE 2 — FONTES DE DADOS ─────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK); bg(s2)
header(s2, "Origem dos dados",
       "Três relatórios exportados da Efcaz — base de todos os cálculos")
footer(s2); logo(s2)

BOXES = [
    ("R1", "Pendências por Solicitação",
     ["Pendências abertas por fornecedor",
      "Tipo de documento exigido",
      "Área (Trabalhadores / Documentos)",
      "Status (EM_ELABORACAO / APROVADO)"], TEAL),
    ("R2", "Terceiros Cadastrados",
     ["Trabalhadores ativos por fornecedor",
      "Trabalhadores inativos",
      "Dados cadastrais dos terceiros"], CINZA_ESC),
    ("R3", "Situação Documental",
     ["Status de cada doc por trabalhador",
      "% Conformidade e Não Conformidade",
      "Docs vencidos / pendentes / conformes",
      "FILTRADO: última solicitação"], LARANJA),
]
XS = [0.25, 3.5, 6.75]
for i, (code, title, items, cor) in enumerate(BOXES):
    x = XS[i]
    rect(s2, x, 0.92, 3.05, 0.58, cor)
    tb(s2, code, x+0.08, 0.94, 0.6, 0.38, size=14, bold=True, color=BRANCO, font=F_TITLE)
    tb(s2, title, x+0.7, 0.97, 2.28, 0.5, size=9, bold=True, color=BRANCO, font=F_TITLE)
    rect(s2, x, 1.5, 3.05, 3.55, BRANCO, border=cor)
    for j, item in enumerate(items):
        is_warn = "FILTRADO" in item
        c = VERMELHO if is_warn else PRETO
        txt = ("! " if is_warn else "• ") + item
        tb(s2, txt, x+0.12, 1.62 + j*0.72, 2.85, 0.6, size=9, bold=is_warn, color=c, font=F_BODY)

# ── SLIDE 3 — CLASSIFICAÇÃO DE STATUS ────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK); bg(s3)
header(s3, "Como cada documento é classificado",
       "Relatório 3 — Situação Documental")
footer(s3); logo(s3)

rect(s3, 0.25, 0.90, 9.5, 0.38, TEAL)
tb(s3, "Cada linha do R3  =  1 documento  ×  1 trabalhador  ×  1 fornecedor",
   0.35, 0.92, 9.3, 0.34, size=10, bold=True, color=BRANCO,
   align=PP_ALIGN.CENTER, font=F_TITLE)

CX = [0.25, 3.35, 7.25]; CW = [2.95, 3.8, 2.5]
HCOLS = ["Status no dado", "Classificação", "Entra no cálculo?"]
rect(s3, 0.25, 1.38, 9.5, 0.4, TEAL_ESCURO)
for i, lbl in enumerate(HCOLS):
    tb(s3, lbl, CX[i]+0.08, 1.41, CW[i], 0.34, size=10, bold=True, color=BRANCO, font=F_TITLE)

ROWS3 = [
    ('"A vencer"',    "CONFORME",                      "SIM",                       VERDE),
    ('"Vencido"',     "VENCIDO  (Não Conforme)",        "SIM",                       VERMELHO),
    ('"Não anexado"', "PENDENTE  (Não Conforme)",       "SIM",                       LARANJA),
    ('"N/A"',         "Ignorar",                        "NÃO — excluído do total",   CINZA_ESC),
]
for ri, (status, classif, calc, cor) in enumerate(ROWS3):
    y = 1.88 + ri * 0.76
    row_bg = BRANCO if ri % 2 == 0 else ECF
    rect(s3, 0.25, y, 9.5, 0.68, row_bg)
    rect(s3, 0.25, y, 0.08, 0.68, cor)
    tb(s3, status,  CX[0]+0.12, y+0.16, CW[0]-0.1, 0.4, size=10, bold=True, font=F_BODY)
    tb(s3, classif, CX[1]+0.08, y+0.16, CW[1]-0.1, 0.4, size=10, bold=True, color=cor, font=F_BODY)
    tb(s3, calc,    CX[2]+0.08, y+0.16, CW[2]-0.1, 0.4, size=9, font=F_BODY)

# ── SLIDE 4 — FÓRMULA DO PERCENTUAL ──────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK); bg(s4)
header(s4, "Como o percentual é calculado", "Base de cálculo e fórmulas")
footer(s4); logo(s4)

rect(s4, 0.25, 0.9, 9.5, 0.72, TEAL)
tb(s4, "UNIVERSO  =  Conforme  +  Vencido  +  Pendente",
   0.35, 0.94, 9.3, 0.38, size=13, bold=True, color=BRANCO,
   align=PP_ALIGN.CENTER, font=F_TITLE)
tb(s4, "(linhas com status N/A são excluídas do total)",
   0.35, 1.28, 9.3, 0.28, size=9, color=TEXT_LIGHT,
   align=PP_ALIGN.CENTER, font=F_BODY)

rect(s4, 0.25, 1.75, 4.45, 3.3, VERMELHO)
tb(s4, "% NÃO CONFORMIDADE", 0.35, 1.82, 4.25, 0.42, size=12,
   color=BRANCO, align=PP_ALIGN.CENTER, font=F_TITLE)
tb(s4, "Vencido  +  Pendente", 0.35, 2.3, 4.25, 0.38, size=12, bold=True,
   color=BRANCO, align=PP_ALIGN.CENTER, font=F_BODY)
rect(s4, 0.45, 2.74, 4.05, 0.05, BRANCO)
tb(s4, "Universo", 0.35, 2.8, 4.25, 0.35, size=12,
   color=BRANCO, align=PP_ALIGN.CENTER, font=F_BODY)
tb(s4, "× 100", 0.35, 3.18, 4.25, 0.35, size=12, bold=True,
   color=BRANCO, align=PP_ALIGN.CENTER, font=F_BODY)
tb(s4, "Resultado atual: 60,3%", 0.35, 3.6, 4.25, 0.32, size=10,
   color=BRANCO, align=PP_ALIGN.CENTER, font=F_BODY)
rect(s4, 0.25, 4.0, 4.45, 0.75, RGBColor(0x8B, 0x1A, 0x1A))
tb(s4, "Quanto do total está irregular\n(vencido ou nunca enviado)",
   0.35, 4.05, 4.25, 0.65, size=9, color=BRANCO, align=PP_ALIGN.CENTER, font=F_BODY)

tb(s4, "=\n100%", 4.75, 3.1, 0.5, 0.6, size=10, bold=True, color=PRETO,
   align=PP_ALIGN.CENTER, font=F_TITLE)

rect(s4, 5.3, 1.75, 4.45, 3.3, VERDE)
tb(s4, "% CONFORMIDADE", 5.4, 1.82, 4.25, 0.42, size=12,
   color=BRANCO, align=PP_ALIGN.CENTER, font=F_TITLE)
tb(s4, "Conforme", 5.4, 2.3, 4.25, 0.38, size=12, bold=True,
   color=BRANCO, align=PP_ALIGN.CENTER, font=F_BODY)
rect(s4, 5.5, 2.74, 4.05, 0.05, BRANCO)
tb(s4, "Universo", 5.4, 2.8, 4.25, 0.35, size=12,
   color=BRANCO, align=PP_ALIGN.CENTER, font=F_BODY)
tb(s4, "× 100", 5.4, 3.18, 4.25, 0.35, size=12, bold=True,
   color=BRANCO, align=PP_ALIGN.CENTER, font=F_BODY)
tb(s4, "Resultado atual: 39,7%", 5.4, 3.6, 4.25, 0.32, size=10,
   color=BRANCO, align=PP_ALIGN.CENTER, font=F_BODY)
rect(s4, 5.3, 4.0, 4.45, 0.75, RGBColor(0x4A, 0x7A, 0x1E))
tb(s4, "Quanto do total está em dia\n(dentro do prazo de validade)",
   5.4, 4.05, 4.25, 0.65, size=9, color=BRANCO, align=PP_ALIGN.CENTER, font=F_BODY)

# ── SLIDE 5 — GERAL vs POR FORNECEDOR ────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK); bg(s5)
header(s5, "Dois níveis de cálculo",
       "O percentual pode ser visto de forma global ou individual por fornecedor")
footer(s5); logo(s5)

rect(s5, 0.25, 0.9, 4.45, 4.15, TEAL)
tb(s5, "VISÃO GERAL", 0.35, 0.97, 4.25, 0.42, size=13,
   color=BRANCO, align=PP_ALIGN.CENTER, font=F_TITLE)
rect(s5, 0.35, 1.48, 4.25, 3.42, BRANCO)
GERAL = [
    ("Universo",  "Todos os documentos\nde todos os fornecedores"),
    ("Resultado", "1 único percentual\npara o aeroporto inteiro"),
    ("Exemplo",   "60,3% de Não Conformidade\ngeral Zurich Airport"),
    ("Uso",       "KPI principal no topo\ndo dashboard"),
]
for j, (lbl, val) in enumerate(GERAL):
    y = 1.56 + j*0.82
    tb(s5, lbl, 0.48, y, 1.1, 0.28, size=8, bold=True, color=TEAL, font=F_TITLE)
    tb(s5, val, 0.48, y+0.26, 4.0, 0.5, size=9, font=F_BODY)

tb(s5, "->", 4.77, 2.7, 0.46, 0.46, size=18, bold=True, color=PRETO,
   align=PP_ALIGN.CENTER, font=F_TITLE)

rect(s5, 5.3, 0.9, 4.45, 4.15, LARANJA)
tb(s5, "POR FORNECEDOR", 5.4, 0.97, 4.25, 0.42, size=13,
   color=BRANCO, align=PP_ALIGN.CENTER, font=F_TITLE)
rect(s5, 5.4, 1.48, 4.25, 3.42, BRANCO)
FORN = [
    ("Universo",  "Apenas os documentos\ndaquele fornecedor"),
    ("Resultado", "1 percentual por empresa\n(denominador individual)"),
    ("Exemplo",   "Fornecedor A: 33% NC\nFornecedor B: 87% NC"),
    ("Uso",       "Gráfico de barras e drill-down\nno dashboard"),
]
for j, (lbl, val) in enumerate(FORN):
    y = 1.56 + j*0.82
    tb(s5, lbl, 5.53, y, 1.1, 0.28, size=8, bold=True, color=LARANJA, font=F_TITLE)
    tb(s5, val, 5.53, y+0.26, 4.0, 0.5, size=9, font=F_BODY)

# ── SLIDE 6 — ÚLTIMA SOLICITAÇÃO ─────────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK); bg(s6)
header(s6, "Filtro de Última Solicitação",
       "Ponto de atenção para replicar no BD")
footer(s6); logo(s6)

rect(s6, 0.25, 0.90, 9.5, 1.05, LARANJA)
tb(s6, "O Relatório 3 já vem pré-filtrado para a ÚLTIMA solicitação de cada fornecedor",
   0.38, 0.95, 9.26, 0.42, size=11, bold=True, color=BRANCO,
   align=PP_ALIGN.CENTER, font=F_TITLE)
tb(s6, "Não considera o histórico completo — apenas o estado mais recente de cada empresa",
   0.38, 1.35, 9.26, 0.34, size=9, color=BRANCO,
   align=PP_ALIGN.CENTER, font=F_BODY)

rect(s6, 0.25, 2.1, 4.45, 2.95, BRANCO, border=VERDE)
rect(s6, 0.25, 2.1, 4.45, 0.42, VERDE)
tb(s6, "SE o filtro for aplicado", 0.35, 2.13, 4.25, 0.36, size=11,
   bold=True, color=BRANCO, align=PP_ALIGN.CENTER, font=F_TITLE)
OK = ["Os percentuais batem com o dashboard",
      "Cada fornecedor reflete o estado atual",
      "Visão consistente com o que a Débora\nvê na plataforma Efcaz"]
for j, line in enumerate(OK):
    tb(s6, f"• {line}", 0.38, 2.65 + j*0.72, 4.2, 0.62, size=9, font=F_BODY)

rect(s6, 5.3, 2.1, 4.45, 2.95, BRANCO, border=VERMELHO)
rect(s6, 5.3, 2.1, 4.45, 0.42, VERMELHO)
tb(s6, "SEM o filtro", 5.4, 2.13, 4.25, 0.36, size=11,
   bold=True, color=BRANCO, align=PP_ALIGN.CENTER, font=F_TITLE)
NOK = ["Solicitações antigas entram no cálculo",
       "Percentuais ficam inflados\nou inconsistentes",
       "Números não batem\ncom o dashboard"]
for j, line in enumerate(NOK):
    tb(s6, f"• {line}", 5.43, 2.65 + j*0.72, 4.2, 0.62, size=9, font=F_BODY)

rect(s6, 0.25, 4.98, 9.5, 0.25, TEAL_ESCURO)
tb(s6, "No BD: filtrar por MAX(solicitacao_id) por fornecedor — ou equivalente na estrutura da Efcaz",
   0.38, 5.00, 9.26, 0.20, size=8, bold=True, color=BRANCO,
   align=PP_ALIGN.CENTER, font=F_TITLE)

# ── SLIDE 7 — KPIs E FONTES ───────────────────────────────────────────────────
s7 = prs.slides.add_slide(BLANK); bg(s7)
header(s7, "De onde vem cada KPI",
       "Mapeamento completo de indicadores, fontes e filtros")
footer(s7); logo(s7)

KPI_ROWS = [
    ("Fornecedores com pendências",  "R1", "COUNT DISTINCT fornecedor"),
    ("Total pendências",             "R1", "COUNT todas as linhas"),
    ("Pendente (não enviado)",       "R1", 'Situação da solicitação = "EM_ELABORACAO"'),
    ("Em análise com pendências",   "R1", 'Situação da solicitação = "APROVADO"'),
    ("Pend. Trabalhadores",         "R1", 'Área da pendência = "TERCEIROS"'),
    ("Pend. Documentais",           "R1", 'Área da pendência = "DOCUMENTOS"'),
    ("Trabalhadores Ativos",        "R2", 'Status = "Ativo"'),
    ("Docs Conformes",              "R3", 'Status = "A vencer"'),
    ("Docs Vencidos",               "R3", 'Status = "Vencido"'),
    ("Docs Pendentes",              "R3", 'Status = "Não anexado"'),
    ("% Não Conformidade",          "R3", "(Vencido + Pendente) / Universo × 100"),
    ("% Conformidade",              "R3", "Conforme / Universo × 100"),
]
SRC_CLR = {"R1": TEAL, "R2": CINZA_ESC, "R3": LARANJA}
KCX = [0.25, 5.15, 6.2]; KCW = [4.78, 0.95, 3.6]
KHCOLS = ["KPI", "Fonte", "Filtro / Operação"]
rect(s7, 0.25, 0.90, 9.5, 0.38, TEAL_ESCURO)
for i, lbl in enumerate(KHCOLS):
    tb(s7, lbl, KCX[i]+0.08, 0.92, KCW[i], 0.32, size=9, bold=True, color=BRANCO, font=F_TITLE)
for ri, (kpi, src, filt) in enumerate(KPI_ROWS):
    y = 1.35 + ri * 0.33
    row_bg = BRANCO if ri % 2 == 0 else ECF
    rect(s7, 0.25, y, 9.5, 0.33, row_bg)
    rect(s7, 0.25, y, 0.07, 0.33, SRC_CLR.get(src, CINZA_ESC))
    tb(s7, kpi, 0.37, y+0.06, 4.65, 0.24, size=9, font=F_BODY)
    rect(s7, 5.25, y+0.05, 0.6, 0.22, SRC_CLR.get(src, CINZA_ESC))
    tb(s7, src, 5.27, y+0.06, 0.56, 0.2, size=8, bold=True, color=BRANCO,
       align=PP_ALIGN.CENTER, font=F_TITLE)
    tb(s7, filt, 6.28, y+0.06, 3.4, 0.24, size=9, font=F_BODY)

# ── SLIDE 8 — FLUXO AUTOMÁTICO ────────────────────────────────────────────────
s8 = prs.slides.add_slide(BLANK); bg(s8)
header(s8, "Modelo de atualização automática",
       "Script agendado consulta o BD e regenera o HTML uma vez por dia às 12h")
footer(s8); logo(s8)

ETAPAS = [
    ("BD\nZurich",      "Credencial\nread-only\n(Ricardo)"),
    ("Script\nPython",  "Consulta o BD\ne gera o HTML"),
    ("Agenda-\nmento",  "Atualiza\n1x por dia\nàs 12h"),
    ("Servidor\nEfcaz", "HTML salvo\nem link fixo"),
    ("Débora",          "Abre o link\nsempre\natualizado"),
]
N = len(ETAPAS)
BOX_W = 1.5; BOX_H = 1.1; ARROW_W = 0.35
TOTAL = N * BOX_W + (N-1) * ARROW_W
START_X = (W - TOTAL) / 2
Y_BOX = 1.7

for i, (titulo, desc) in enumerate(ETAPAS):
    x = START_X + i * (BOX_W + ARROW_W)
    rect(s8, x, Y_BOX, BOX_W, 0.45, TEAL)
    tb(s8, titulo, x, Y_BOX, BOX_W, 0.45, size=10, bold=True,
       color=BRANCO, align=PP_ALIGN.CENTER, font=F_TITLE)
    rect(s8, x, Y_BOX+0.45, BOX_W, BOX_H-0.45, BRANCO, border=TEAL, bw=0.8)
    tb(s8, desc, x+0.05, Y_BOX+0.50, BOX_W-0.1, BOX_H-0.55, size=8,
       color=PRETO, align=PP_ALIGN.CENTER, font=F_BODY)
    if i < N-1:
        ax = x + BOX_W + 0.04
        tb(s8, "->", ax, Y_BOX+0.25, ARROW_W-0.08, 0.42, size=14, bold=True,
           color=TEAL, align=PP_ALIGN.CENTER, font=F_TITLE)

rect(s8, 0.25, 3.22, 9.5, 0.04, TEAL_CLARO)
tb(s8, "O HTML, o visual e o drill-down não mudam — só a fonte dos dados muda (BD em vez de CSV)",
   0.25, 3.34, 9.5, 0.32, size=9, color=CINZA_MED,
   align=PP_ALIGN.CENTER, font=F_BODY)

rect(s8, 0.25, 3.75, 9.5, 1.22, BRANCO, border=TEAL_CLARO, bw=0.8)
rect(s8, 0.25, 3.75, 9.5, 0.35, TEAL_CLARO)
tb(s8, "Pré-requisitos para implementação", 0.35, 3.78, 9.3, 0.28,
   size=9, bold=True, color=BRANCO, font=F_TITLE)
NEEDS = [
    "Credencial read-only para o BD do ambiente Zurich (host, porta, usuário, senha)",
    "Endereço do servidor para publicação do HTML e configuração do agendador (cron / Task Scheduler)",
]
for j, n in enumerate(NEEDS):
    tb(s8, f"• {n}", 0.4, 4.18 + j*0.35, 9.2, 0.3, size=9, color=PRETO, font=F_BODY)

# ── SLIDE 9 — ACESSO CONTÍNUO ────────────────────────────────────────────────
s9 = prs.slides.add_slide(BLANK); bg(s9)
header(s9, "Acesso contínuo",
       "Usuário acessa link fixo — HTML atualizado automaticamente pelo agendador")
footer(s9); logo(s9)

rect(s9, 0.25, 0.90, 6.1, 4.15, BRANCO, border=TEAL)
rect(s9, 0.25, 0.90, 6.1, 0.42, TEAL)
tb(s9, "Premissas do modelo", 0.38, 0.93, 5.9, 0.36,
   size=11, color=BRANCO, font=F_TITLE)

REASONS9 = [
    ("Ambiente isolado por cliente",
     "Dados da Zurich separados dos demais clientes na base — sem cruzamento"),
    ("Credencial read-only",
     "Script só lê o BD — sem permissão de escrita"),
    ("Agendamento diário às 12h",
     "Execução automática sem intervenção manual — frequência configurável"),
    ("Link fixo",
     "Usuário acessa sempre o mesmo endereço — o HTML é sobrescrito a cada ciclo"),
    ("Dado na infraestrutura Efcaz",
     "Não transita por serviços externos — LGPD: dado permanece no ambiente da Efcaz"),
]
for j, (title, desc) in enumerate(REASONS9):
    y = 1.42 + j*0.72
    rect(s9, 0.35, y+0.08, 0.06, 0.24, TEAL)
    tb(s9, title, 0.52, y+0.06, 5.7, 0.26, size=10, bold=True, font=F_TITLE)
    tb(s9, desc,  0.52, y+0.30, 5.7, 0.32, size=8, color=CINZA_MED, font=F_BODY)

rect(s9, 6.6, 0.90, 3.15, 4.15, BRANCO, border=TEAL_CLARO)
rect(s9, 6.6, 0.90, 3.15, 0.42, TEAL_CLARO)
tb(s9, "Horários sugeridos", 6.7, 0.93, 2.95, 0.36, size=11,
   color=BRANCO, font=F_TITLE)

rect(s9, 6.72, 1.55, 0.88, 0.52, TEAL)
tb(s9, "12h00", 6.74, 1.63, 0.84, 0.36, size=13, bold=True,
   color=BRANCO, align=PP_ALIGN.CENTER, font=F_TITLE)
tb(s9, "Atualização diária\nno horário de almoço —\ndados disponíveis para\na tarde de trabalho.",
   7.7, 1.55, 1.92, 0.9, size=9, color=PRETO, font=F_BODY)

tb(s9, "Horário ajustável conforme\npreferência da operação.",
   6.62, 4.1, 3.05, 0.85, size=8, color=CINZA_MED, font=F_BODY)

rect(s9, 0.25, 4.97, 9.5, 0.26, TEAL_ESCURO)
tb(s9, "Débora acessa o link a qualquer momento — os dados refletem sempre a última atualização automática.",
   0.38, 4.99, 9.26, 0.2, size=8, color=TEXT_LIGHT,
   align=PP_ALIGN.CENTER, font=F_BODY)

# ── SALVAR ────────────────────────────────────────────────────────────────────
OUTPUT = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\Zurich_Criterios_Calculo_Mai2026.pptx"
prs.save(OUTPUT)
print(f"PPTX salvo: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
