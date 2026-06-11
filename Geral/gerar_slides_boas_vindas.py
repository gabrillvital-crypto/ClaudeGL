"""
Gera 2 slides adicionais para o deck de Boas-vindas do Onboarding Efcaz.

Slide 1 — Jornada de Onboarding  : timeline visual dos 7 encontros
Slide 2 — Cronograma Sugerido    : tabela com datas (+7 dias por encontro)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paleta Efcaz ──────────────────────────────────────────────────────────────
TEAL       = RGBColor(0x0E, 0x8F, 0xA3)
TEAL_ESC   = RGBColor(0x0A, 0x6A, 0x7A)
TEAL_CLR   = RGBColor(0x14, 0xB3, 0xCC)
BRANCO     = RGBColor(0xFF, 0xFF, 0xFF)
PRETO_SUV  = RGBColor(0x1A, 0x2A, 0x2A)
CINZA_CLR  = RGBColor(0xF2, 0xF4, 0xF4)
CINZA_MED  = RGBColor(0x85, 0x92, 0x9E)
VERDE      = RGBColor(0x27, 0xAE, 0x60)
LARANJA    = RGBColor(0xF3, 0x9C, 0x12)

ENCONTROS = [
    ("1º",  "Boas-vindas",                                       TEAL_ESC),
    ("2º",  "Fluxo de Homologação\ne Análise de Cadastro",       TEAL),
    ("3º",  "Preenchimento\nde Planilhas",                        TEAL_CLR),
    ("4º",  "Configurações da Base,\nRelatórios e Outros",        TEAL),
    ("5º",  "Prática Assistida",                                  TEAL_ESC),
    ("6º",  "Avaliações, Ocorrências\ne Gestão de Terceiros",     TEAL),
    ("7º",  "Check-In Final,\nDúvidas e Próximos Passos",         VERDE),
]

# ── Helpers ───────────────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=12, bold=False, color=PRETO_SUV,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb

def header(slide, titulo, subtitulo=""):
    rect(slide, 0, 0, 13.33, 1.1, fill=TEAL)
    rect(slide, 0, 1.05, 13.33, 0.08, fill=TEAL_CLR)
    txt(slide, titulo, 0.5, 0.15, 10.5, 0.75, size=26, bold=True, color=BRANCO)
    if subtitulo:
        txt(slide, subtitulo, 0.5, 1.18, 10.5, 0.38, size=11, color=CINZA_MED)
    txt(slide, "Efcaz SRM", 10.8, 0.32, 2.2, 0.45,
        size=11, bold=True, color=BRANCO, align=PP_ALIGN.RIGHT)

def rodape(slide, texto="Efcaz SRM  •  Onboarding"):
    rect(slide, 0, 7.15, 13.33, 0.35, fill=TEAL_ESC)
    txt(slide, texto, 0.4, 7.2, 12.5, 0.28, size=10, color=BRANCO)

# ── Apresentação ─────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Jornada de Onboarding (timeline visual)
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Nossa Jornada de Onboarding",
       "7 encontros para que você extraia o máximo da plataforma Efcaz")

# Linha da timeline
rect(slide, 0.55, 3.85, 12.25, 0.08, fill=TEAL_CLR)

# Nós + cards por encontro
card_w   = 1.6
card_gap = 0.155
start_x  = 0.55

for i, (num, nome, cor) in enumerate(ENCONTROS):
    cx = start_x + i * (card_w + card_gap)
    cy = 3.61   # centro Y da linha de nó

    # Nó circular (usando retângulo arredondado — aproximação via retângulo cheio pequeno)
    node = slide.shapes.add_shape(
        9,  # RoundedRectangle
        Inches(cx + card_w/2 - 0.16),
        Inches(cy),
        Inches(0.32), Inches(0.32)
    )
    node.fill.solid(); node.fill.fore_color.rgb = cor
    node.line.fill.background()
    adj = node.adjustments
    if adj: adj[0] = 50000  # totalmente arredondado

    # Card acima (ímpares) ou abaixo (pares) alternando
    if i % 2 == 0:
        card_top = 1.4
        connector_top = card_top + 1.85
        connector_h   = cy - connector_top
    else:
        card_top = 4.35
        connector_top = cy + 0.32
        connector_h   = card_top - connector_top

    # Linha conectora
    rect(slide, cx + card_w/2 - 0.03, connector_top, 0.06, connector_h, fill=cor)

    # Card de fundo
    rect(slide, cx, card_top, card_w, 1.85, fill=CINZA_CLR)
    rect(slide, cx, card_top, card_w, 0.38, fill=cor)

    # Número do encontro
    txt(slide, f"{num} Enc.", cx + 0.08, card_top + 0.04,
        card_w - 0.16, 0.3, size=11, bold=True, color=BRANCO)

    # Nome do módulo (multiline)
    tb = slide.shapes.add_textbox(
        Inches(cx + 0.1), Inches(card_top + 0.45),
        Inches(card_w - 0.2), Inches(1.3)
    )
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    for j, line in enumerate(nome.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.size  = Pt(10)
        r.font.name  = "Calibri"
        r.font.color.rgb = PRETO_SUV

rodape(slide, "Efcaz SRM  •  Onboarding  •  [Nome do Cliente]")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Cronograma Sugerido (tabela com datas)
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Cronograma Sugerido de Encontros",
       "Datas de referência — ajuste conforme disponibilidade")

# Legenda explicativa
rect(slide, 0.5, 1.25, 12.33, 0.38, fill=RGBColor(0xE8, 0xF8, 0xFA))
txt(slide, "  O 1º encontro (Boas-vindas) marca o início oficial da jornada."
    " Os demais são agendados com intervalo de 7 dias — edite a data do 1º encontro no Playbook e as demais se preenchem automaticamente.",
    0.55, 1.27, 12.1, 0.34, size=10, italic=True, color=TEAL_ESC)

# Cabeçalho da tabela
cols = [("Encontro", 1.8), ("Tema", 5.8), ("Data Sugerida", 2.5), ("Status", 2.05)]
row_h = 0.54
hdr_t = 1.75
x = 0.5
for col_txt, col_w in cols:
    rect(slide, x, hdr_t, col_w, row_h, fill=TEAL_ESC)
    txt(slide, col_txt, x + 0.1, hdr_t + 0.12, col_w - 0.2, row_h - 0.15,
        size=12, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    x += col_w

# Linhas de dados
TEMAS = [
    "Boas-vindas",
    "Fluxo de Homologação e Análise de Cadastro",
    "Preenchimento de Planilhas",
    "Configurações da Base, Relatórios e Outros",
    "Prática Assistida",
    "Avaliações, Ocorrências e Gestão de Terceiros",
    "Check-In Final, Dúvidas e Próximos Passos",
]

# Cores de status por encontro (todos "A agendar" inicialmente)
STATUS = ["A agendar"] * 7

for i, (enc, nome, tema, status) in enumerate(
        zip([e[0] for e in ENCONTROS],
            [e[1].replace("\n", " ") for e in ENCONTROS],
            TEMAS,
            STATUS)):
    row_t = hdr_t + (i + 1) * row_h
    bg = CINZA_CLR if i % 2 == 0 else BRANCO

    # Cor do número
    cor_enc = ENCONTROS[i][2]
    delta_days = i * 7

    if i == 0:
        date_str = "[Definir com o cliente]"
        date_color = LARANJA
    else:
        date_str = f"1º encontro + {delta_days} dias"
        date_color = CINZA_MED

    row_data = [
        (f"{enc} Encontro", 1.8, BRANCO, cor_enc),
        (tema,              5.8, PRETO_SUV, bg),
        (date_str,          2.5, date_color, bg),
        (status,            2.05, CINZA_MED, bg),
    ]

    x = 0.5
    for cell_txt, col_w, font_col, bg_col in row_data:
        rect(slide, x, row_t, col_w, row_h, fill=bg_col)
        txt(slide, cell_txt, x + 0.1, row_t + 0.12,
            col_w - 0.2, row_h - 0.15,
            size=10, color=font_col,
            bold=(bg_col == cor_enc))
        x += col_w

    # Separador
    rect(slide, 0.5, row_t + row_h - 0.02, 12.33, 0.02, fill=TEAL_CLR)

# Nota de rodapé
txt(slide,
    "  Dica: preencha a data do 1º encontro no Playbook de Onboarding "
    "e o cronograma completo será calculado automaticamente.",
    0.5, 7.05, 12.33, 0.28, size=9, italic=True, color=CINZA_MED)

rodape(slide, "Efcaz SRM  •  Onboarding  •  [Nome do Cliente]")

# ── Salvar ────────────────────────────────────────────────────────────────────
OUTPUT = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\Slides_Boas_Vindas_Onboarding.pptx"
prs.save(OUTPUT)
print(f"Slides gerados: {OUTPUT}")
