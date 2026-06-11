from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Paleta Efcaz
EFCAZ_PRIM    = RGBColor(0x0E, 0x8F, 0xA3)  # #0e8fa3
EFCAZ_SEC     = RGBColor(0x14, 0xB3, 0xCC)  # #14b3cc
EFCAZ_ESCURO  = RGBColor(0x0A, 0x6A, 0x7A)  # versão mais escura da primária
BRANCO        = RGBColor(0xFF, 0xFF, 0xFF)
PRETO_SUAVE   = RGBColor(0x1A, 0x2A, 0x2A)
CINZA_CLARO   = RGBColor(0xF2, 0xF4, 0xF4)
CINZA_MEDIO   = RGBColor(0x85, 0x92, 0x9E)
VERDE         = RGBColor(0x27, 0xAE, 0x60)
VERDE_CLARO   = RGBColor(0xD5, 0xF5, 0xE3)
VERMELHO      = RGBColor(0xE7, 0x4C, 0x3C)
LARANJA       = RGBColor(0xF3, 0x9C, 0x12)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def rect(slide, l, t, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    s.line.fill.background() if not line else None
    if line: s.line.color.rgb = line
    return s

def txt(slide, text, l, t, w, h, size=13, bold=False, color=PRETO_SUAVE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb

def multiline(slide, lines, l, t, w, h, size=12, color=PRETO_SUAVE, bold=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Calibri"

def header(slide, titulo, subtitulo=""):
    rect(slide, 0, 0, 13.33, 1.1, fill=EFCAZ_PRIM)
    rect(slide, 0, 1.05, 13.33, 0.08, fill=EFCAZ_SEC)
    txt(slide, titulo, 0.5, 0.15, 12, 0.75, size=28, bold=True, color=BRANCO)
    if subtitulo:
        txt(slide, subtitulo, 0.5, 1.2, 12, 0.38, size=12, color=CINZA_MEDIO)

def kpi_card(slide, l, t, w, h, number, label, color=EFCAZ_PRIM):
    rect(slide, l, t, w, h, fill=color)
    txt(slide, number, l+0.1, t+0.12, w-0.2, h*0.52,
        size=30, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    txt(slide, label, l+0.05, t+h*0.58, w-0.1, h*0.38,
        size=10, color=BRANCO, align=PP_ALIGN.CENTER, wrap=True)

def rodape(slide, texto):
    rect(slide, 0, 7.15, 13.33, 0.35, fill=EFCAZ_ESCURO)
    txt(slide, texto, 0.4, 7.2, 12.5, 0.28, size=10, color=BRANCO)

# ─────────────────────────────────────────────
# SLIDE 1 — CAPA
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=EFCAZ_PRIM)
rect(slide, 0, 4.8, 13.33, 2.7, fill=EFCAZ_ESCURO)
rect(slide, 0.5, 2.5, 0.1, 1.8, fill=BRANCO)

txt(slide, "[NOME DO CLIENTE]", 0.85, 1.3, 11, 1.1,
    size=46, bold=True, color=BRANCO)
txt(slide, "[Razão Social / Segmento]", 0.85, 2.5, 9, 0.55,
    size=20, color=EFCAZ_SEC)
txt(slide, "Análise de Renovação  |  [Mês/Ano]", 0.85, 5.15, 9, 0.5,
    size=18, color=BRANCO)
txt(slide, "Gabriel Vital  •  Customer Success Specialist  •  Efcaz SRM",
    0.85, 5.75, 10, 0.45, size=13, color=EFCAZ_SEC)
txt(slide, "Confidencial", 10.5, 7.0, 2.5, 0.38,
    size=10, color=RGBColor(0xAE, 0xD6, 0xF1), align=PP_ALIGN.RIGHT, italic=True)

# ─────────────────────────────────────────────
# SLIDE 2 — AGENDA
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Agenda", "O que vamos abordar hoje")

itens = [
    ("01", "Uso e Adoção da Plataforma",    "KPIs de uso, fornecedores e terceiros"),
    ("02", "Crescimento no Período",         "Evolução de solicitações e terceiros cadastrados"),
    ("03", "Pontos de Alerta",               "Módulos subutilizados e riscos identificados"),
    ("04", "Health Score",                   "Avaliação geral da saúde do cliente"),
    ("05", "Oportunidades de Expansão",      "Capacidade, módulos e integrações"),
    ("06", "Próximos Passos",                "Plano de ação e renovação"),
]

for i, (num, title, desc) in enumerate(itens):
    top = 1.45 + i * 0.95
    bg = CINZA_CLARO if i % 2 == 0 else BRANCO
    rect(slide, 0.5, top, 12.33, 0.82, fill=bg)
    rect(slide, 0.5, top, 0.7, 0.82, fill=EFCAZ_PRIM)
    txt(slide, num, 0.52, top+0.2, 0.65, 0.42, size=15, bold=True,
        color=BRANCO, align=PP_ALIGN.CENTER)
    txt(slide, title, 1.35, top+0.08, 5.5, 0.36, size=14, bold=True, color=EFCAZ_ESCURO)
    txt(slide, desc,  1.35, top+0.46, 5.5, 0.32, size=11, color=CINZA_MEDIO)

# ─────────────────────────────────────────────
# SLIDE 3 — KPIs DE USO
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Uso e Adoção da Plataforma", "Dados extraídos do BI de usabilidade  •  [Mês/Ano]")

kpis_top = [
    ("[000]", "Fornecedores\nAtivos",      EFCAZ_ESCURO),
    ("[000]", "Terceiros\nCadastrados",    EFCAZ_PRIM),
    ("[000]", "Solicitações\nAbertas",     EFCAZ_SEC),
    ("[00%]", "Taxa de\nAprovação",        VERDE),
]
cw, ch = 2.8, 1.65
sx = 0.55
for j, (n, l, c) in enumerate(kpis_top):
    kpi_card(slide, sx + j*(cw+0.25), 1.72, cw, ch, n, l, c)

kpis_bot = [
    ("[00]",  "Usuários\nCadastrados",     LARANJA),
    ("[000]", "Fornecedores\nContratados", CINZA_MEDIO),
    ("[00]",  "Módulos\nAtivos",           EFCAZ_PRIM),
    ("[000]", "Média de\nHomologação",     RGBColor(0x7F, 0x8C, 0x8D)),
]
for j, (n, l, c) in enumerate(kpis_bot):
    kpi_card(slide, sx + j*(cw+0.25), 3.6, cw, ch, n, l, c)

rect(slide, 0.5, 5.4, 12.33, 0.48, fill=RGBColor(0xFE, 0xF9, 0xE7))
txt(slide, "⚠  [Insira aqui observação de alerta sobre capacidade ou limite de usuários]",
    0.65, 5.47, 12, 0.36, size=11, color=LARANJA, bold=True)

rodape(slide, "Efcaz SRM  •  Análise de Renovação  •  [Nome do Cliente]  •  [Mês/Ano]")

# ─────────────────────────────────────────────
# SLIDE 4 — CRESCIMENTO
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Crescimento no Período", "Evolução mensal de solicitações e terceiros cadastrados")

# Área de gráfico — solicitações (placeholder visual)
rect(slide, 0.5, 1.45, 12.33, 0.38, fill=CINZA_CLARO)
txt(slide, "Solicitações por Mês  •  [período]", 0.65, 1.5, 8, 0.28,
    size=13, bold=True, color=EFCAZ_ESCURO)

# Barras placeholder (8 barras representativas)
bar_vals  = [0.3, 0.5, 0.45, 0.55, 0.85, 0.72, 0.82, 1.0]
bar_label = ["M1","M2","M3","M4","M5","M6","M7","M8"]
bw, bh_max = 1.35, 1.8
for i, (v, lbl) in enumerate(zip(bar_vals, bar_label)):
    bh = bh_max * v
    bx = 0.55 + i*(bw+0.07)
    by = 1.95 + bh_max - bh
    col = EFCAZ_SEC if i < 6 else VERDE
    rect(slide, bx, by, bw, bh, fill=col)
    txt(slide, lbl, bx, 3.82, bw, 0.25, size=9, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

rect(slide, 0.5, 4.15, 12.33, 0.05, fill=CINZA_MEDIO)

# Barras placeholder — terceiros
rect(slide, 0.5, 4.28, 12.33, 0.38, fill=CINZA_CLARO)
txt(slide, "Terceiros Cadastrados por Mês  •  [período]", 0.65, 4.33, 9, 0.28,
    size=13, bold=True, color=EFCAZ_ESCURO)

terc_vals = [0.4, 0.58, 0.62, 0.68, 0.83, 0.8, 0.82, 1.0]
for i, (v, lbl) in enumerate(zip(terc_vals, bar_label)):
    bh = 1.5 * v
    bx = 0.55 + i*(bw+0.07)
    by = 4.75 + 1.5 - bh
    col = VERDE if i == 7 else EFCAZ_PRIM
    rect(slide, bx, by, bw, bh, fill=col)
    txt(slide, lbl, bx, 6.3, bw, 0.22, size=8, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

txt(slide, "[Insira aqui o insight de crescimento do período — ex: X vezes mais solicitações em Y meses]",
    0.5, 7.08, 12.5, 0.3, size=10, italic=True, color=CINZA_MEDIO)

# ─────────────────────────────────────────────
# SLIDE 5 — PONTOS DE ALERTA
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Pontos de Alerta — Módulos Subutilizados", "")

alertas = [
    (VERMELHO, "🔴 [Título do Alerta 1]",
     "[Descrição do alerta — ex: módulo parado, integração inativa]\n[Pergunta estratégica a levantar na reunião]"),
    (VERMELHO, "🔴 [Título do Alerta 2]",
     "[Descrição do alerta — ex: integrações automáticas inativas]\n[Impacto operacional para o cliente]"),
    (LARANJA,  "🟡 [Título do Alerta 3]",
     "[Descrição — ex: recurso disponível mas não configurado]\n[Oportunidade de ativação e redução de retrabalho]"),
    (LARANJA,  "🟡 [Título do Alerta 4]",
     "[Descrição — ex: usuários próximos do limite]\n[Ação preventiva recomendada]"),
]

for i, (cor, titulo, desc) in enumerate(alertas):
    top = 1.35 + i * 1.45
    rect(slide, 0.5, top, 12.33, 1.3, fill=CINZA_CLARO)
    rect(slide, 0.5, top, 0.18, 1.3, fill=cor)
    txt(slide, titulo, 0.85, top+0.1, 11.5, 0.36, size=13, bold=True, color=PRETO_SUAVE)
    multiline(slide, desc.split("\n"), 0.85, top+0.52, 11.5, 0.72, size=11, color=CINZA_MEDIO)

rodape(slide, "Efcaz SRM  •  Análise de Renovação  •  [Nome do Cliente]  •  [Mês/Ano]")

# ─────────────────────────────────────────────
# SLIDE 6 — HEALTH SCORE
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Health Score — [Nome do Cliente]", "")

rect(slide, 4.9, 1.3, 3.5, 2.2, fill=EFCAZ_ESCURO)
txt(slide, "[0/10]", 4.9, 1.42, 3.5, 1.2, size=52, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
txt(slide, "[STATUS GERAL DO CLIENTE]", 4.9, 2.72, 3.5, 0.55,
    size=10, bold=True, color=EFCAZ_SEC, align=PP_ALIGN.CENTER)

dimensoes = [
    ("Adoção Operacional",     "[STATUS]", VERDE,    "[Descrição do nível de uso operacional]"),
    ("Engajamento de Módulos", "[STATUS]", VERMELHO, "[Descrição dos módulos em uso ou parados]"),
    ("Capacidade",             "[STATUS]", LARANJA,  "[Descrição da ocupação vs. contrato]"),
    ("Relacionamento",         "[STATUS]", LARANJA,  "[Descrição do relacionamento e proximidade do vencimento]"),
]

for i, (dim, status, cor, desc) in enumerate(dimensoes):
    top = 1.38 + i * 1.42
    side = 0.5 if i % 2 == 0 else 8.38
    rect(slide, side, top, 4.2, 1.25, fill=CINZA_CLARO)
    rect(slide, side, top, 4.2, 0.38, fill=cor)
    txt(slide, dim, side+0.15, top+0.05, 4.0, 0.3, size=12, bold=True, color=BRANCO)
    txt(slide, "● " + status, side+0.15, top+0.45, 1.8, 0.3, size=12, bold=True, color=cor)
    txt(slide, desc, side+0.15, top+0.78, 3.9, 0.42, size=10, color=CINZA_MEDIO, wrap=True)

rect(slide, 0.5, 7.0, 12.33, 0.38, fill=CINZA_CLARO)
txt(slide, "[Insira aqui o diagnóstico síntese: risco de churn, ponto forte, recomendação principal]",
    0.65, 7.06, 12, 0.28, size=11, italic=True, color=EFCAZ_ESCURO, bold=True)

# ─────────────────────────────────────────────
# SLIDE 7 — OPORTUNIDADES
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Oportunidades de Expansão", "")

opps = [
    ("ALTA",   VERMELHO, "[Oportunidade 1]",
     "[Descrição — ex: cliente no teto de fornecedores, crescimento mensal constante]"),
    ("ALTA",   VERMELHO, "[Oportunidade 2]",
     "[Descrição — ex: usuários no limite, qualquer admissão bloqueia]"),
    ("MÉDIA",  LARANJA,  "[Oportunidade 3]",
     "[Descrição — ex: ativação de integrações contratadas e inativas]"),
    ("MÉDIA",  LARANJA,  "[Oportunidade 4]",
     "[Descrição — ex: reengajamento de módulo parado]"),
    ("BAIXA",  EFCAZ_PRIM, "[Oportunidade 5]",
     "[Descrição — ex: módulo não contratado com fit ao perfil do cliente]"),
]

for i, (urg, cor, titulo, desc) in enumerate(opps):
    top = 1.35 + i * 1.18
    rect(slide, 0.5, top, 12.33, 1.05, fill=CINZA_CLARO)
    rect(slide, 0.5, top, 1.5, 1.05, fill=cor)
    txt(slide, urg, 0.5, top+0.32, 1.5, 0.4, size=12, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    txt(slide, titulo, 2.2, top+0.08, 9.8, 0.36, size=13, bold=True, color=PRETO_SUAVE)
    txt(slide, desc, 2.2, top+0.52, 9.8, 0.45, size=11, color=CINZA_MEDIO)

rodape(slide, "Efcaz SRM  •  Análise de Renovação  •  [Nome do Cliente]  •  [Mês/Ano]")

# ─────────────────────────────────────────────
# SLIDE 8 — PLANO DE AÇÃO
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Plano de Ação — [X] Dias para a Renovação", "")

etapas = [
    ("Semana 1\n[data]", VERDE,
     "Agendar QBR / Reunião de Renovação",
     ["[Nome do contato principal — champion identificado]",
      "Pauta: resultados do período + oportunidades + renovação",
      "Duração sugerida: 40 minutos"]),
    ("Semana 2\n[data]", EFCAZ_PRIM,
     "Conduzir Reunião e Enviar Proposta",
     ["Apresentar crescimento e valor gerado no período",
      "Levantar causas de módulos subutilizados",
      "Propor plano de ativação e enviar proposta com opção de upgrade"]),
    ("Semana 3-4\n[data]", LARANJA,
     "Negociação e Assinatura",
     ["Follow-up da proposta (máx. 3 dias sem resposta)",
      "Alinhar jurídico e financeiro do cliente",
      "Assinatura antes do vencimento — garantir continuidade dos bônus"]),
]

for i, (periodo, cor, titulo, bullets) in enumerate(etapas):
    top = 1.35 + i * 1.95
    rect(slide, 0.5, top, 2.0, 1.75, fill=cor)
    txt(slide, periodo, 0.5, top+0.55, 2.0, 0.65,
        size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    rect(slide, 2.65, top, 10.18, 1.75, fill=CINZA_CLARO)
    txt(slide, titulo, 2.8, top+0.08, 10.0, 0.38, size=13, bold=True, color=PRETO_SUAVE)
    for j, b in enumerate(bullets):
        txt(slide, "• " + b, 2.85, top+0.5+j*0.3, 9.9, 0.3, size=10, color=CINZA_MEDIO)

rect(slide, 0.5, 7.05, 12.33, 0.35, fill=EFCAZ_ESCURO)
txt(slide, "Vencimento: [data]  •  Valor mensal: R$ [valor]  •  Contato: [nome] — [email]",
    0.65, 7.1, 12, 0.28, size=11, color=BRANCO)

# ─────────────────────────────────────────────
# SLIDE 9 — ENCERRAMENTO
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=EFCAZ_PRIM)
rect(slide, 0, 3.5, 13.33, 0.06, fill=BRANCO)
rect(slide, 0.5, 1.5, 0.1, 2.0, fill=BRANCO)

txt(slide, "Obrigado", 0.9, 1.5, 11, 1.2, size=60, bold=True, color=BRANCO)
txt(slide, "Gabriel Vital  •  Customer Success Specialist  •  Efcaz SRM",
    0.9, 2.8, 11, 0.5, size=16, color=EFCAZ_SEC)
txt(slide, "vitallgabriel@outlook.com", 0.9, 4.0, 6, 0.4, size=14, color=BRANCO)
txt(slide, "Efcaz SRM  •  efcaz.com.br", 0.9, 4.5, 6, 0.4, size=14, color=EFCAZ_SEC)

output = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\Template_Efcaz_Renovacao.pptx"
prs.save(output)
print(f"Template gerado: {output}")
