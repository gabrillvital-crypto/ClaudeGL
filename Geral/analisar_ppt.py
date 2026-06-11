# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.dml.color import RGBColor

pptx_path = r"C:\Users\gabriel.evangelista\Downloads\Plano de Ação Efcaz _ DockBrasil - modelo.pptx"
prs = Presentation(pptx_path)

print(f"Slides: {len(prs.slides)}")
print(f"Largura: {prs.slide_width.inches:.2f}\"")
print(f"Altura: {prs.slide_height.inches:.2f}\"")
print()

for si, slide in enumerate(prs.slides):
    print(f"=== SLIDE {si+1} ===")
    for shape in slide.shapes:
        info = f"  {shape.name} | tipo={shape.shape_type} | x={shape.left.inches:.2f} y={shape.top.inches:.2f} w={shape.width.inches:.2f} h={shape.height.inches:.2f}"
        try:
            fill = shape.fill
            if fill.type == 1:
                try:
                    info += f" | fill=#{fill.fore_color.rgb}"
                except:
                    pass
        except:
            pass
        print(info)
        if shape.shape_type == 13:
            print(f"    [IMAGEM]")
        try:
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        try:
                            fc = f"#{run.font.color.rgb}"
                        except:
                            fc = "inherited"
                        fn = run.font.name or "inherited"
                        fs = run.font.size
                        print(f"    TXT: \"{run.text[:80]}\" color={fc} font={fn} size={fs} bold={run.font.bold}")
        except:
            pass
