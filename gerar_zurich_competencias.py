"""
Gerador de slides: Plano de Divulgação — Competências Zurich
2 slides | Paleta Efcaz | Calibri
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Paleta ──────────────────────────────────────────────────────────────────
TEAL        = RGBColor(0x0E, 0x8F, 0xA3)
TEAL_ESCURO = RGBColor(0x0A, 0x6A, 0x7A)
TEAL_CLARO  = RGBColor(0xD6, 0xF0, 0xF4)
VERDE       = RGBColor(0x27, 0xAE, 0x60)
VERDE_CLARO = RGBColor(0xD5, 0xF5, 0xE3)
LARANJA     = RGBColor(0xF3, 0x9C, 0x12)
LARANJA_CLARO = RGBColor(0xFE, 0xF3, 0xCD)
VERMELHO    = RGBColor(0xC0, 0x39, 0x2B)
BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLARO = RGBColor(0xF4, 0xF6, 0xF7)
CINZA_MEDIO = RGBColor(0x85, 0x92, 0x9E)
CINZA_ESCURO= RGBColor(0x2C, 0x3E, 0x50)
PRETO_SUAVE = RGBColor(0x17, 0x20, 0x2A)

W = Inches(10)
H = Inches(5.62)

# ── Helpers ──────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, x, y, w, h, text, font_size=10, bold=False,
                color=PRETO_SUAVE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_para(tf, text, font_size=9, bold=False, color=PRETO_SUAVE,
             align=PP_ALIGN.LEFT, space_before=0):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

def header_bar(slide, title, subtitle=None):
    """Barra de cabeçalho teal escuro."""
    bar = add_rect(slide, Inches(0), Inches(0), W, Inches(0.68), TEAL_ESCURO)
    add_textbox(slide, Inches(0.18), Inches(0.08), Inches(7.5), Inches(0.52),
                title, font_size=18, bold=True, color=BRANCO)
    if subtitle:
        add_textbox(slide, Inches(7.7), Inches(0.18), Inches(2.1), Inches(0.38),
                    subtitle, font_size=9, color=RGBColor(0xCC, 0xEE, 0xF2),
                    align=PP_ALIGN.RIGHT)

def footer_bar(slide):
    bar = add_rect(slide, Inches(0), Inches(5.36), W, Inches(0.26), TEAL_ESCURO)
    add_textbox(slide, Inches(0.18), Inches(5.38), Inches(9.5), Inches(0.20),
                "Gabriel Vital  |  Customer Success Specialist  |  Efcaz SRM  |  Zurich Airport Brazil",
                font_size=7, color=RGBColor(0xCC, 0xEE, 0xF2))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — O que muda e Por que agora
# ════════════════════════════════════════════════════════════════════════════
def build_slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ── Header ──
    header_bar(slide,
               "Ativação de Competências — O que Muda na Plataforma",
               "Zurich Airport Brazil | Set/2026")

    # ── Coluna A: O que muda (tabela antes/depois) ──
    col_a_x = Inches(0.18)
    col_a_w = Inches(4.30)
    top = Inches(0.82)

    # Título coluna
    lbl = add_rect(slide, col_a_x, top, col_a_w, Inches(0.32), TEAL)
    lbl.line.fill.background()
    add_textbox(slide, col_a_x + Inches(0.08), top + Inches(0.05),
                col_a_w - Inches(0.16), Inches(0.25),
                "O QUE MUDA NA PLATAFORMA", font_size=9, bold=True, color=BRANCO)

    # Linhas antes/depois
    rows = [
        ("ANTES", "DEPOIS"),
        ("Documentos em listagem única", "Organizados por janela de competência\n(ex: Janeiro, Fevereiro…)"),
        ("Navegação linear", "Abas separadas por período\nVisão mais clara e intuitiva"),
        ("Reclassificação manual dispersa", "Migração automática do histórico\nvia script + Janaína continua na nova estrutura"),
    ]

    row_top = top + Inches(0.35)
    row_h = Inches(0.56)
    col_w = (col_a_w - Inches(0.04)) / 2

    for r_idx, (antes, depois) in enumerate(rows):
        y = row_top + r_idx * (row_h + Inches(0.04))
        is_header = r_idx == 0

        bg_a = TEAL_CLARO if not is_header else RGBColor(0xA8, 0xD8, 0xE0)
        bg_d = CINZA_CLARO if not is_header else RGBColor(0xD5, 0xD8, 0xDC)

        rect_a = add_rect(slide, col_a_x, y, col_w - Inches(0.02), row_h, bg_a)
        rect_d = add_rect(slide, col_a_x + col_w + Inches(0.02), y,
                          col_w - Inches(0.02), row_h, bg_d)

        fs = 8 if not is_header else 8
        bd = is_header
        add_textbox(slide, col_a_x + Inches(0.06), y + Inches(0.06),
                    col_w - Inches(0.14), row_h - Inches(0.1),
                    antes, font_size=fs, bold=bd, color=TEAL_ESCURO if is_header else PRETO_SUAVE)
        add_textbox(slide, col_a_x + col_w + Inches(0.08), y + Inches(0.06),
                    col_w - Inches(0.14), row_h - Inches(0.1),
                    depois, font_size=fs, bold=bd, color=TEAL_ESCURO if is_header else PRETO_SUAVE)

    # ── Coluna B: Escopo desta fase ──
    col_b_x = Inches(4.62)
    col_b_w = Inches(2.40)

    lbl2 = add_rect(slide, col_b_x, top, col_b_w, Inches(0.32), TEAL)
    lbl2.line.fill.background()
    add_textbox(slide, col_b_x + Inches(0.08), top + Inches(0.05),
                col_b_w - Inches(0.12), Inches(0.25),
                "ESCOPO DESTA FASE", font_size=9, bold=True, color=BRANCO)

    # Terceiros — verde
    box_t = add_rect(slide, col_b_x, top + Inches(0.37), col_b_w, Inches(0.88), VERDE_CLARO)
    add_textbox(slide, col_b_x + Inches(0.10), top + Inches(0.42),
                col_b_w - Inches(0.14), Inches(0.25),
                "TERCEIROS (PRESTADORES)", font_size=8, bold=True, color=VERDE)
    add_textbox(slide, col_b_x + Inches(0.10), top + Inches(0.67),
                col_b_w - Inches(0.14), Inches(0.52),
                "Ativado nesta janela\n(set/2026)", font_size=8, color=PRETO_SUAVE)

    # Separador
    add_textbox(slide, col_b_x + Inches(0.10), top + Inches(1.32),
                col_b_w - Inches(0.14), Inches(0.25),
                "FORNECEDORES DE PRODUTO", font_size=8, bold=True,
                color=CINZA_MEDIO)
    box_f = add_rect(slide, col_b_x, top + Inches(1.28), col_b_w, Inches(0.80),
                     CINZA_CLARO)
    add_textbox(slide, col_b_x + Inches(0.10), top + Inches(1.32),
                col_b_w - Inches(0.14), Inches(0.52),
                "Próxima fase\n(a definir)", font_size=8, color=CINZA_MEDIO)

    # Nota técnica
    nota = add_rect(slide, col_b_x, top + Inches(2.18), col_b_w, Inches(0.60),
                    RGBColor(0xFF, 0xF3, 0xCD))
    nota.line.color.rgb = LARANJA
    nota.line.width = Pt(0.75)
    add_textbox(slide, col_b_x + Inches(0.08), top + Inches(2.22),
                col_b_w - Inches(0.12), Inches(0.52),
                "Dados históricos migrados\nautomaticamente via script\n(time Ricardo)",
                font_size=7.5, color=RGBColor(0x7D, 0x60, 0x08))

    # ── Coluna C: Por que agora ──
    col_c_x = Inches(7.16)
    col_c_w = Inches(2.68)

    lbl3 = add_rect(slide, col_c_x, top, col_c_w, Inches(0.32), TEAL)
    lbl3.line.fill.background()
    add_textbox(slide, col_c_x + Inches(0.08), top + Inches(0.05),
                col_c_w - Inches(0.12), Inches(0.25),
                "POR QUE AGORA?", font_size=9, bold=True, color=BRANCO)

    motivos = [
        (VERMELHO,  "Solicitação recorrente da Zurich — ponto crítico no relacionamento"),
        (LARANJA,   "Contrato vence em NOV/2026 — precisamos de 2 janelas limpas antes da renovação"),
        (VERDE,     "Feature pronta e validada pelo time de produto"),
    ]

    y_m = top + Inches(0.38)
    for cor, texto in motivos:
        dot = add_rect(slide, col_c_x + Inches(0.10), y_m + Inches(0.10),
                       Inches(0.14), Inches(0.14), cor)
        add_textbox(slide, col_c_x + Inches(0.30), y_m,
                    col_c_w - Inches(0.36), Inches(0.80),
                    texto, font_size=8.5, color=PRETO_SUAVE)
        y_m += Inches(0.88)

    footer_bar(slide)
    return slide


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Fluxo de Divulgação / Linha do Tempo
# ════════════════════════════════════════════════════════════════════════════
def build_slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    header_bar(slide,
               "Plano de Comunicação — Linha do Tempo",
               "Ago–Set 2026")

    # ── 5 etapas em caixas horizontais ──────────────────────────────────────
    etapas = [
        {
            "num": "01",
            "titulo": "Alinhamento\ncom Zurich",
            "data": "20–21/ago",
            "resp": "Gabriel + Ricardo",
            "canal": "Reunião dedicada",
            "acao": "Apresentar plano + demo em trial da nova visualização",
            "cor": TEAL_ESCURO,
        },
        {
            "num": "02",
            "titulo": "Briefing\nSuporte",
            "data": "~20/ago",
            "resp": "Gabriel + Thais",
            "canal": "Reunião interna",
            "acao": "Alinhar time de suporte (Jamil/Carlos) sobre o que muda",
            "cor": TEAL,
        },
        {
            "num": "03",
            "titulo": "Comunicação\nFornecedores",
            "data": "24–28/ago",
            "resp": "CS Efcaz",
            "canal": "E-mail marketing + vídeo",
            "acao": "Enviar aviso + link do vídeo passo a passo + link webinar",
            "cor": LARANJA,
        },
        {
            "num": "04",
            "titulo": "Virada do\nSistema",
            "data": "28/ago — EOD",
            "resp": "Time Produto (Ricardo)",
            "canal": "Script de migração",
            "acao": "Ativar competências em produção — 31/ago: novo formato ao vivo",
            "cor": VERMELHO,
        },
        {
            "num": "05",
            "titulo": "Suporte\nIntensivo",
            "data": "31/ago – 10/set",
            "resp": "Thais + CS",
            "canal": "Webinar + Plantão",
            "acao": "Horários abertos para dúvidas — 3 sessões programadas",
            "cor": VERDE,
        },
    ]

    n = len(etapas)
    margin_x = Inches(0.18)
    margin_top = Inches(0.80)
    box_gap = Inches(0.06)
    arrow_w = Inches(0.16)
    total_w = W - 2 * margin_x
    box_w = (total_w - (n - 1) * (box_gap + arrow_w)) / n
    box_h = Inches(2.62)

    for i, e in enumerate(etapas):
        x = margin_x + i * (box_w + box_gap + arrow_w)

        # Fundo da caixa
        box = add_rect(slide, x, margin_top, box_w, box_h, CINZA_CLARO)
        box.line.color.rgb = e["cor"]
        box.line.width = Pt(1.5)

        # Cabeçalho colorido
        header = add_rect(slide, x, margin_top, box_w, Inches(0.38), e["cor"])

        # Número + título no header
        add_textbox(slide, x + Inches(0.06), margin_top + Inches(0.03),
                    box_w - Inches(0.10), Inches(0.32),
                    e["num"], font_size=10, bold=True, color=BRANCO)

        # Título (2 linhas)
        add_textbox(slide, x + Inches(0.06), margin_top + Inches(0.38),
                    box_w - Inches(0.10), Inches(0.52),
                    e["titulo"], font_size=9, bold=True, color=e["cor"])

        # Data
        add_textbox(slide, x + Inches(0.06), margin_top + Inches(0.92),
                    box_w - Inches(0.10), Inches(0.22),
                    e["data"], font_size=8, bold=True, color=PRETO_SUAVE)

        # Responsável
        add_textbox(slide, x + Inches(0.06), margin_top + Inches(1.14),
                    box_w - Inches(0.10), Inches(0.20),
                    e["resp"], font_size=7.5, color=CINZA_MEDIO)

        # Linha separadora fina
        sep = add_rect(slide, x + Inches(0.06), margin_top + Inches(1.36),
                       box_w - Inches(0.12), Inches(0.02), e["cor"])

        # Canal
        add_textbox(slide, x + Inches(0.06), margin_top + Inches(1.42),
                    box_w - Inches(0.10), Inches(0.20),
                    e["canal"], font_size=7.5, bold=True, color=e["cor"])

        # Ação
        add_textbox(slide, x + Inches(0.06), margin_top + Inches(1.62),
                    box_w - Inches(0.10), Inches(0.88),
                    e["acao"], font_size=7.5, color=PRETO_SUAVE)

        # Seta entre boxes (exceto última)
        if i < n - 1:
            ax = x + box_w + box_gap / 2
            ay = margin_top + box_h / 2 - Inches(0.10)
            add_textbox(slide, ax, ay, arrow_w, Inches(0.22),
                        "→", font_size=14, bold=True, color=TEAL,
                        align=PP_ALIGN.CENTER)

    # ── Tabela de webinars ──────────────────────────────────────────────────
    tbl_top = margin_top + box_h + Inches(0.14)
    tbl_h = Inches(0.62)

    # Fundo tabela
    tbl_bg = add_rect(slide, margin_x, tbl_top,
                      W - 2 * margin_x, tbl_h, RGBColor(0xEA, 0xF6, 0xF8))
    tbl_bg.line.color.rgb = TEAL
    tbl_bg.line.width = Pt(0.75)

    # Header da tabela
    add_textbox(slide, margin_x + Inches(0.10), tbl_top + Inches(0.04),
                Inches(2.0), Inches(0.22),
                "SUPORTE POS-VIRADA (31/ago – 10/set)", font_size=7.5,
                bold=True, color=TEAL_ESCURO)

    sessoes = [
        ("28/ago  Webinar   11:30–12:30   Todos os terceiros"),
        ("04/set  Webinar   11:30–12:30   Todos os terceiros"),
        ("10/set  Webinar   11:30–12:30   + 3 slots tarde"),
        ("Sextas  Plantão   14:00–15:00   Dúvidas abertas por link"),
        ("Quartas  Exclusivo  A definir   Por CNPJ (individual)"),
    ]

    col_w2 = (W - 2 * margin_x - Inches(0.20)) / len(sessoes)
    for j, s in enumerate(sessoes):
        cx = margin_x + Inches(0.10) + j * col_w2
        partes = s.split("  ")
        add_textbox(slide, cx, tbl_top + Inches(0.28),
                    col_w2 - Inches(0.04), Inches(0.32),
                    partes[0], font_size=7.5, bold=True, color=TEAL_ESCURO)
        if len(partes) > 1:
            resto = "  ".join(partes[1:])
            add_textbox(slide, cx, tbl_top + Inches(0.38),
                        col_w2 - Inches(0.04), Inches(0.22),
                        resto, font_size=7, color=CINZA_ESCURO)

        # Divisor vertical (exceto último)
        if j < len(sessoes) - 1:
            div_x = cx + col_w2 - Inches(0.04)
            add_rect(slide, div_x, tbl_top + Inches(0.10),
                     Inches(0.01), tbl_h - Inches(0.20), TEAL)

    footer_bar(slide)
    return slide


# ════════════════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

build_slide1(prs)
build_slide2(prs)

output = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\Zurich_Competencias_Ago2026.pptx"
prs.save(output)
print(f"Salvo: {output}")
