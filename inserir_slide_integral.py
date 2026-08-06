import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
from pptx import Presentation
import copy
import zipfile
import os

source_path = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\Documentos\Cópia de Proposta Comercial  Efcaz SRM - Transportes Cavalinho v.2.pptx"
target_path = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\Renovacoes\Integral_Medica\Integral_Medica_Renovacao_Jun2026_v2.pptx"
output_path = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\Renovacoes\Integral_Medica\Integral_Medica_Renovacao_Jun2026_v3.pptx"
tmp_img_path = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\_temp_slide_img.png"

prs_source = Presentation(source_path)
prs_target = Presentation(target_path)

source_slide = prs_source.slides[0]

# Extrai imagem do PPTX source diretamente
with zipfile.ZipFile(source_path, "r") as z:
    img_bytes = z.read("ppt/media/image2.png")
with open(tmp_img_path, "wb") as f:
    f.write(img_bytes)
print(f"Imagem extraída: {len(img_bytes)} bytes")

# Localiza shape da imagem para obter posição e tamanho originais
img_shape_info = None
for shape in source_slide.shapes:
    if shape.shape_type == 13:  # PICTURE
        img_shape_info = (shape.left, shape.top, shape.width, shape.height)
        print(f"Shape imagem: left={shape.left}, top={shape.top}, w={shape.width}, h={shape.height}")
        break

# Adiciona slide em branco no final
blank_layout = prs_target.slide_layouts[6]
new_slide = prs_target.slides.add_slide(blank_layout)

# Remove shapes padrão
spTree = new_slide.shapes._spTree
for child in list(spTree)[2:]:
    spTree.remove(child)

# Copia shapes de texto (não-imagem) preservando formatação
for shape in source_slide.shapes:
    if shape.shape_type != 13:  # Pula PICTURE
        el = copy.deepcopy(shape.element)
        spTree.append(el)
        print(f"Shape copiado: {shape.name} ({shape.shape_type})")

# Insere imagem com posição e tamanho originais
if img_shape_info:
    left, top, width, height = img_shape_info
    new_slide.shapes.add_picture(tmp_img_path, left, top, width, height)
    print(f"Imagem inserida via add_picture.")

# Copia background do slide de origem
NS_PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
bg = source_slide._element.find(f"{{{NS_PML}}}bg")
if bg is not None:
    existing_bg = new_slide._element.find(f"{{{NS_PML}}}bg")
    if existing_bg is not None:
        new_slide._element.remove(existing_bg)
    new_slide._element.insert(2, copy.deepcopy(bg))
    print("Background copiado.")

# Move slide para posição 5 (entre slides 5 e 6)
sldIdLst = prs_target.slides._sldIdLst
ultimo = sldIdLst[-1]
sldIdLst.remove(ultimo)
sldIdLst.insert(5, ultimo)

prs_target.save(output_path)

# Remove imagem temporária
os.remove(tmp_img_path)

print(f"\nSlides no arquivo final: {len(prs_target.slides)}")
print(f"Salvo em: {output_path}")
