import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
from pptx import Presentation

path = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\Documentos\Cópia de Proposta Comercial  Efcaz SRM - Transportes Cavalinho v.2.pptx"
prs = Presentation(path)
slide = prs.slides[0]

print("=== Shapes do slide ===")
for shape in slide.shapes:
    tem_texto = shape.has_text_frame
    texto = shape.text_frame.text[:80] if tem_texto else ""
    print(f"  {shape.name} | shape_type={shape.shape_type} | texto={texto}")

print("\n=== Relacionamentos ===")
for rel in slide.part.rels.values():
    tipo_rel = rel.reltype.split("/")[-1]
    print(f"  {rel.rId} | {tipo_rel} | {rel.target_ref}")

print("\n=== Arquivos de mídia no PPTX ===")
import zipfile
with zipfile.ZipFile(path, "r") as z:
    for f in z.namelist():
        if "media" in f or "image" in f.lower():
            print(f"  {f}")
