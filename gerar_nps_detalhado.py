#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path

# Paleta Efcaz
TEAL = RGBColor(0x0E, 0x8F, 0xA3)
TEAL_CLARO = RGBColor(0x14, 0xB3, 0xCC)
VERDE = RGBColor(0x27, 0xAE, 0x60)
VERMELHO = RGBColor(0xE7, 0x4C, 0x3C)
LARANJA = RGBColor(0xF3, 0x9C, 0x12)
CINZA_CLARO = RGBColor(0xF2, 0xF4, 0xF4)
PRETO = RGBColor(0x17, 0x20, 0x2A)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)

# Criar apresentação
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_slide_blank():
    """Adiciona slide em branco"""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def add_title_shape(slide, text, top, left, width, height, size=54, color=TEAL, bold=True):
    """Adiciona título formatado"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Calibri'

    return shape

def add_subtitle_shape(slide, text, top, left, width, height, size=24, color=PRETO, bold=False):
    """Adiciona subtítulo"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = shape.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Calibri'

    return shape

def add_metric_box(slide, title, value, top, left, width=2.2, height=1.8, color_bg=CINZA_CLARO, color_text=TEAL):
    """Adiciona box de métrica"""
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color_bg
    shape.line.color.rgb = color_text
    shape.line.width = Pt(2)

    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRETO
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = value
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color_text
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(10)

def add_bullet_point(slide, text, top, left, width=9, size=16, color=PRETO, bold=False):
    """Adiciona ponto de lista"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.5))
    text_frame = shape.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "• " + text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Calibri'

print("Criando PPTX detalhado de NPS...\n")

# ====== SLIDE 1: CAPA ======
print("Slide 1: Capa")
slide1 = add_slide_blank()
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BRANCO

add_title_shape(slide1, "NPS Carteira Efcaz", 2, 0.5, 9, 1.5, size=60, color=TEAL)
add_subtitle_shape(slide1, "Análise Detalhada de Engajamento e Satisfação", 3.7, 0.5, 9, 0.6, size=28, color=PRETO)
add_subtitle_shape(slide1, "Junho 2026 | Gabriel Vital • Customer Success Specialist", 5.5, 0.5, 9, 0.5, size=16, color=CINZA_CLARO)

# ====== SLIDE 2: RESUMO EXECUTIVO ======
print("Slide 2: Resumo Executivo")
slide2 = add_slide_blank()
background = slide2.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BRANCO

add_title_shape(slide2, "Resumo Executivo", 0.5, 0.5, 9, 0.7, size=44)

add_metric_box(slide2, "Respondidos", "15\n(60%)", 1.5, 0.8, color_bg=CINZA_CLARO, color_text=VERDE)
add_metric_box(slide2, "Sem Resposta", "8\n(32%)", 1.5, 3.3, color_bg=CINZA_CLARO, color_text=LARANJA)
add_metric_box(slide2, "Sem Acesso", "2\n(8%)", 1.5, 5.8, color_bg=CINZA_CLARO, color_text=VERMELHO)

add_subtitle_shape(slide2, "25 clientes avaliados | NPS Médio +73 | Taxa resposta 22.6% (37 de 164 acessos)", 4, 0.8, 8.4, 0.5, size=13, color=PRETO)
add_subtitle_shape(slide2, "Detratores: 2 | Passivos: 1 | Promoters: 12", 4.65, 0.8, 8.4, 0.5, size=13, color=PRETO, bold=True)

# ====== SLIDE 3: SEGMENTAÇÃO NPS COMPLETA ======
print("Slide 3: Segmentação NPS Completa")
slide3 = add_slide_blank()
background = slide3.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BRANCO

add_title_shape(slide3, "Segmentação NPS — Respondentes Únicos", 0.5, 0.5, 9, 0.7, size=44)

# Promoters
add_subtitle_shape(slide3, "🟢 PROMOTERS (9-10) — 12 respondentes | NPS +100", 1.4, 0.8, 9, 0.35, size=16, color=VERDE, bold=True)
add_bullet_point(slide3, "Unimed Sou, Geistlich (2), Engesp (2), Hospital Adventista (4), Unimed Campo Grande (2), Dock Brasil (3), DOF, Ponsse, ZAB, Unimed Brasil", 1.85, 1.2, width=8.2, size=13)

# Passivos
add_subtitle_shape(slide3, "🟡 PASSIVOS (7-8) — 1 respondente | Feedback técnico", 2.9, 0.8, 9, 0.35, size=16, color=LARANJA, bold=True)
add_bullet_point(slide3, "Ida Laiber (Integral Médica, nota 7): 'Conexão cai — tenho que reabrir em curto período'", 3.35, 1.2, width=8.2, size=13)

# Detratores
add_subtitle_shape(slide3, "🔴 DETRATORES (0-6) — 2 respondentes | Feedback específico", 4.4, 0.8, 9, 0.35, size=16, color=VERMELHO, bold=True)
add_bullet_point(slide3, "Thais Barroso (Afonso França, nota 6): 'Navegabilidade, suporte, praticidade operacional'", 4.85, 1.2, width=8.2, size=13)
add_bullet_point(slide3, "Bruno Machado (Cebrace, nota 1): 'Não faz sentido avaliar mesmo fornecedor em período curto' (feedback sobre frequência, não produto)", 5.5, 1.2, width=8.2, size=13)

# Insight
add_subtitle_shape(slide3, "⚠ Insight Crítico", 6.3, 0.8, 9, 0.3, size=14, color=PRETO, bold=True)
add_bullet_point(slide3, "Detratores não questionam qualidade do produto — feedback é sobre UX e frequência de uso", 6.65, 1.2, width=8.2, size=13, color=PRETO)

# ====== SLIDE 4: DISTRIBUIÇÃO POR STATUS ======
print("Slide 4: Distribuição por Status")
slide4 = add_slide_blank()
background = slide4.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BRANCO

add_title_shape(slide4, "Distribuição por Status — Visão da Carteira (25)", 0.5, 0.5, 9, 0.7, size=44)

add_bullet_point(slide4, "15 clientes (60%) — Com resposta: Segmentados em 12 promoters + 1 passivo + 2 detratores", 1.5, 0.8, width=8.4, size=14)
add_bullet_point(slide4, "8 clientes (32%) — Acessam mas não responderam: Heavy users desengajados", 2.15, 0.8, width=8.4, size=14)
add_bullet_point(slide4, "2 clientes (8%) — Sem nenhum acesso: ADVtec + Pacco (verificar credenciais)", 2.8, 0.8, width=8.4, size=14)

add_subtitle_shape(slide4, "Métricas-chave", 3.6, 0.8, 8.4, 0.35, size=14, color=PRETO, bold=True)
add_bullet_point(slide4, "Taxa de resposta geral: 22.6% (37 de 164 usuários que acessaram)", 4.05, 1.2, width=8.2, size=13)
add_bullet_point(slide4, "Problema: Engajamento com pesquisa, não satisfação", 4.55, 1.2, width=8.2, size=13)

add_subtitle_shape(slide4, "Heavy Users Não Respondentes", 5.25, 0.8, 8.4, 0.35, size=14, color=PRETO, bold=True)
add_bullet_point(slide4, "Integral Médica: 9 users com 5+ views, Federação Paulista: Filipe com 10 views, Tarkett: 5 users com 5+ views", 5.7, 1.2, width=8.2, size=12)

# ====== SLIDE 5: CRÍTICOS POR TIER ======
print("Slide 5: Críticos por Tier")
slide5 = add_slide_blank()
background = slide5.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BRANCO

add_title_shape(slide5, "Análise por Tier — Distribuição de Problemas", 0.5, 0.5, 9, 0.7, size=44)

add_subtitle_shape(slide5, "Tier A", 1.4, 0.8, 2, 0.3, size=13, color=PRETO, bold=True)
add_bullet_point(slide5, "2 sem resposta (Bom Futuro, Eucatex — transferido a Wagner)", 1.8, 1.2, width=7.8, size=12)

add_subtitle_shape(slide5, "Tier B+", 2.4, 0.8, 2, 0.3, size=13, color=PRETO, bold=True)
add_bullet_point(slide5, "1 sem resposta: Bunker One (6 acessos, desengajado apesar de Ongoing)", 2.8, 1.2, width=7.8, size=12)

add_subtitle_shape(slide5, "Tier B", 3.5, 0.8, 2, 0.3, size=13, color=VERMELHO, bold=True)
add_bullet_point(slide5, "4 críticos: 1 sem acesso (Alumetaf) + 3 sem resposta", 3.9, 1.2, width=7.8, size=12)
add_bullet_point(slide5, "Destaque: Integral Médica (17% resposta, renovação 30/06), Federação Paulista (0% resposta, Filipe 10 views)", 4.45, 1.2, width=7.8, size=12)

add_subtitle_shape(slide5, "Tier C", 5.35, 0.8, 2, 0.3, size=13, color=PRETO, bold=True)
add_bullet_point(slide5, "3 críticos: 1 sem acesso (ADVtec) + 2 sem resposta (Asso Marítima, Amboretto)", 5.75, 1.2, width=7.8, size=12)

# ====== SLIDE 6: RESUMO FINAL ======
print("Slide 6: Resumo Final")
slide6 = add_slide_blank()
background = slide6.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BRANCO

add_title_shape(slide6, "Resumo — Situação Geral", 0.5, 0.5, 9, 0.7, size=44)

add_subtitle_shape(slide6, "✓ Positivo", 1.4, 0.8, 9, 0.3, size=14, color=VERDE, bold=True)
add_bullet_point(slide6, "NPS médio +73 — 12 promoters (80% dos respondentes)", 1.85, 1.2, width=8.2, size=12)
add_bullet_point(slide6, "Detratores não questionam qualidade do produto — feedback é operacional", 2.3, 1.2, width=8.2, size=12)

add_subtitle_shape(slide6, "⚠ Crítico", 3, 0.8, 9, 0.3, size=14, color=VERMELHO, bold=True)
add_bullet_point(slide6, "10 clientes com problema: 2 sem acesso + 8 acessando mas não respondendo", 3.45, 1.2, width=8.2, size=12)
add_bullet_point(slide6, "Tier B concentra maioria (4 de 12) — Integral Médica urgente (renovação 30/06)", 3.9, 1.2, width=8.2, size=12)
add_bullet_point(slide6, "Heavy users consistentemente não respondendo — padrão de desengajamento com pesquisa", 4.35, 1.2, width=8.2, size=12)

add_subtitle_shape(slide6, "→ Próxima Onda", 5.1, 0.8, 9, 0.3, size=14, color=TEAL, bold=True)
add_bullet_point(slide6, "Simplificar pesquisa (2-3 perguntas) para aumentar taxa de 22.6% — alvo: 35%+", 5.55, 1.2, width=8.2, size=12)

# ====== SALVAR ======
output_dir = Path(r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\clientes\NPS")
output_dir.mkdir(parents=True, exist_ok=True)
output_path = output_dir / "NPS_Carteira_Efcaz_2026_06_19.pptx"

prs.save(str(output_path))
print(f"\n✅ PPTX detalhado salvo em: {output_path}")
print(f"   6 slides com segmentação NPS completa (Detratores + Passivos + Promoters + Feedback)")
