from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Paleta de cores
AZUL_ESCURO   = RGBColor(0x1B, 0x3A, 0x5C)
AZUL_MEDIO    = RGBColor(0x2E, 0x6D, 0xA4)
VERDE         = RGBColor(0x27, 0xAE, 0x60)
VERDE_CLARO   = RGBColor(0xD5, 0xF5, 0xE3)
VERMELHO      = RGBColor(0xE7, 0x4C, 0x3C)
VERMELHO_CLARO= RGBColor(0xFD, 0xED, 0xEC)
LARANJA       = RGBColor(0xF3, 0x9C, 0x12)
LARANJA_CLARO = RGBColor(0xFE, 0xF9, 0xE7)
BRANCO        = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLARO   = RGBColor(0xF2, 0xF4, 0xF4)
CINZA_MEDIO   = RGBColor(0x85, 0x92, 0x9E)
PRETO_SUAVE   = RGBColor(0x17, 0x20, 0x2A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # layout em branco

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height, font_size=14, bold=False,
             color=PRETO_SUAVE, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def add_multiline(slide, lines, left, top, width, height, font_size=13,
                  color=PRETO_SUAVE, bold=False, line_spacing=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox

def kpi_card(slide, left, top, width, height, number, label, color=AZUL_ESCURO):
    add_rect(slide, left, top, width, height, fill_color=color)
    add_text(slide, number, left+0.1, top+0.12, width-0.2, height*0.5,
             font_size=32, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_text(slide, label, left+0.05, top+height*0.55, width-0.1, height*0.4,
             font_size=11, color=BRANCO, align=PP_ALIGN.CENTER, wrap=True)

def status_row(slide, left, top, width, label, status, status_color, desc, row_color=CINZA_CLARO):
    add_rect(slide, left, top, width, 0.42, fill_color=row_color)
    add_text(slide, label, left+0.15, top+0.06, width*0.35, 0.32, font_size=12, bold=True, color=PRETO_SUAVE)
    dot = "●"
    add_text(slide, dot + " " + status, left+width*0.38, top+0.06, width*0.22, 0.32, font_size=12, bold=True, color=status_color)
    add_text(slide, desc, left+width*0.62, top+0.06, width*0.36, 0.32, font_size=11, color=CINZA_MEDIO)

# ─────────────────────────────────────────────
# SLIDE 1 — CAPA
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=AZUL_ESCURO)
add_rect(slide, 0, 5.2, 13.33, 2.3, fill_color=AZUL_MEDIO)
add_rect(slide, 0.5, 2.8, 0.08, 1.5, fill_color=LARANJA)

add_text(slide, "DOCK BRASIL", 0.8, 1.6, 11, 1.0,
         font_size=48, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)
add_text(slide, "Engenharia e Serviços S.A.", 0.8, 2.65, 9, 0.6,
         font_size=22, color=RGBColor(0xAE, 0xD6, 0xF1), align=PP_ALIGN.LEFT)
add_text(slide, "Análise de Renovação  |  Maio 2026", 0.8, 5.5, 9, 0.55,
         font_size=18, color=BRANCO, align=PP_ALIGN.LEFT)
add_text(slide, "Efcaz SRM  •  Customer Success", 0.8, 6.15, 9, 0.45,
         font_size=14, color=RGBColor(0xAE, 0xD6, 0xF1), align=PP_ALIGN.LEFT)
add_text(slide, "Confidencial", 10.5, 6.9, 2.5, 0.4,
         font_size=10, color=RGBColor(0x85, 0x92, 0x9E), align=PP_ALIGN.RIGHT, italic=True)

# ─────────────────────────────────────────────
# SLIDE 2 — AGENDA
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 1.1, fill_color=AZUL_ESCURO)
add_text(slide, "Agenda", 0.5, 0.18, 12, 0.75,
         font_size=28, bold=True, color=BRANCO)

items = [
    ("01", "Uso e Adoção da Plataforma",   "KPIs de uso, fornecedores e terceiros"),
    ("02", "Crescimento no Período",        "Evolução de solicitações e terceiros cadastrados"),
    ("03", "Pontos de Alerta",              "Módulos subutilizados e riscos identificados"),
    ("04", "Health Score",                  "Avaliação geral da saúde do cliente"),
    ("05", "Oportunidades de Expansão",     "Capacidade, módulos e integrações"),
    ("06", "Próximos Passos",               "Plano de ação para renovação"),
]

for i, (num, title, desc) in enumerate(items):
    row = 1.35 + i * 0.95
    bg = CINZA_CLARO if i % 2 == 0 else BRANCO
    add_rect(slide, 0.5, row, 12.33, 0.82, fill_color=bg)
    add_rect(slide, 0.5, row, 0.7, 0.82, fill_color=AZUL_ESCURO)
    add_text(slide, num, 0.52, row+0.18, 0.65, 0.5, font_size=16, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_text(slide, title, 1.35, row+0.08, 5.5, 0.38, font_size=15, bold=True, color=AZUL_ESCURO)
    add_text(slide, desc,  1.35, row+0.45, 5.5, 0.32, font_size=12, color=CINZA_MEDIO)

# ─────────────────────────────────────────────
# SLIDE 3 — KPIs DE USO
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 1.1, fill_color=AZUL_ESCURO)
add_text(slide, "Uso e Adoção da Plataforma", 0.5, 0.18, 12, 0.75,
         font_size=28, bold=True, color=BRANCO)
add_text(slide, "Dados extraídos do BI de usabilidade  •  Maio 2026", 0.5, 1.25, 12, 0.4,
         font_size=12, color=CINZA_MEDIO)

kpis_top = [
    ("763",   "Fornecedores\nAtivos",       AZUL_ESCURO),
    ("4.620", "Terceiros\nCadastrados",     AZUL_MEDIO),
    ("3.919", "Solicitações\nAbertas",      RGBColor(0x16, 0x7A, 0xC6)),
    ("89,9%", "Taxa de\nAprovação",         VERDE),
]

card_w, card_h = 2.8, 1.7
start_x = 0.55
for j, (num, label, col) in enumerate(kpis_top):
    kpi_card(slide, start_x + j*(card_w+0.25), 1.78, card_w, card_h, num, label, col)

kpis_bot = [
    ("24",        "Usuários\nCadastrados",    LARANJA),
    ("800",       "Fornecedores\nContratados",CINZA_MEDIO),
    ("5",         "Módulos\nAtivos",          AZUL_MEDIO),
    ("265 dias",  "Média de\nHomologação",    RGBColor(0x7F, 0x8C, 0x8D)),
]
for j, (num, label, col) in enumerate(kpis_bot):
    kpi_card(slide, start_x + j*(card_w+0.25), 3.72, card_w, card_h, num, label, col)

add_rect(slide, 0.5, 5.55, 12.33, 0.5, fill_color=RGBColor(0xFE, 0xF9, 0xE7))
add_text(slide,
         "⚠  Usuários no limite: 24 cadastrados vs. 23 permitidos em contrato (15 + 8 auxiliares). Qualquer novo colaborador exige expansão.",
         0.65, 5.6, 12, 0.38, font_size=11, color=LARANJA, bold=True)

# ─────────────────────────────────────────────
# SLIDE 4 — CRESCIMENTO
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 1.1, fill_color=AZUL_ESCURO)
add_text(slide, "Crescimento no Período", 0.5, 0.18, 12, 0.75,
         font_size=28, bold=True, color=BRANCO)

# Solicitações
add_text(slide, "Solicitações por Mês (Jul/2024 → Abr/2026)", 0.5, 1.25, 8, 0.4,
         font_size=14, bold=True, color=AZUL_ESCURO)

sol_data = [
    ("Jul/24","11"), ("Out/24","109"), ("Jan/25","80"),
    ("Abr/25","86"), ("Jul/25","292"), ("Out/25","204"),
    ("Jan/26","189"), ("Mar/26","377"),
]
max_val = 377
bar_area_h = 2.2
bar_area_top = 1.75
bar_w = 1.38
gap = 0.05
for i, (mes, val) in enumerate(sol_data):
    v = int(val)
    ratio = v / max_val
    bh = bar_area_h * ratio
    bx = 0.5 + i * (bar_w + gap)
    by = bar_area_top + bar_area_h - bh
    col = AZUL_MEDIO if i < 6 else VERDE
    add_rect(slide, bx, by, bar_w, bh, fill_color=col)
    add_text(slide, val, bx, by - 0.28, bar_w, 0.28, font_size=10, bold=True, color=AZUL_ESCURO, align=PP_ALIGN.CENTER)
    add_text(slide, mes,  bx, bar_area_top+bar_area_h+0.05, bar_w, 0.25, font_size=9, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

add_rect(slide, 0.5, 4.25, 12.0, 0.04, fill_color=CINZA_MEDIO)

# Terceiros
add_text(slide, "Terceiros Cadastrados por Mês (Out/2024 → Abr/2026)", 0.5, 4.38, 9, 0.38,
         font_size=14, bold=True, color=AZUL_ESCURO)

terc_data = [
    ("Out/24","177"), ("Jan/25","253"), ("Abr/25","268"),
    ("Jul/25","293"), ("Out/25","363"), ("Jan/26","348"),
    ("Mar/26","355"), ("Abr/26","436"),
]
bar_area_top2 = 4.85
bar_area_h2 = 1.9
for i, (mes, val) in enumerate(terc_data):
    v = int(val)
    ratio = v / 436
    bh = bar_area_h2 * ratio
    bx = 0.5 + i * (bar_w + gap)
    by = bar_area_top2 + bar_area_h2 - bh
    col = VERDE if i == 7 else AZUL_MEDIO
    add_rect(slide, bx, by, bar_w, bh, fill_color=col)
    add_text(slide, val, bx, by - 0.27, bar_w, 0.27, font_size=9, bold=True, color=AZUL_ESCURO, align=PP_ALIGN.CENTER)
    add_text(slide, mes,  bx, bar_area_top2+bar_area_h2+0.05, bar_w, 0.22, font_size=8, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

add_text(slide, "34x mais solicitações em 21 meses  •  Crescimento constante de terceiros — cliente em plena expansão operacional",
         0.5, 7.1, 12, 0.3, font_size=11, italic=True, color=CINZA_MEDIO)

# ─────────────────────────────────────────────
# SLIDE 5 — PONTOS DE ALERTA
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 1.1, fill_color=AZUL_ESCURO)
add_text(slide, "Pontos de Alerta — Módulos Subutilizados", 0.5, 0.18, 12, 0.75,
         font_size=28, bold=True, color=BRANCO)

alertas = [
    (VERMELHO, "🔴 Módulo de Avaliação (RFI) — Parado desde Set/2025",
     "Pico de uso em Jul/2025 (170 avaliações). De setembro/2025 a abril/2026: praticamente zero.\n"
     "Pergunta-chave: o que mudou internamente? Troca de responsável? Processo pausado?"),
    (VERMELHO, "🔴 Buscas Automáticas sem Validade — Renovação Interrompida",
     "4 documentos com Background Check ativo (CNDT, CND, CNPJ, FGTS), mas sem validade configurada. O sistema buscou uma vez e parou — as certidões não são renovadas automaticamente.\n"
     "Resultado: 53 fornecedores com documentos vencidos. Configurar a validade nesses documentos resolve o ciclo completo de automação."),
    (LARANJA,  "🟡 Usuários no Limite do Contrato",
     "24 usuários cadastrados vs. 23 permitidos (15 + 8 aux.). Qualquer admissão exige expansão.\n"
     "Sinalizar proativamente evita bloqueio operacional do cliente."),
]

for i, (cor, titulo, desc) in enumerate(alertas):
    top = 1.35 + i * 1.95
    add_rect(slide, 0.5, top, 12.33, 1.6, fill_color=CINZA_CLARO)
    add_rect(slide, 0.5, top, 0.18, 1.6, fill_color=cor)
    add_text(slide, titulo, 0.85, top+0.12, 11.5, 0.38, font_size=13, bold=True, color=PRETO_SUAVE)
    add_multiline(slide, desc.split("\n"), 0.85, top+0.58, 11.5, 0.88, font_size=11, color=CINZA_MEDIO)

# ─────────────────────────────────────────────
# SLIDE 6 — HEALTH SCORE
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 1.1, fill_color=AZUL_ESCURO)
add_text(slide, "Health Score — Dock Brasil", 0.5, 0.18, 12, 0.75,
         font_size=28, bold=True, color=BRANCO)

# Score central
add_rect(slide, 4.9, 1.3, 3.5, 2.2, fill_color=AZUL_ESCURO)
add_text(slide, "6 / 10", 4.9, 1.45, 3.5, 1.2, font_size=52, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
add_text(slide, "ESTÁVEL COM RISCOS LATENTES", 4.9, 2.75, 3.5, 0.55,
         font_size=10, bold=True, color=LARANJA, align=PP_ALIGN.CENTER)

dimensoes = [
    ("Adoção Operacional",     "ALTO",      VERDE,    "Solicitações, fornecedores e terceiros em crescimento acelerado"),
    ("Engajamento de Módulos", "CRÍTICO",   VERMELHO, "Avaliação parada há 8 meses; 53 fornecedores com documentos vencidos"),
    ("Capacidade",             "ATENÇÃO",   LARANJA,  "Fornecedores (95%) e usuários no teto do contrato"),
    ("Relacionamento",         "ATENÇÃO",   LARANJA,  "28 dias para vencer; necessidade de QBR imediato"),
]

for i, (dim, status, cor, desc) in enumerate(dimensoes):
    top = 1.38 + i * 1.42
    side = 0.5 if i % 2 == 0 else 8.38  # esquerda ou direita
    add_rect(slide, side, top, 4.2, 1.25, fill_color=CINZA_CLARO)
    add_rect(slide, side, top, 4.2, 0.38, fill_color=cor)
    add_text(slide, dim, side+0.15, top+0.05, 4.0, 0.3, font_size=12, bold=True, color=BRANCO)
    add_text(slide, "● " + status, side+0.15, top+0.45, 1.5, 0.3, font_size=12, bold=True, color=cor)
    add_text(slide, desc, side+0.15, top+0.78, 3.9, 0.42, font_size=10, color=CINZA_MEDIO, wrap=True)

add_rect(slide, 0.5, 7.0, 12.33, 0.38, fill_color=CINZA_CLARO)
add_text(slide, "Risco de churn não é por insatisfação — é por falta de engajamento estratégico recente. A renovação tem argumentos sólidos.",
         0.65, 7.06, 12, 0.28, font_size=11, italic=True, color=AZUL_ESCURO, bold=True)

# ─────────────────────────────────────────────
# SLIDE 7 — OPORTUNIDADES
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 1.1, fill_color=AZUL_ESCURO)
add_text(slide, "Oportunidades de Expansão", 0.5, 0.18, 12, 0.75,
         font_size=28, bold=True, color=BRANCO)

opps = [
    ("ALTA",   VERMELHO, "Expansão de Fornecedores",
     "95,4% de ocupação (763/800). Crescimento mensal constante indica que o teto será atingido em breve."),
    ("ALTA",   VERMELHO, "Expansão de Usuários",
     "24 de 23 contratados. Qualquer novo usuário bloqueia o cliente — ação imediata necessária."),
    ("MÉDIA",  LARANJA,  "Reativação do Módulo de Avaliação (RFI)",
     "Parado há 8 meses. Entender causa e propor plano de reengajamento com suporte da Efcaz."),
    ("MÉDIA",  LARANJA,  "Buscas Automáticas — Configurar e Expandir",
     "Ativar validade nos 4 documentos existentes para renovação contínua + ampliar escopo de automação nos 65 documentos hoje gerenciados manualmente."),
    ("MÉDIA",  LARANJA,  "BPO de Gestão Documental",
     "Com 65 documentos manuais e 53 fornecedores com vencidos, a Efcaz pode assumir a gestão documental — reduzindo retrabalho e garantindo compliance contínuo."),
]

for i, (urg, cor, titulo, desc) in enumerate(opps):
    top = 1.35 + i * 1.18
    add_rect(slide, 0.5, top, 12.33, 1.05, fill_color=CINZA_CLARO)
    add_rect(slide, 0.5, top, 1.5, 1.05, fill_color=cor)
    add_text(slide, urg, 0.5, top+0.32, 1.5, 0.4, font_size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_text(slide, titulo, 2.2, top+0.08, 9.8, 0.38, font_size=13, bold=True, color=PRETO_SUAVE)
    add_text(slide, desc, 2.2, top+0.52, 9.8, 0.45, font_size=11, color=CINZA_MEDIO)

# ─────────────────────────────────────────────
# SLIDE 8 — PLANO DE AÇÃO / PRÓXIMOS PASSOS
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 1.1, fill_color=AZUL_ESCURO)
add_text(slide, "Plano de Ação — 28 Dias para a Renovação", 0.5, 0.18, 12, 0.75,
         font_size=28, bold=True, color=BRANCO)

etapas = [
    ("Semana 1\naté 10/05", VERDE,
     "Agendar QBR / Reunião de Renovação",
     ["Contato: Arthur Sumihara (champion identificado no BI)",
      "Pauta: resultados do período + oportunidades + renovação",
      "Duração sugerida: 40 minutos"]),
    ("Semana 2\naté 17/05", AZUL_MEDIO,
     "Conduzir Reunião e Enviar Proposta",
     ["Apresentar crescimento (34x em solicitações, 95% de fornecedores)",
      "Levantar causa da parada no módulo de Avaliação",
      "Propor plano de ativação das integrações dormentes",
      "Enviar proposta com opção de upgrade de capacidade"]),
    ("Semana 3-4\naté 31/05", LARANJA,
     "Negociação e Assinatura",
     ["Follow-up da proposta (máx. 3 dias sem resposta)",
      "Alinhar jurídico e financeiro do cliente",
      "Assinatura antes do vencimento — garantir continuidade dos bônus"]),
]

for i, (periodo, cor, titulo, bullets) in enumerate(etapas):
    top = 1.35 + i * 1.95
    add_rect(slide, 0.5, top, 2.0, 1.75, fill_color=cor)
    add_text(slide, periodo, 0.5, top+0.55, 2.0, 0.65,
             font_size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_rect(slide, 2.65, top, 10.18, 1.75, fill_color=CINZA_CLARO)
    add_text(slide, titulo, 2.8, top+0.08, 10.0, 0.38, font_size=13, bold=True, color=PRETO_SUAVE)
    for j, b in enumerate(bullets):
        add_text(slide, "• " + b, 2.85, top+0.5+j*0.3, 9.9, 0.3, font_size=10, color=CINZA_MEDIO)

add_rect(slide, 0.5, 7.05, 12.33, 0.35, fill_color=AZUL_ESCURO)
add_text(slide, "Vencimento: 31/05/2026  •  Valor mensal: R$ 3.644,00  •  Contato principal: Arthur Sumihara — arthur.sumihara@dockbrasil.com.br",
         0.65, 7.1, 12, 0.28, font_size=11, color=BRANCO)

# ─────────────────────────────────────────────
# SLIDE 9 — ENCERRAMENTO
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=AZUL_ESCURO)
add_rect(slide, 0, 3.5, 13.33, 0.06, fill_color=LARANJA)
add_rect(slide, 0.5, 1.5, 0.1, 2.0, fill_color=LARANJA)

add_text(slide, "Obrigado", 0.9, 1.5, 11, 1.2,
         font_size=60, bold=True, color=BRANCO)
add_text(slide, "Gabriel Vital  •  Customer Success Specialist  •  Efcaz SRM",
         0.9, 2.8, 11, 0.5, font_size=16, color=RGBColor(0xAE, 0xD6, 0xF1))
add_text(slide, "vitallgabriel@outlook.com", 0.9, 4.0, 6, 0.4,
         font_size=14, color=BRANCO)
add_text(slide, "Efcaz SRM  •  efcaz.com.br", 0.9, 4.5, 6, 0.4,
         font_size=14, color=RGBColor(0xAE, 0xD6, 0xF1))

output = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\DockBrasil_Renovacao_Maio2026.pptx"
prs.save(output)
print(f"Arquivo gerado: {output}")
