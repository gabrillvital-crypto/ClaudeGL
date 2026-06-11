import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

TEAL        = RGBColor(0x0E, 0x8F, 0xA3)
TEAL_CLARO  = RGBColor(0x14, 0xB3, 0xCC)
TEAL_ESCURO = RGBColor(0x0A, 0x6A, 0x7A)
VERDE       = RGBColor(0x27, 0xAE, 0x60)
VERMELHO    = RGBColor(0xE7, 0x4C, 0x3C)
LARANJA     = RGBColor(0xF3, 0x9C, 0x12)
BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLARO = RGBColor(0xF2, 0xF4, 0xF4)
CINZA_MEDIO = RGBColor(0x85, 0x92, 0x9E)
PRETO_SUAVE = RGBColor(0x17, 0x20, 0x2A)

LOGO_DOCK  = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\logo_dock.png"
LOGO_EFCAZ = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\logo_efcaz.png"


def _remove_white_bg(src):
    """Remove fundo branco com gradiente suave nas bordas anti-aliased."""
    if not os.path.exists(src):
        return src
    base, _ = os.path.splitext(src)
    dst = base + "_proc.png"
    img = Image.open(src).convert("RGBA")
    px = img.load()
    w, h = img.size
    for y in range(h):
        for x in range(w):
            r, g, b, a = px[x, y]
            whiteness = min(r, g, b)
            if whiteness >= 245:
                px[x, y] = (r, g, b, 0)
            elif whiteness >= 210:
                alpha = int((245 - whiteness) / 35 * 255)
                px[x, y] = (r, g, b, alpha)
    img.save(dst, "PNG")
    return dst


LOGO_DOCK_P  = _remove_white_bg(LOGO_DOCK)
LOGO_EFCAZ_P = _remove_white_bg(LOGO_EFCAZ)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def rect(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h, size=13, bold=False, color=PRETO_SUAVE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb


def multiline(slide, lines, l, t, w, h, size=12, color=PRETO_SUAVE):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"


def efcaz_logo(slide):
    """Logo Efcaz pequeno no canto superior direito — aparece em todos os slides."""
    if os.path.exists(LOGO_EFCAZ_P):
        # Efcaz logo: 223x78 px, ratio 2.86. width=1.8" → right edge = 11.35+1.8 = 13.15" (OK)
        slide.shapes.add_picture(LOGO_EFCAZ_P, Inches(11.35), Inches(0.15), width=Inches(1.8))


def header(slide, titulo, subtitulo=""):
    rect(slide, 0, 0, 13.33, 1.0, fill=TEAL)
    rect(slide, 0, 0.95, 13.33, 0.06, fill=TEAL_CLARO)
    txt(slide, titulo, 0.5, 0.13, 11.2, 0.72, size=26, bold=True, color=BRANCO)
    if subtitulo:
        txt(slide, subtitulo, 0.5, 1.12, 12, 0.35, size=11, color=CINZA_MEDIO)
    efcaz_logo(slide)


def rodape(slide):
    rect(slide, 0, 7.15, 13.33, 0.35, fill=TEAL_ESCURO)
    txt(slide, "Efcaz SRM  •  Customer Success  •  Gabriel Vital  •  Mai/2026",
        0.4, 7.2, 12.5, 0.28, size=9, color=BRANCO)


# ─────────────────────────────────────────────
# SLIDE 1 — CAPA
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=TEAL)
rect(slide, 0, 4.5, 13.33, 3.0, fill=TEAL_ESCURO)
rect(slide, 0.5, 2.1, 0.1, 1.8, fill=BRANCO)

txt(slide, "DOCK BRASIL", 0.85, 0.95, 7.5, 1.0, size=48, bold=True, color=BRANCO)
txt(slide, "Engenharia e Serviços S.A.", 0.85, 2.0, 7.5, 0.52,
    size=20, color=RGBColor(0xD6, 0xF5, 0xF8))
txt(slide, "Próximos Passos — Renovação Dock Brasil  |  Mai/2026", 0.85, 4.7, 10, 0.52,
    size=17, bold=True, color=BRANCO)
txt(slide, "Gabriel Vital  •  Customer Success Specialist  •  Efcaz SRM", 0.85, 5.32, 10, 0.42,
    size=12, color=RGBColor(0xD6, 0xF5, 0xF8))

# Logo Dock — centro-direito, abaixo do título, acima da faixa escura
# ratio 2.12, width=4.0" → height≈1.89", right edge=12.8", bottom=4.29" — tudo dentro do slide
if os.path.exists(LOGO_DOCK_P):
    slide.shapes.add_picture(LOGO_DOCK_P, Inches(8.8), Inches(2.4), width=Inches(4.0))

# Logo Efcaz no canto superior direito
efcaz_logo(slide)

# ─────────────────────────────────────────────
# SLIDE 2 — SITUAÇÃO ATUAL
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Situação Atual — Dock Brasil", "Diagnóstico do momento de renovação  •  Mai/2026")

cards = [
    (VERMELHO,    "31/mai/2026",         "Vencimento do contrato"),
    (LARANJA,     "763 / 800",           "Fornecedores ativos (95,4%)"),
    (LARANJA,     "25 / 15",             "Usuários em uso vs. contratados"),
    (LARANJA,     "8 usu. + 4 buscas",  "Bônus que expiram em 31/mai"),
    (LARANJA,     "53 fornecedores",     "Com documentos vencidos"),
    (CINZA_MEDIO, "Parado set/2025",     "Módulo de Avaliação (RFI)"),
]

cw, ch = 3.8, 1.52
for j, (cor, numero, label) in enumerate(cards):
    col = j % 3
    row = j // 3
    cx = 0.42 + col * (cw + 0.22)
    cy = 1.55 + row * (ch + 0.22)
    rect(slide, cx, cy, cw, ch, fill=cor)
    txt(slide, numero, cx + 0.12, cy + 0.1, cw - 0.24, ch * 0.52,
        size=20, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    txt(slide, label, cx + 0.1, cy + ch * 0.62, cw - 0.2, ch * 0.34,
        size=10, color=BRANCO, align=PP_ALIGN.CENTER, wrap=True)

rodape(slide)

# ─────────────────────────────────────────────
# SLIDE 3 — CONTEXTO DA RENOVAÇÃO
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Contexto da Renovação — O Momento de Decidir", "")

pontos = [
    (LARANJA, "ATENÇÃO", "Acesso da equipe — garantir continuidade",
     ["25 usuários ativos hoje. O contrato cobre 15 + 8 auxiliares de bônus, que expiram em 31/mai.",
      "Renovar agora garante que toda a equipe continua operando sem nenhuma interrupção."]),
    (LARANJA, "ATENÇÃO", "Buscas automáticas — manter ativas",
     ["CNDT, CND, CNPJ e FGTS são bônus do aditivo atual — renovar mantém tudo funcionando.",
      "Com a renovação, as certidões continuam sendo atualizadas automaticamente."]),
    (LARANJA,  "ATENÇÃO", "Capacidade de fornecedores próxima do limite",
     ["763 de 800 fornecedores ativos — 95,4% de ocupação com crescimento constante.",
      "Expandir agora para 1.120 garante continuidade das homologações sem interrupção."]),
]

for i, (cor, urgencia, titulo, linhas) in enumerate(pontos):
    top = 1.5 + i * 1.85
    rect(slide, 0.5, top, 12.33, 1.65, fill=CINZA_CLARO)
    rect(slide, 0.5, top, 1.6, 1.65, fill=cor)
    txt(slide, urgencia, 0.5, top + 0.55, 1.6, 0.5,
        size=11, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    txt(slide, titulo, 2.3, top + 0.1, 10.2, 0.4, size=14, bold=True, color=PRETO_SUAVE)
    multiline(slide, linhas, 2.3, top + 0.58, 10.2, 0.95, size=11, color=CINZA_MEDIO)

rodape(slide)

# ─────────────────────────────────────────────
# SLIDE 4 — PORTFÓLIO DE EXPANSÃO
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Portfólio de Expansão", "3 oportunidades concretas para esta renovação")

upsells = [
    (TEAL, "01", "Expansão de Usuários", "15 → 25 usuários",
     ["Formaliza o que já é a realidade operacional da Dock Brasil",
      "Garante acesso contínuo para toda a equipe após 31/mai",
      "Investimento adicional: R$ 300,00/mês"]),
    (VERDE, "02", "Expansão de Fornecedores", "800 → 1.120 (+40%)",
     ["763 ativos hoje — ritmo constante de novas homologações",
      "Expansão de 40% acompanha o crescimento da operação",
      "Investimento adicional: R$ 200,00/mês"]),
    (LARANJA, "03", "Buscas IBAMA", "3 consultas automáticas/mês",
     ["Certidão IBAMA: obrigatória para fornecedores do setor",
      "Automação protege em auditorias e processos licitatórios",
      "Investimento: R$ 250,00/mês — bonificado nesta proposta"]),
]

for i, (cor, num, titulo, subtitulo, bullets) in enumerate(upsells):
    cx = 0.42 + i * 4.3
    cy = 1.5
    cw2, ch2 = 4.0, 5.4
    rect(slide, cx, cy, cw2, ch2, fill=CINZA_CLARO)
    rect(slide, cx, cy, cw2, 0.55, fill=cor)
    txt(slide, num, cx + 0.15, cy + 0.1, 0.5, 0.38, size=16, bold=True, color=BRANCO)
    txt(slide, titulo, cx + 0.72, cy + 0.1, cw2 - 0.85, 0.38, size=13, bold=True, color=BRANCO)
    txt(slide, subtitulo, cx + 0.15, cy + 0.72, cw2 - 0.3, 0.5, size=16, bold=True, color=cor)
    for j, b in enumerate(bullets):
        txt(slide, "• " + b, cx + 0.15, cy + 1.38 + j * 0.85, cw2 - 0.3, 0.75,
            size=11, color=PRETO_SUAVE, wrap=True)

rodape(slide)

# ─────────────────────────────────────────────
# SLIDE 5 — POR QUE EXPANDIR AGORA
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Por que Expandir Agora", "O valor de cada decisão para a Dock Brasil")

itens = [
    (TEAL, "Usuários (15 → 25)",
     ["Sua operação já funciona com 25 usuários — formalizar é só alinhar o contrato à realidade.",
      "Com o contrato atualizado, toda a equipe segue operando normalmente, sem nenhuma interrupção.",
      "Investimento adicional: R$ 300,00/mês para garantir o que já está funcionando."]),
    (VERDE, "Fornecedores (800 → 1.120)",
     ["Crescimento consistente: 763 ativos hoje, com ritmo que demanda mais capacidade em breve.",
      "Ampliar agora garante que o crescimento da Dock Brasil não para por limite de plataforma.",
      "Expansão planejada, com custo previsível — sem surpresas no meio da operação."]),
    (LARANJA, "Buscas IBAMA — bonificadas nesta proposta",
     ["Engenharia e serviços: fornecedores sujeitos a auditorias ambientais e exigências em licitações.",
      "Automatizar a certidão IBAMA elimina o risco de um fornecedor irregular comprometer um contrato.",
      "Nesta renovação, as 3 consultas mensais são 100% bonificadas — valor agregado sem custo adicional."]),
]

for i, (cor, titulo, linhas) in enumerate(itens):
    top = 1.5 + i * 1.82
    rect(slide, 0.5, top, 12.33, 1.62, fill=CINZA_CLARO)
    rect(slide, 0.5, top, 3.0, 0.42, fill=cor)
    rect(slide, 0.5, top, 0.18, 1.62, fill=cor)
    txt(slide, titulo, 0.75, top + 0.07, 2.8, 0.3, size=12, bold=True, color=BRANCO)
    multiline(slide, linhas, 0.85, top + 0.52, 12.0, 1.0, size=11, color=CINZA_MEDIO)

rodape(slide)

# ─────────────────────────────────────────────
# SLIDE 6 — PROPOSTA DE RENOVAÇÃO
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Proposta de Renovação", "Composição do investimento  •  Mai/2026")

# Painel esquerdo — Hoje
rect(slide, 0.5, 1.45, 3.2, 5.45, fill=TEAL)
txt(slide, "HOJE", 0.5, 1.6, 3.2, 0.5,
    size=14, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
rect(slide, 0.7, 2.22, 2.8, 0.04, fill=RGBColor(0xFF, 0xFF, 0xFF))
linhas_hoje = [
    ("R$ 3.644,00/mês", 18, True),
    ("800 fornecedores", 12, False),
    ("15 usuários", 12, False),
    ("4 buscas automáticas", 12, False),
    ("(bônus — expirando)", 10, False),
]
for j, (texto, sz, negrito) in enumerate(linhas_hoje):
    txt(slide, texto, 0.6, 2.38 + j * 0.72, 3.0, 0.62,
        size=sz, bold=negrito, color=BRANCO, align=PP_ALIGN.CENTER)

# Painel direito — Nova Proposta
rect(slide, 3.9, 1.45, 9.0, 5.45, fill=CINZA_CLARO)
txt(slide, "NOVA PROPOSTA", 3.9, 1.55, 9.0, 0.42,
    size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

linhas_prop = [
    ("Base com reajuste IPCA (4,14%)",          "R$ 3.794,86", PRETO_SUAVE, False),
    ("+ Expansão de Usuários (15 → 25)",         "+ R$ 300,00", PRETO_SUAVE, False),
    ("+ Expansão de Fornecedores (800 → 1.120)", "+ R$ 200,00", PRETO_SUAVE, False),
    ("+ Buscas Automáticas IBAMA",               "+ R$ 250,00", PRETO_SUAVE, False),
    ("Subtotal",                                 "R$ 4.544,86", CINZA_MEDIO, False),
]

for j, (desc, valor, cor_txt, negrito) in enumerate(linhas_prop):
    top_l = 2.1 + j * 0.62
    bg = BRANCO if j % 2 == 0 else CINZA_CLARO
    rect(slide, 3.9, top_l, 9.0, 0.58, fill=bg)
    txt(slide, desc,  4.1, top_l + 0.12, 6.5, 0.36, size=11, color=cor_txt, bold=negrito)
    txt(slide, valor, 10.6, top_l + 0.12, 2.1, 0.36, size=11, color=cor_txt,
        bold=negrito, align=PP_ALIGN.RIGHT)

# Bonificação
rect(slide, 3.9, 5.2, 9.0, 0.55, fill=RGBColor(0xD5, 0xF5, 0xE3))
txt(slide, "Bonificação IBAMA — Condição especial Dock Brasil",
    4.1, 5.28, 6.5, 0.35, size=11, color=VERDE, bold=True)
txt(slide, "- R$ 250,00", 10.6, 5.28, 2.1, 0.35, size=11, color=VERDE,
    bold=True, align=PP_ALIGN.RIGHT)

# Valor final
rect(slide, 3.9, 5.85, 9.0, 1.05, fill=TEAL)
txt(slide, "VALOR FINAL DE RENOVAÇÃO", 3.9, 5.92, 9.0, 0.38,
    size=11, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
txt(slide, "R$ 4.294,86/mês", 3.9, 6.3, 9.0, 0.55,
    size=28, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

rodape(slide)

# ─────────────────────────────────────────────
# SLIDE 7 — PLANO DE AÇÃO
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Plano de Ação — 27 Dias para a Renovação", "")

etapas = [
    ("Semana 1\naté 10/05", VERDE,
     "Agendar QBR com Arthur Sumihara",
     ["Contato: arthur.sumihara@dockbrasil.com.br",
      "Pauta: resultados do período + contexto de renovação + proposta de expansão",
      "Duração sugerida: 40 minutos"]),
    ("Semana 2\naté 17/05", TEAL,
     "Conduzir Reunião + Enviar Proposta Atualizada",
     ["Apresentar crescimento (34x solicitações, 95% fornecedores)",
      "Demonstrar impacto do vencimento dos bônus em 31/mai",
      "Apresentar proposta: 25 usuários + 1.120 fornecedores + IBAMA bonificado",
      "Retomar módulo de Avaliação (RFI) — parado desde set/2025"]),
    ("Semana 3-4\naté 31/05", LARANJA,
     "Negociação e Assinatura",
     ["Follow-up máximo 3 dias após envio da proposta",
      "Alinhar com jurídico e financeiro do cliente",
      "Garantir assinatura antes de 31/mai — continuidade operacional sem interrupção"]),
]

for i, (periodo, cor, titulo, bullets) in enumerate(etapas):
    top = 1.3 + i * 1.93
    rect(slide, 0.5, top, 2.0, 1.72, fill=cor)
    txt(slide, periodo, 0.5, top + 0.52, 2.0, 0.68,
        size=12, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    rect(slide, 2.65, top, 10.18, 1.72, fill=CINZA_CLARO)
    txt(slide, titulo, 2.8, top + 0.08, 10.0, 0.38, size=13, bold=True, color=PRETO_SUAVE)
    for j, b in enumerate(bullets):
        txt(slide, "• " + b, 2.85, top + 0.5 + j * 0.3, 9.9, 0.3, size=10, color=CINZA_MEDIO)

rect(slide, 0.5, 7.05, 12.33, 0.35, fill=TEAL_ESCURO)
txt(slide, "Vencimento: 31/mai/2026  •  Investimento atual: R$ 3.644,00/mês  •  Champion: Arthur Sumihara — arthur.sumihara@dockbrasil.com.br",
    0.65, 7.1, 12, 0.28, size=10, color=BRANCO)

# ─────────────────────────────────────────────
# SLIDE 8 — ENCERRAMENTO
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=TEAL)
rect(slide, 0, 3.5, 13.33, 0.06, fill=BRANCO)
rect(slide, 0.5, 1.5, 0.1, 2.0, fill=BRANCO)

txt(slide, "Juntos,", 0.9, 1.4, 11, 0.85, size=52, bold=True, color=BRANCO)
txt(slide, "Dock Brasil cresce. A Efcaz cresce junto.", 0.9, 2.25, 11, 0.85, size=36, bold=False, color=BRANCO)
txt(slide, "Gabriel Vital  •  Customer Success Specialist  •  Efcaz SRM",
    0.9, 3.72, 11, 0.5, size=15, color=RGBColor(0xD6, 0xF5, 0xF8))
txt(slide, "vitallgabriel@outlook.com", 0.9, 4.35, 6, 0.42, size=13, color=BRANCO)
txt(slide, "efcaz.com.br", 0.9, 4.85, 6, 0.42, size=13, color=RGBColor(0xD6, 0xF5, 0xF8))

efcaz_logo(slide)

output = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\DockBrasil_PlanoAcao_Mai2026.pptx"
prs.save(output)
print(f"Arquivo gerado: {output}")
