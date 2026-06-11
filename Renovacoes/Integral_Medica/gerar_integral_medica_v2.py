"""
Integral Médica — Proposta de Renovação Jun/2026
Abordagem: editar template Dock Brasil preservando layout, cores e estilo
Gabriel Vital | Efcaz CS
"""

import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

TEMPLATE = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\Documentos\Plano de Ação Efcaz _ DockBrasil - modelo.pptx"
LOGO_IM  = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\Documentos\logos_extraidas\slide1_img6.png"
OUTPUT   = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\Integral_Medica_Renovacao_Jun2026_v2.pptx"

prs = Presentation(TEMPLATE)

# ─── MAPA DE SUBSTITUIÇÕES ────────────────────────────────────────────────────
substituicoes = {

    # ── Identificação — ESPECÍFICOS PRIMEIRO, genérico por último ────────────
    "Renovação Dock Brasil":            "Renovação Integral Médica",
    "Bonificação IBAMA — Condição especial Dock Brasil": "",
    "O valor de cada decisão para a Dock Brasil":
        "O valor de cada decisão para a Integral Médica",
    "Ampliar agora garante que o crescimento da Dock Brasil não para por limite de plataforma.":
        "Uma sessão de ativação com o CS é suficiente para colocar a equipe em funcionamento.",
    "Formaliza o que já é a realidade operacional da Dock Brasil":
        "Formaliza o que já é a realidade operacional da Integral Médica",
    "Juntos, Dock Brasil cresce. ":     "Juntos, Integral Médica cresce. ",
    # Genérico — deve vir por último
    "Dock Brasil":                      "Integral Médica",

    # ── Datas ─────────────────────────────────────────────────────────────────
    "31/mai/2026":                      "30/jun/2026",
    "Mai/2026":                         "Jun/2026",
    "após 31/mai":                      "após 30/jun",

    # ── Slide 2 — KPIs ────────────────────────────────────────────────────────
    "Fornecedores ativos (95,4%)":      "Fornecedores cadastrados (80,7%)",
    "763 / 800":                        "242 / 300",
    "25 / 15":                          "35 / 21",
    "Bônus que expiram em 31/mai":      "Fabricantes cadastrados",
    "8 usu. + 4 buscas":               "242 ativos",
    "53 fornecedores":                  "80 fornecedores",
    "com documentos vencidos":          "com certidões vencidas",
    "Parado set/2025":                  "Zero registros",

    # ── Slide 3 — Contexto ────────────────────────────────────────────────────
    "Contexto da Renovação — O Momento de Decidir":
        "Contexto da Renovação — O Momento de Agir",

    "Acesso da equipe — garantir continuidade":
        "Acesso da equipe — formalizar o que já existe",

    "25 usuários ativos hoje. O contrato cobre 15 + 8 auxiliares de bônus, que expiram em 31/mai.":
        "35 usuários operam na plataforma hoje. O contrato cobre apenas 21.",

    "Renovar agora garante que toda a equipe continua operando sem nenhuma interrupção.":
        "Renovar e regularizar garante que toda a equipe continua operando sem nenhuma interrupção.",

    "Buscas automáticas — manter ativas":
        "Gestão de Fabricantes — foco central da operação",

    "CNDT, CND, CNPJ e FGTS são bônus do aditivo atual — renovar mantém tudo funcionando.":
        "Com 242 fabricantes cadastrados e crescimento expressivo desde 2024, a plataforma é o núcleo da gestão de matéria-prima.",

    "Com a renovação, as certidões continuam sendo atualizadas automaticamente.":
        "A operação da Integral Médica depende diretamente da continuidade da plataforma.",

    "Capacidade de fornecedores próxima do limite":
        "Certidões vencidas — risco de compliance",

    "763 de 800 fornecedores ativos — 95,4% de ocupação com crescimento constante.":
        "80 fornecedores com certidões vencidas representam exposição jurídica.",

    "Expandir agora para 1.120 garante continuidade das homologações sem interrupção.":
        "Ativar alertas automáticos de vencimento resolve de forma preventiva — sem custo adicional.",

    # ── Slide 4 — Portfólio ───────────────────────────────────────────────────
    "01. Expansão de usuários":         "01. Regularização de Usuários",
    "02. Expansão de fornecedores":     "02. Ativar Avaliações",
    "03. Buscas IBAMA":                 "03. Gestão de Certidões",
    "15 → 25 usuários":                 "21 → 35 usuários",
    "Garante acesso contínuo para toda a equipe após 31/mai":
        "Garante acesso contínuo para toda a equipe após 30/jun",
    "R$ 300,00/mês":                    "R$ 420,00/mês",
    "800 → 1.120 (+40%)":              "Módulo já contratado",
    "763 ativos hoje — ritmo constante de novas homologações":
        "Zero avaliações registradas até hoje",
    "Expansão de 40% acompanha o crescimento da operação":
        "Ativar agora: sem custo adicional",
    "R$ 200,00/mês":                    "R$ 0,00 (já incluso)",
    "3 consultas automáticas/mês":      "80 fornecedores em risco",
    "Certidão IBAMA: obrigatória para fornecedores do setor":
        "Certidões vencidas: risco jurídico e compliance",
    "Automação protege em auditorias e processos licitatórios":
        "Alertas automáticos já inclusos no plano",
    "R$ 250,00/mês":                    "R$ 0,00",
    " bonificado nesta proposta":       "(sem custo adicional)",
    "3 oportunidades concretas para esta renovação":
        "2 oportunidades concretas para esta renovação",

    # ── Slide 5 — Por que agir ────────────────────────────────────────────────
    "Por que Expandir Agora":           "Por que Regularizar Agora",
    "Fornecedores (800 → 1.120)":       "Módulo de Avaliações",
    # "Buscas IBAMA" já coberto acima (slide 4)

    "Sua operação já funciona com 25 usuários — formalizar é só alinhar o contrato à realidade.":
        "Sua operação já funciona com 35 usuários — formalizar é só alinhar o contrato à realidade.",

    "R$ 300,00/mês para garantir o que já está funcionando.":
        "R$ 420,00/mês para garantir o que já está funcionando.",

    "O valor de cada decisão para a Dock Brasil":
        "O valor de cada decisão para a Integral Médica",

    "Crescimento consistente: 763 ativos hoje, com ritmo que demanda mais capacidade em breve.":
        "Módulo já pago e nunca utilizado. Ativar agora significa extrair mais valor do investimento já feito.",

    "Ampliar agora garante que o crescimento da Dock Brasil não para por limite de plataforma.":
        "Uma sessão de ativação com o CS é suficiente para colocar a equipe em funcionamento.",

    "Expansão planejada, com custo previsível — sem surpresas no meio da operação.":
        "Sem custo extra — o módulo já está incluso no contrato atual.",

    "Engenharia e serviços: fornecedores sujeitos a auditorias ambientais e exigências em licitações.":
        "80 fornecedores com certidões vencidas é um risco real de compliance e contratos.",

    "Automatizar a certidão IBAMA elimina o risco de um fornecedor irregular comprometer um contrato.":
        "Configurar alertas automáticos de vencimento resolve de forma preventiva.",

    "Nesta renovação, as 3 consultas mensais são 100% bonificadas — valor agregado sem custo adicional.":
        "Ação imediata, sem custo adicional — os alertas já fazem parte da plataforma.",

    # ── Slide 6 — Proposta financeira ─────────────────────────────────────────
    "Composição do investimento  •  Mai/2026":
        "Composição do investimento  •  Jun/2026",

    "800 fornecedores":                 "300 fornecedores",
    "15 usuários":                      "21 usuários",
    "4 buscas automáticas":             "(contratados)",
    "(bônus — expirando)":              "",
    "R$ 3.644,00/mês":                  "R$ 2.181,92/mês",
    "Base com reajuste IPCA (4,14%)":   "Base com reajuste IPCA (4,39%)",
    "+ Expansão de Usuários (15 → 25)": "+ Regularização de Usuários (21 → 35)",
    "+ Expansão de Fornecedores (800 → 1.120)": "",
    "R$ 3.794,86":                      "R$ 2.277,73",
    "+ R$ 300,00":                      "+ R$ 420,00",
    "+ R$ 200,00":                      "",
    "Subtotal":                         "",
    "R$ 4.544,86":                      "",
    "+ Buscas Automáticas IBAMA":       "",
    "- R$ 250,00":                      "",
    "+ R$ 250,00":                      "",
    "R$ 4.294,86/mês":                  "R$ 2.697,73/mês",

    # ── Slide 7 — Encerramento ────────────────────────────────────────────────
    "Juntos, Dock Brasil cresce. ":     "Juntos, Integral Médica cresce. ",
}

# ─── FUNÇÃO DE SUBSTITUIÇÃO ───────────────────────────────────────────────────
def substituir_textos(prs, subs):
    total = 0
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                # Concatenar todos os runs do parágrafo
                original = "".join(r.text for r in para.runs)
                novo = original
                for old, new in subs.items():
                    novo = novo.replace(old, new)
                if novo != original and para.runs:
                    # Escreve tudo no primeiro run, zera os demais
                    para.runs[0].text = novo
                    for r in para.runs[1:]:
                        r.text = ""
                    total += 1
                    print(f"  [S{i+1}] {repr(original[:60])} -> {repr(novo[:60])}")
    return total

print("Substituindo textos...")
n = substituir_textos(prs, substituicoes)
print(f"\nTotal de parágrafos alterados: {n}")

# ─── SUBSTITUIR LOGO DO CLIENTE ───────────────────────────────────────────────
print("\nSubstituindo logo do cliente...")
slide1 = prs.slides[0]
substituida = False
for shape in slide1.shapes:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        left_in = shape.left / 914400
        top_in  = shape.top  / 914400
        # Logo do cliente está em (2.704, 3.502) no template
        if abs(left_in - 2.704) < 0.3 and abs(top_in - 3.502) < 0.3:
            # Encontrar o blob da imagem via relacionamento
            blip = shape._pic.find(
                './/{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
            )
            if blip is not None:
                rId = blip.get(
                    '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                )
                if rId:
                    img_part = slide1.part.related_part(rId)
                    with open(LOGO_IM, 'rb') as f:
                        img_part._blob = f.read()
                    print(f"  Logo substituída (rId={rId}, pos={left_in:.2f},{top_in:.2f})")
                    substituida = True

if not substituida:
    print("  AVISO: logo do cliente não encontrada no slide 1 — verifique posição")

# ─── SALVAR ───────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"\n✓ Arquivo salvo em:\n  {OUTPUT}")
