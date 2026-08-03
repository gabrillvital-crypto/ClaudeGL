"""
Relatório do Ciclo Zurich Airport — 10/07 a 27/07/2026
Gerado para o time de BPO montar apresentação.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import datetime

# ── Paleta Efcaz ──────────────────────────────────────────────────────────────
TEAL        = RGBColor(0x0E, 0x8F, 0xA3)
TEAL_ESCURO = RGBColor(0x0A, 0x6A, 0x7A)
TEAL_CLARO  = RGBColor(0xE0, 0xF5, 0xF8)
VERDE       = RGBColor(0x27, 0xAE, 0x60)
VERDE_BG    = RGBColor(0xD4, 0xED, 0xDA)
VERMELHO    = RGBColor(0xE7, 0x4C, 0x3C)
LARANJA     = RGBColor(0xF3, 0x9C, 0x12)
BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLARO = RGBColor(0xF2, 0xF4, 0xF4)
CINZA_MEDIO = RGBColor(0x85, 0x92, 0x9E)
PRETO_SUAVE = RGBColor(0x17, 0x20, 0x2A)
AMARELO_BG  = RGBColor(0xFF, 0xF3, 0xCD)

SW = Inches(13.33)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

BLANK = prs.slide_layouts[6]  # layout em branco

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=14, bold=False, color=PRETO_SUAVE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return txBox

def add_header_band(slide, title, subtitle=None):
    """Banda teal no topo com título branco."""
    band_h = Inches(1.3)
    add_rect(slide, 0, 0, SW, band_h, TEAL)
    add_text(slide, title,
             Inches(0.4), Inches(0.15), Inches(11), Inches(0.7),
             font_size=28, bold=True, color=BRANCO)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.85), Inches(10), Inches(0.4),
                 font_size=14, color=RGBColor(0xCC, 0xEE, 0xF3))
    # rodapé discreto
    add_text(slide, "Efcaz SRM  |  Ciclo 10/07–27/07/2026",
             Inches(0.3), Inches(7.15), Inches(8), Inches(0.3),
             font_size=9, color=CINZA_MEDIO)

def kpi_card(slide, x, y, w, h, value, label, value_color=TEAL, bg_color=CINZA_CLARO, delta=None):
    """Card de KPI com valor grande + label pequeno."""
    add_rect(slide, x, y, w, h, bg_color, RGBColor(0xCC, 0xCC, 0xCC))
    cy = y + Inches(0.12)
    add_text(slide, value,
             x + Inches(0.1), cy, w - Inches(0.2), Inches(0.7),
             font_size=30, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    cy2 = y + Inches(0.75)
    add_text(slide, label,
             x + Inches(0.05), cy2, w - Inches(0.1), Inches(0.45),
             font_size=11, color=PRETO_SUAVE, align=PP_ALIGN.CENTER)
    if delta:
        cy3 = y + Inches(1.18)
        col = VERDE if delta.startswith('+') else VERMELHO
        add_text(slide, delta,
                 x + Inches(0.05), cy3, w - Inches(0.1), Inches(0.25),
                 font_size=10, bold=True, color=col, align=PP_ALIGN.CENTER)

def table_header_row(table, headers, header_color=TEAL, text_color=BRANCO, font_size=11):
    row = table.rows[0]
    for i, h in enumerate(headers):
        cell = row.cells[i]
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.bold = True
        run.font.size = Pt(font_size)
        run.font.color.rgb = text_color
        run.font.name = 'Calibri'

def table_data_row(table, row_idx, values, alt=False, bold_last=False, highlight_color=None):
    row = table.rows[row_idx]
    bg = RGBColor(0xF0, 0xF8, 0xFA) if alt else BRANCO
    for i, v in enumerate(values):
        cell = row.cells[i]
        cell.text = str(v)
        cell.fill.solid()
        cell.fill.fore_color.rgb = highlight_color if highlight_color else bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(10)
        run.font.bold = bold_last and i == len(values) - 1
        run.font.name = 'Calibri'
        run.font.color.rgb = PRETO_SUAVE

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SW, SH, TEAL_ESCURO)
add_rect(sl, 0, SH - Inches(2.2), SW, Inches(2.2), TEAL)
# Logo text placeholder
add_text(sl, "EFCAZ", Inches(0.5), Inches(0.3), Inches(3), Inches(0.6),
         font_size=22, bold=True, color=BRANCO)
add_text(sl, "Supplier Relationship Management",
         Inches(0.5), Inches(0.85), Inches(6), Inches(0.35),
         font_size=11, color=RGBColor(0xCC, 0xEE, 0xF3))

add_text(sl, "Zurich Airport Brasil",
         Inches(0.5), Inches(2.2), Inches(12), Inches(1.0),
         font_size=44, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
add_text(sl, "Relatório do Ciclo — 10 a 27 de Julho de 2026",
         Inches(0.5), Inches(3.2), Inches(12), Inches(0.6),
         font_size=22, color=RGBColor(0xCC, 0xEE, 0xF3), align=PP_ALIGN.CENTER)
add_text(sl, "Preparado pelo time Efcaz para o BPO",
         Inches(0.5), Inches(3.85), Inches(12), Inches(0.45),
         font_size=15, italic=True, color=RGBColor(0xCC, 0xEE, 0xF3), align=PP_ALIGN.CENTER)

add_text(sl, f"Gerado em {datetime.date.today().strftime('%d/%m/%Y')}",
         Inches(9.5), SH - Inches(1.0), Inches(3.5), Inches(0.4),
         font_size=10, color=RGBColor(0xCC, 0xEE, 0xF3), align=PP_ALIGN.RIGHT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — RESUMO EXECUTIVO (KPIs grandes)
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_header_band(sl, "Resumo do Ciclo", "Principais números — 10/07 a 27/07/2026")

cw = Inches(2.4)
ch = Inches(1.6)
gap = Inches(0.22)
y_top = Inches(1.55)
x_start = Inches(0.35)

kpis = [
    ("44,9%",  "Conformidade\nFornecedores (R4)",  TEAL,    CINZA_CLARO, "+3,1 p.p."),
    ("24,6%",  "Conformidade\nTerceiros (R3)",      TEAL,    CINZA_CLARO, "+7,3 p.p."),
    ("2.696",  "Não Resolvidas\nem 27/07",          LARANJA, AMARELO_BG,  "3.013 no início"),
    ("-10,5%", "Redução de\nNão Resolvidas",        VERDE,   VERDE_BG,    "3.013 → 2.696"),
    ("+16",    "Novos Terceiros\nCadastrados",       TEAL,    CINZA_CLARO, "1.251 → 1.267"),
]

for idx, (val, lbl, vc, bg, delta) in enumerate(kpis):
    x = x_start + idx * (cw + gap)
    kpi_card(sl, x, y_top, cw, ch, val, lbl, vc, bg, delta)

# frase de destaque
add_rect(sl, Inches(0.35), Inches(3.35), Inches(12.63), Inches(0.95), TEAL_CLARO)
add_text(sl, "Das 3.013 pendências não resolvidas em 10/07, 317 foram regularizadas até o fechamento em 27/07 — "
             "encerrando o ciclo com 2.696 não resolvidas, uma redução de 10,5% no período.",
         Inches(0.55), Inches(3.42), Inches(12.2), Inches(0.8),
         font_size=13, color=TEAL_ESCURO)

# Segunda linha de KPIs — documentos
y2 = Inches(4.5)
kpis2 = [
    ("1.771", "Docs Fornecedores\nna Base (R4)", TEAL, CINZA_CLARO, "+219 no ciclo"),
    ("9.374", "Docs Terceiros\nna Base (R3)",    TEAL, CINZA_CLARO, "+531 no ciclo"),
    ("796",   "Docs Aprovados\nR4 (27/07)",      VERDE, VERDE_BG,   "+147 no ciclo"),
    ("2.306", "Docs Aprovados\nR3 (27/07)",      VERDE, VERDE_BG,   "+773 no ciclo"),
    ("51",    "Fornecedores\nAtivos",             TEAL, CINZA_CLARO, None),
]
for idx, (val, lbl, vc, bg, delta) in enumerate(kpis2):
    x = x_start + idx * (cw + gap)
    kpi_card(sl, x, y2, cw, Inches(1.5), val, lbl, vc, bg, delta)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — R4: DOCUMENTOS DE FORNECEDORES
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_header_band(sl, "R4 — Documentos de Fornecedores", "Comparativo início × fechamento do ciclo")

# Tabela principal
cols = ["Indicador", "10/07 (Início)", "27/07 (Fechamento)", "Variação"]
col_widths = [Inches(4.2), Inches(2.5), Inches(2.8), Inches(2.8)]
rows_data = [
    ("Total de documentos", "1.552", "1.771", "+219"),
    ("Aprovados",           "649",   "796",   "+147"),
    ("Reprovados",          "371",   "428",   "+57"),
    ("Não Anexados",        "195",   "239",   "+44"),
    ("Conformidade",        "41,8%", "44,9%", "+3,1 p.p. ✓"),
]

n_rows = len(rows_data) + 1
tbl = sl.shapes.add_table(n_rows, 4, Inches(0.35), Inches(1.55),
                           sum(col_widths), Inches(2.6)).table
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

table_header_row(tbl, cols)
for ri, (ind, v1, v2, var) in enumerate(rows_data):
    highlight = VERDE_BG if ri == 4 else None
    table_data_row(tbl, ri + 1, [ind, v1, v2, var],
                   alt=(ri % 2 == 1), highlight_color=highlight)
    if ri == 4:
        row = tbl.rows[ri + 1]
        for ci in range(4):
            p = row.cells[ci].text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x15, 0x57, 0x24)

# Gráfico de barras horizontais — conformidade início vs fim
add_rect(sl, Inches(7.7), Inches(1.55), Inches(5.3), Inches(4.5), CINZA_CLARO)
add_text(sl, "Conformidade R4 — Início × Fechamento",
         Inches(7.8), Inches(1.65), Inches(5.0), Inches(0.4),
         font_size=12, bold=True, color=TEAL)

# barra início
add_rect(sl, Inches(7.8), Inches(2.2), Inches(2.79), Inches(0.55), CINZA_MEDIO)
add_text(sl, "10/07  41,8%", Inches(10.7), Inches(2.3), Inches(2.5), Inches(0.35),
         font_size=12, bold=False, color=PRETO_SUAVE)
# barra fim
add_rect(sl, Inches(7.8), Inches(3.0), Inches(2.99), Inches(0.55), TEAL)
add_text(sl, "27/07  44,9%", Inches(10.9), Inches(3.1), Inches(2.5), Inches(0.35),
         font_size=12, bold=True, color=TEAL)

add_text(sl, "▲ +3,1 p.p. no ciclo",
         Inches(7.8), Inches(3.8), Inches(5.0), Inches(0.45),
         font_size=14, bold=True, color=VERDE)

# observação
add_rect(sl, Inches(0.35), Inches(4.35), Inches(12.63), Inches(0.85), TEAL_CLARO)
add_text(sl, "💡  Os 4 documentos com mais não-conformidades são trabalhistas mensais (FOPAG, DCTFWEB, GFD, Comprovante Bancário) — "
             "recomenda-se comunicação preventiva no início do próximo ciclo.",
         Inches(0.55), Inches(4.42), Inches(12.2), Inches(0.75),
         font_size=11, color=TEAL_ESCURO)

# Top docs problemáticos
add_text(sl, "Top documentos não conformes (R4):",
         Inches(0.35), Inches(5.35), Inches(6.5), Inches(0.3),
         font_size=11, bold=True, color=TEAL)
docs_r4 = [
    ("FOPAG — Folha de Pagamento + Resumo", "129"),
    ("DCTFWEB", "128"),
    ("GFD — Guia do FGTS Digital Mensal", "123"),
    ("Comprovante Bancário de Pagamento", "121"),
    ("CIPA", "44"),
]
for i, (doc, qtd) in enumerate(docs_r4):
    y_d = Inches(5.7) + i * Inches(0.28)
    add_text(sl, f"• {doc}", Inches(0.45), y_d, Inches(5.8), Inches(0.26),
             font_size=10, color=PRETO_SUAVE)
    add_text(sl, qtd, Inches(6.3), y_d, Inches(0.6), Inches(0.26),
             font_size=10, bold=True, color=VERMELHO, align=PP_ALIGN.RIGHT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — TENDÊNCIA DE CONFORMIDADE R4
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_header_band(sl, "Tendência de Conformidade — R4 Fornecedores",
                "Evolução diária ao longo do ciclo")

chart_data = ChartData()
chart_data.categories = ['10/07', '15/07', '17/07', '20/07', '22/07', '27/07']
chart_data.add_series('Conformidade R4 (%)', (41.8, 41.3, 40.9, 41.0, 41.8, 44.9))

chart = sl.shapes.add_chart(
    XL_CHART_TYPE.LINE_MARKERS,
    Inches(0.5), Inches(1.5),
    Inches(8.5), Inches(5.2),
    chart_data
).chart

chart.has_legend = False
chart.has_title = False
plot = chart.plots[0]
series = plot.series[0]
series.format.line.color.rgb = TEAL
series.format.line.width = Pt(2.5)
series.marker.format.fill.solid()
series.marker.format.fill.fore_color.rgb = TEAL
series.marker.size = 8

# Caixa de destaque lateral
add_rect(sl, Inches(9.3), Inches(1.5), Inches(3.7), Inches(5.2), CINZA_CLARO)
add_text(sl, "O que o gráfico mostra",
         Inches(9.45), Inches(1.6), Inches(3.4), Inches(0.4),
         font_size=12, bold=True, color=TEAL)
insights = [
    ("📉", "Queda nos 1ºs dias", "Conformidade recuou até 40,9% em 17/07 — fornecedores ainda processando as solicitações abertas no início do ciclo."),
    ("📈", "Virada na última semana", "De 22/07 a 27/07, +3,1 p.p. em 5 dias — concentração de entregas perto do fechamento."),
    ("✅", "Fechamento positivo", "44,9% — melhor índice do ciclo, com 147 novos documentos aprovados."),
]
y_ins = Inches(2.15)
for icon, title, body in insights:
    add_text(sl, f"{icon} {title}", Inches(9.45), y_ins, Inches(3.4), Inches(0.3),
             font_size=11, bold=True, color=PRETO_SUAVE)
    y_ins += Inches(0.32)
    add_text(sl, body, Inches(9.55), y_ins, Inches(3.3), Inches(0.65),
             font_size=10, color=CINZA_MEDIO, wrap=True)
    y_ins += Inches(0.78)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — R3: DOCUMENTOS DE TERCEIROS
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_header_band(sl, "R3 — Documentos de Terceiros",
                "Comparativo início × fechamento do ciclo")

rows_r3 = [
    ("Total de documentos", "8.843", "9.374", "+531"),
    ("Aprovados",           "1.533", "2.306", "+773"),
    ("Reprovados",          "4.454", "4.474", "+20"),
    ("Não Anexados",        "1.278", "1.317", "+39"),
    ("Conformidade",        "17,3%", "24,6%", "+7,3 p.p. ✓"),
]
col_widths_r3 = [Inches(4.2), Inches(2.5), Inches(2.8), Inches(2.8)]
tbl_r3 = sl.shapes.add_table(6, 4, Inches(0.35), Inches(1.55),
                              sum(col_widths_r3), Inches(2.6)).table
for i, w in enumerate(col_widths_r3):
    tbl_r3.columns[i].width = w
table_header_row(tbl_r3, ["Indicador", "10/07 (Início)", "27/07 (Fechamento)", "Variação"])
for ri, row_d in enumerate(rows_r3):
    highlight = VERDE_BG if ri == 4 else None
    table_data_row(tbl_r3, ri + 1, list(row_d), alt=(ri % 2 == 1), highlight_color=highlight)
    if ri == 4:
        for ci in range(4):
            p = tbl_r3.rows[ri + 1].cells[ci].text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x15, 0x57, 0x24)

# Terceiros cadastrados
add_rect(sl, Inches(7.7), Inches(1.55), Inches(5.3), Inches(2.6), CINZA_CLARO)
add_text(sl, "Terceiros Cadastrados",
         Inches(7.8), Inches(1.65), Inches(5.0), Inches(0.4),
         font_size=13, bold=True, color=TEAL)
add_text(sl, "1.267", Inches(7.8), Inches(2.1), Inches(5.0), Inches(0.85),
         font_size=48, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_text(sl, "terceiros ativos ao fechamento",
         Inches(7.8), Inches(2.95), Inches(5.0), Inches(0.35),
         font_size=12, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)
add_text(sl, "+16 novos no ciclo  (1.251 → 1.267)",
         Inches(7.8), Inches(3.35), Inches(5.0), Inches(0.4),
         font_size=12, bold=True, color=VERDE, align=PP_ALIGN.CENTER)

# observação + top docs R3
add_rect(sl, Inches(0.35), Inches(4.35), Inches(6.8), Inches(0.85), TEAL_CLARO)
add_text(sl, "💡  A conformidade R3 avançou +7,3 p.p. no ciclo — o maior salto entre os dois relatórios. "
             "O volume de aprovados cresceu 50% (1.533 → 2.306).",
         Inches(0.55), Inches(4.42), Inches(6.6), Inches(0.75),
         font_size=11, color=TEAL_ESCURO)

add_text(sl, "Top documentos não conformes (R3):",
         Inches(0.35), Inches(5.35), Inches(6.5), Inches(0.3),
         font_size=11, bold=True, color=TEAL)
docs_r3 = [
    ("Cartão Ponto (horas extras/noturnas)", "3.121"),
    ("Ficha de EPI", "1.388"),
    ("ASO", "1.237"),
    ("Capacitação conforme ordem de serviço", "1.056"),
    ("Ordens de Serviço", "1.037"),
]
for i, (doc, qtd) in enumerate(docs_r3):
    y_d = Inches(5.7) + i * Inches(0.28)
    add_text(sl, f"• {doc}", Inches(0.45), y_d, Inches(5.8), Inches(0.26),
             font_size=10, color=PRETO_SUAVE)
    add_text(sl, qtd, Inches(6.3), y_d, Inches(0.6), Inches(0.26),
             font_size=10, bold=True, color=VERMELHO, align=PP_ALIGN.RIGHT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PENDÊNCIAS — VISÃO DO CICLO
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_header_band(sl, "Pendências — Visão do Ciclo",
                "Movimentação de 10/07 a 27/07/2026")

# Card superior: pendências ativas (fila BPO)
add_rect(sl, Inches(0.35), Inches(1.55), Inches(4.5), Inches(1.45), VERDE_BG)
add_text(sl, "−83%", Inches(0.35), Inches(1.6), Inches(4.5), Inches(0.85),
         font_size=54, bold=True, color=VERDE, align=PP_ALIGN.CENTER)
add_text(sl, "Fila BPO (Em Elaboração)  |  842 → 138",
         Inches(0.35), Inches(2.45), Inches(4.5), Inches(0.4),
         font_size=11, bold=True, color=VERDE, align=PP_ALIGN.CENTER)

# Separador
add_rect(sl, Inches(0.35), Inches(3.07), Inches(4.5), Inches(0.05), CINZA_MEDIO)

# Card inferior: total não resolvidas (visão completa dashboard)
add_rect(sl, Inches(0.35), Inches(3.12), Inches(4.5), Inches(1.55), AMARELO_BG)
add_text(sl, "−10,5%", Inches(0.35), Inches(3.17), Inches(4.5), Inches(0.85),
         font_size=46, bold=True, color=LARANJA, align=PP_ALIGN.CENTER)
add_text(sl, "Não Resolvidas (dashboard)  |  3.013 → 2.696",
         Inches(0.35), Inches(4.0), Inches(4.5), Inches(0.4),
         font_size=11, bold=True, color=LARANJA, align=PP_ALIGN.CENTER)
add_text(sl, "Inclui pendências formalmente fechadas cujo\ndocumento ainda está irregular",
         Inches(0.35), Inches(4.42), Inches(4.5), Inches(0.35),
         font_size=9, italic=True, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

# tabela de pendências
cols_p = ["Indicador", "10/07", "27/07", "Variação"]
rows_p = [
    ("Total geral",               "5.812", "6.417", "+605"),
    ("Fila BPO (Em elaboração)",   "842",   "138",   "−704 (−83%) ✅"),
    ("Não Resolvidas (dashboard)", "3.013", "2.696", "−317 (−10,5%)"),
    ("Área Fornecedores",         "705",   "847",   "+142"),
    ("Área Terceiros",            "5.087", "5.550", "+463"),
]
col_w_p = [Inches(3.8), Inches(1.6), Inches(1.8), Inches(2.4)]
tbl_p = sl.shapes.add_table(6, 4, Inches(5.1), Inches(1.55),
                             sum(col_w_p), Inches(3.2)).table
for i, w in enumerate(col_w_p):
    tbl_p.columns[i].width = w
table_header_row(tbl_p, cols_p)
for ri, rd in enumerate(rows_p):
    highlight = VERDE_BG if ri == 1 else (AMARELO_BG if ri == 2 else None)
    table_data_row(tbl_p, ri + 1, list(rd), alt=(ri % 2 == 1), highlight_color=highlight)
    if ri in (1, 2):
        for ci in range(4):
            p = tbl_p.rows[ri + 1].cells[ci].text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x15, 0x57, 0x24) if ri == 1 else RGBColor(0x92, 0x60, 0x00)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PANORAMA DE REJEIÇÕES
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_header_band(sl, "Panorama de Rejeições", "Ciclo 10/07 a 27/07/2026")

col_h  = Inches(5.5)
col_y  = Inches(1.45)
col_w  = Inches(4.0)
gap    = Inches(0.315)
cols_x = [Inches(0.35), Inches(0.35) + col_w + gap, Inches(0.35) + 2 * (col_w + gap)]

for cx in cols_x:
    add_rect(sl, cx, col_y, col_w, col_h, CINZA_CLARO)

# ── Coluna 1: Documentos Mais Críticos ────────────────────────────────────
cx = cols_x[0]
add_text(sl, "Documentos Mais Críticos",
         cx + Inches(0.12), col_y + Inches(0.1), col_w - Inches(0.24), Inches(0.38),
         font_size=12, bold=True, color=TEAL)
docs_crit = [
    ("Cartão Ponto (HE/Noturno)",
     "Maior frequência de reprovação ao longo de várias competências."),
    ("Ordens de Serviço",
     "Volume expressivo de rejeições em múltiplos prestadores."),
    ("Ficha de EPI",
     "Alto índice de não conformidade mensal e geral."),
    ("ASO e Capacitações",
     "Volume considerável de pendências e recusas."),
]
y_d = col_y + Inches(0.58)
for titulo, desc in docs_crit:
    add_text(sl, f"• {titulo}", cx + Inches(0.15), y_d, col_w - Inches(0.28), Inches(0.26),
             font_size=11, bold=True, color=PRETO_SUAVE)
    y_d += Inches(0.28)
    add_text(sl, desc, cx + Inches(0.25), y_d, col_w - Inches(0.38), Inches(0.48),
             font_size=10, color=CINZA_MEDIO, wrap=True)
    y_d += Inches(0.57)

# ── Coluna 2: Competências Críticas ───────────────────────────────────────
cx = cols_x[1]
add_text(sl, "Competências Críticas",
         cx + Inches(0.12), col_y + Inches(0.1), col_w - Inches(0.24), Inches(0.38),
         font_size=12, bold=True, color=TEAL)
comps = [
    ("🔴  12/25 e 05/26", VERMELHO, RGBColor(0xFF, 0xE5, 0xE5),
     "Maior taxa de reprovação direta — foco em Cartão Ponto, EPI e Capacitação."),
    ("⚠️  02/26 e 03/26", LARANJA, AMARELO_BG,
     "Acúmulo de documentos com Aguardando Submissão e reprovações pontuais."),
]
y_c = col_y + Inches(0.58)
for label, cor, bg, desc in comps:
    add_rect(sl, cx + Inches(0.12), y_c, col_w - Inches(0.24), Inches(0.38), bg)
    add_text(sl, label, cx + Inches(0.18), y_c + Inches(0.04), col_w - Inches(0.36), Inches(0.3),
             font_size=12, bold=True, color=cor)
    y_c += Inches(0.45)
    add_text(sl, desc, cx + Inches(0.18), y_c, col_w - Inches(0.36), Inches(0.65),
             font_size=10, color=CINZA_MEDIO, wrap=True)
    y_c += Inches(1.1)

# ── Coluna 3: Fornecedores com Maior Impacto ──────────────────────────────
cx = cols_x[2]
add_text(sl, "Fornecedores com Maior Impacto",
         cx + Inches(0.12), col_y + Inches(0.1), col_w - Inches(0.24), Inches(0.38),
         font_size=12, bold=True, color=TEAL)
forn_imp = [
    ("Fokus Consultoria",
     "Maior volume acumulado — ASO, OS, EPI e Capacitação reprovados em quase todos os terceiros."),
    ("Centro",
     "Índice próximo a 100% de reprovação na grade de colaboradores."),
    ("Ecotech",
     "Reprovações sequenciais em cartões ponto de 12/25 a 05/26."),
    ("Calhas Reis",
     "100% de reprovação em todos os colaboradores e documentos."),
]
y_f = col_y + Inches(0.58)
for forn, desc in forn_imp:
    add_text(sl, f"• {forn}", cx + Inches(0.15), y_f, col_w - Inches(0.28), Inches(0.26),
             font_size=11, bold=True, color=PRETO_SUAVE)
    y_f += Inches(0.28)
    add_text(sl, desc, cx + Inches(0.25), y_f, col_w - Inches(0.38), Inches(0.5),
             font_size=10, color=CINZA_MEDIO, wrap=True)
    y_f += Inches(0.62)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CASOS ESPECÍFICOS
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_header_band(sl, "Pauta — Casos Específicos", "Acompanhamento por fornecedor  |  27/07/2026")

casos = [
    ("Leve Estacionamento", VERMELHO,
     "Documento ilegível. Manter reprovado até que seja enviado o documento original "
     "em alta resolução e legível."),
    ("Teletex", LARANJA,
     "Verificação de CCT pela Atividade Preponderante do Colaborador vs. CNAE Principal da Empresa. "
     "Pela CLT Art. 511 §2º, o enquadramento sindical é determinado pela atividade econômica "
     "preponderante do empregador (salvo categorias diferenciadas)."),
    ("Karuana", VERMELHO,
     "Ponto manual com intervalo de almoço reduzido a 30 min sem respaldo de Acordo Coletivo/CCT. "
     "Recusa no envio de documentos sob alegação de contrato com a Zurich."),
    ("Orbenk", LARANJA,
     "Reunião marcada sem envio prévio do escopo ou solicitação. "
     "Possibilidade de datas enviada sem retorno até o momento."),
    ("Vanderlande", CINZA_MEDIO,
     "Reunião em definição de data, entre quinta-feira e sexta-feira."),
]

y_caso = Inches(1.52)
for i, (forn, cor, desc) in enumerate(casos):
    if i > 0:
        add_rect(sl, Inches(0.35), y_caso - Inches(0.06),
                 Inches(12.63), Inches(0.01), RGBColor(0xDD, 0xDD, 0xDD))
    add_rect(sl, Inches(0.35), y_caso, Inches(0.15), Inches(0.9), cor)
    add_text(sl, forn,
             Inches(0.62), y_caso, Inches(12.36), Inches(0.32),
             font_size=13, bold=True, color=PRETO_SUAVE)
    add_text(sl, desc,
             Inches(0.62), y_caso + Inches(0.33), Inches(12.36), Inches(0.58),
             font_size=10, color=CINZA_MEDIO, wrap=True)
    y_caso += Inches(1.05)

# ═════════════════════════════════════════════════════════════════════════════
# SALVAR
# ═════════════════════════════════════════════════════════════════════════════
output = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\clientes\Zurich\apresentacoes\Zurich_Ciclo_10_27jul2026_27-07-2026_v6.pptx"
prs.save(output)
print(f"✓ Arquivo salvo em: {output}")
print(f"  Slides gerados: {len(prs.slides)}")
