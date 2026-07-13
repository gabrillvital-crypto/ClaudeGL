#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path

# Paleta Efcaz
TEAL = RGBColor(0x0E, 0x8F, 0xA3)
TEAL_CLARO = RGBColor(0x14, 0xB3, 0xCC)
VERDE = RGBColor(0x27, 0xAE, 0x60)
VERMELHO = RGBColor(0xE7, 0x4C, 0x3C)
CINZA_CLARO = RGBColor(0xF2, 0xF4, 0xF4)
PRETO = RGBColor(0x17, 0x20, 0x2A)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)

# Criar apresentação
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_slide_blank():
    """Adiciona slide em branco"""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(blank_layout)

def add_title_shape(slide, text, top, left, width, height, size=54, color=TEAL, bold=True):
    """Adiciona título formatado"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Calibri'

    return shape

def add_subtitle_shape(slide, text, top, left, width, height, size=24, color=PRETO):
    """Adiciona subtítulo"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = shape.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = 'Calibri'

    return shape

def add_metric_box(slide, title, value, top, left, width=2.2, height=1.8):
    """Adiciona box de métrica"""
    # Background
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CINZA_CLARO
    shape.line.color.rgb = TEAL
    shape.line.width = Pt(2)

    # Título da métrica
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRETO
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

    # Valor
    p = text_frame.add_paragraph()
    p.text = value
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(10)

def add_bullet_point(slide, text, top, left, width=9, size=18, color=PRETO):
    """Adiciona ponto de lista"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.4))
    text_frame = shape.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "• " + text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = 'Calibri'

print("Criando PPTX de NPS do zero...\n")

# ====== SLIDE 1: CAPA ======
print("Slide 1: Capa")
slide1 = add_slide_blank()
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BRANCO

add_title_shape(slide1, "NPS Carteira Efcaz", 2, 0.5, 9, 1.5, size=60, color=TEAL)
add_subtitle_shape(slide1, "Análise de Engajamento e Satisfação", 3.7, 0.5, 9, 0.6, size=28, color=PRETO)
add_subtitle_shape(slide1, "Junho 2026 | Gabriel Vital • Customer Success Specialist", 5.5, 0.5, 9, 0.5, size=16, color=CINZA_CLARO)

# ====== SLIDE 2: RESUMO EXECUTIVO ======
print("Slide 2: Resumo Executivo")
slide2 = add_slide_blank()
background = slide2.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BRANCO

add_title_shape(slide2, "Resumo Executivo", 0.5, 0.5, 9, 0.7, size=44)

# 3 Boxes de métricas
add_metric_box(slide2, "Respondidos", "15\n(60%)", 1.5, 0.8)
add_metric_box(slide2, "Sem Resposta", "8\n(32%)", 1.5, 3.3)
add_metric_box(slide2, "Sem Acesso", "2\n(8%)", 1.5, 5.8)

# Contexto
add_subtitle_shape(slide2, "Total de 25 clientes avaliados | NPS Médio: +73 | Taxa de Resposta: 22.6%", 4, 0.8, 8.4, 0.5, size=14, color=PRETO)

# ====== SLIDE 3: DISTRIBUIÇÃO POR STATUS ======
print("Slide 3: Distribuição por Status")
slide3 = add_slide_blank()
background = slide3.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BRANCO

add_title_shape(slide3, "Distribuição por Status — Visão da Carteira", 0.5, 0.5, 9, 0.7, size=44)

add_bullet_point(slide3, "15 clientes (60%) — Com resposta: Promotores mapeados com NPS +73", 1.5, 0.8, size=18)
add_bullet_point(slide3, "8 clientes (32%) — Acessam mas não responderam: Heavy users sem engajamento", 2.2, 0.8, size=18)
add_bullet_point(slide3, "2 clientes (8%) — Sem nenhum acesso: Verificar onboarding/credenciais", 2.9, 0.8, size=18)

add_subtitle_shape(slide3, "Insight Principal", 3.8, 0.8, 8.4, 0.4, size=16, color=PRETO)
add_bullet_point(slide3, "Problema não é satisfação — é engajamento com pesquisa", 4.3, 0.8, size=18, color=TEAL)

add_subtitle_shape(slide3, "Métricas de Participação", 5.1, 0.8, 8.4, 0.4, size=16, color=PRETO)
add_bullet_point(slide3, "Taxa Geral: 22.6% (37 de 164 usuários responderam)", 5.6, 0.8, size=18)

# ====== SLIDE 4: QUALIDADE DOS RESPONDENTES ======
print("Slide 4: Qualidade dos Respondentes")
slide4 = add_slide_blank()
background = slide4.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BRANCO

add_title_shape(slide4, "Qualidade dos Respondentes", 0.5, 0.5, 9, 0.7, size=44)

add_metric_box(slide4, "NPS Médio", "+73", 1.5, 0.8)
add_metric_box(slide4, "Detratores", "0", 1.5, 3.3)
add_metric_box(slide4, "Promoters", "100%", 1.5, 5.8)

add_subtitle_shape(slide4, "Conclusão: Quando respondem, respondentes são 100% promoters (notas 9-10).", 4.3, 0.8, 8.4, 0.6, size=16, color=PRETO)

# ====== SLIDE 5: CRÍTICOS POR TIER ======
print("Slide 5: Críticos por Tier")
slide5 = add_slide_blank()
background = slide5.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BRANCO

add_title_shape(slide5, "Análise por Tier — Distribuição de Problemas", 0.5, 0.5, 9, 0.7, size=44)

# Table-like structure
add_subtitle_shape(slide5, "Tier A", 1.4, 0.8, 2, 0.3, size=14, color=PRETO)
add_bullet_point(slide5, "2 sem resposta (transferido a Wagner)", 1.8, 1.2, width=8, size=14)

add_subtitle_shape(slide5, "Tier B+", 2.5, 0.8, 2, 0.3, size=14, color=PRETO)
add_bullet_point(slide5, "1 sem resposta: Bunker One (desengajado apesar de Ongoing)", 2.9, 1.2, width=8, size=14)

add_subtitle_shape(slide5, "Tier B", 3.9, 0.8, 2, 0.3, size=14, color=VERMELHO)
add_bullet_point(slide5, "4 críticos: 1 sem acesso + 3 sem resposta", 4.3, 1.2, width=8, size=14, color=VERMELHO)
add_bullet_point(slide5, "Destaque: Integral Médica (renovação 30/06), Federação Paulista (Filipe Silva 10 views)", 4.85, 1.2, width=8, size=14)

add_subtitle_shape(slide5, "Tier C", 5.8, 0.8, 2, 0.3, size=14, color=PRETO)
add_bullet_point(slide5, "3 críticos: 1 sem acesso + 2 sem resposta", 6.2, 1.2, width=8, size=14)

# ====== SLIDE 6: RESUMO FINAL ======
print("Slide 6: Resumo Final")
slide6 = add_slide_blank()
background = slide6.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BRANCO

add_title_shape(slide6, "Resumo — Situação Geral", 0.5, 0.5, 9, 0.7, size=44)

add_subtitle_shape(slide6, "✓ Positivo", 1.5, 0.8, 9, 0.3, size=18, color=VERDE)
add_bullet_point(slide6, "NPS médio de +73 indica altíssima satisfação entre respondentes", 1.95, 1.2, width=8.2, size=14)
add_bullet_point(slide6, "Zero detratores genuínos — problema não é produto, é participação", 2.45, 1.2, width=8.2, size=14)

add_subtitle_shape(slide6, "⚠ Crítico", 3.2, 0.8, 9, 0.3, size=18, color=VERMELHO)
add_bullet_point(slide6, "10 clientes com problema: 2 sem acesso + 8 acessando mas não respondendo", 3.65, 1.2, width=8.2, size=14)
add_bullet_point(slide6, "Tier B concentra maioria dos problemas (4 de 12 clientes)", 4.15, 1.2, width=8.2, size=14)
add_bullet_point(slide6, "Heavy users consistentemente não respondendo — padrão de desengajamento", 4.65, 1.2, width=8.2, size=14)

add_subtitle_shape(slide6, "→ Próxima Onda de NPS", 5.5, 0.8, 9, 0.3, size=18, color=TEAL)
add_bullet_point(slide6, "Simplificar pesquisa (2-3 perguntas vs. atual) para aumentar taxa de resposta de 22.6%", 5.95, 1.2, width=8.2, size=14)

# ====== SALVAR ======
output_dir = Path(r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\clientes\NPS")
output_dir.mkdir(parents=True, exist_ok=True)
output_path = output_dir / "NPS_Carteira_Efcaz_2026_06_19.pptx"

prs.save(str(output_path))
print(f"\n✅ PPTX criado do zero e salvo em: {output_path}")
print(f"   6 slides com dados puros de NPS")
