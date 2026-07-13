#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path

LOGO = r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\Geral\logo_efcaz_clean.png"

# Paleta Efcaz
TEAL        = RGBColor(0x0E, 0x8F, 0xA3)
TEAL_ESC    = RGBColor(0x0A, 0x6A, 0x7A)
TEAL_CLARO  = RGBColor(0xE6, 0xF7, 0xF9)   # fundo suave
VERDE       = RGBColor(0x27, 0xAE, 0x60)
VERMELHO    = RGBColor(0xE7, 0x4C, 0x3C)
LARANJA     = RGBColor(0xF3, 0x9C, 0x12)
CINZA_CLARO = RGBColor(0xF2, 0xF4, 0xF4)
CINZA_MED   = RGBColor(0x85, 0x92, 0x9E)
PRETO       = RGBColor(0x17, 0x20, 0x2A)
BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, left, top, width, height, fill_color, line_color=None, line_pt=0):
    sh = slide.shapes.add_shape(1,
         Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill_color
    if line_color:
        sh.line.color.rgb = line_color; sh.line.width = Pt(line_pt)
    else:
        sh.line.fill.background()
    return sh

def txt(slide, text, top, left, width, height,
        size=14, color=PRETO, bold=False, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    sh = slide.shapes.add_textbox(
         Inches(left), Inches(top), Inches(width), Inches(height))
    tf = sh.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.bold = bold
    p.font.italic = italic; p.font.color.rgb = color
    p.font.name = "Calibri"; p.alignment = align
    return sh

def bullet(slide, text, top, left=1.0, width=11.8, size=13, color=PRETO, bold=False):
    return txt(slide, "• " + text, top, left, width, 0.45, size=size, color=color, bold=bold)

def logo(slide, left=11.8, top=0.1, width=1.3):
    slide.shapes.add_picture(LOGO, Inches(left), Inches(top), width=Inches(width))

def header(slide, titulo_txt, size=22):
    """Barra teal no topo com título e logo."""
    rect(slide, 0, 0, 13.33, 1.0, TEAL)
    txt(slide, titulo_txt, 0.18, 0.4, 10.5, 0.7,
        size=size, color=BRANCO, bold=True, align=PP_ALIGN.LEFT)
    # logo sobre fundo branco pequeno para contraste
    rect(slide, 11.7, 0.08, 1.45, 0.84, BRANCO)
    logo(slide, left=11.75, top=0.15, width=1.35)

def rodape(slide):
    rect(slide, 0, 7.2, 13.33, 0.3, TEAL)
    txt(slide, "Gabriel Vital  •  Customer Success Specialist  •  Efcaz SRM  •  Julho/2026",
        7.22, 0.4, 12.5, 0.26, size=10, color=BRANCO, align=PP_ALIGN.CENTER)

def divisor(slide, top, cor=TEAL, alpha=False):
    sh = rect(slide, 0.4, top, 12.5, 0.03, CINZA_MED if alpha else cor)

def box_kpi(slide, label, valor, top, left, width=2.8, height=1.85,
            cor_borda=TEAL, cor_valor=TEAL):
    r = rect(slide, left, top, width, height, CINZA_CLARO, cor_borda, 1.5)
    sh = slide.shapes.add_textbox(
         Inches(left+0.1), Inches(top+0.12), Inches(width-0.2), Inches(height-0.24))
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = label; p.font.size = Pt(12); p.font.bold = True
    p.font.color.rgb = PRETO; p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = valor; p2.font.size = Pt(28); p2.font.bold = True
    p2.font.color.rgb = cor_valor; p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(6)

print("Criando PPTX NPS Julho/2026 — Identidade Visual Efcaz...\n")

# ═══════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ═══════════════════════════════════════════════════════════
print("Slide 1: Capa")
s1 = blank()
s1.background.fill.solid(); s1.background.fill.fore_color.rgb = BRANCO

# Faixa teal superior cobrindo 60% do slide
rect(s1, 0, 0, 13.33, 4.6, TEAL)

# Faixa teal escura no fundo
rect(s1, 0, 7.2, 13.33, 0.3, TEAL_ESC)

# Logo sobre box branco no canto sup. direito
rect(s1, 11.5, 0.1, 1.65, 0.9, BRANCO)
logo(s1, left=11.55, top=0.18, width=1.55)

# Símbolo decorativo (círculo semi-transparente)
sh_circ = s1.shapes.add_shape(9,  # oval
    Inches(9.5), Inches(-1.5), Inches(6), Inches(6))
sh_circ.fill.solid(); sh_circ.fill.fore_color.rgb = TEAL_ESC
sh_circ.line.fill.background()

# Textos na área teal
txt(s1, "NPS", 1.1, 0.6, 12.0, 1.2, size=72, color=BRANCO, bold=True)
txt(s1, "Carteira Efcaz SRM", 2.35, 0.6, 10.0, 0.9, size=34, color=BRANCO, bold=False)
txt(s1, "Análise de Engajamento e Satisfação", 3.2, 0.6, 10.0, 0.6, size=18, color=RGBColor(0xB2,0xEB,0xF2))

# Área branca inferior
txt(s1, "Base: 07/07/2026  |  Comparativo: 19/06 → 07/07/2026",
    5.0, 0.6, 9.0, 0.5, size=15, color=TEAL, bold=True)
txt(s1, "Gabriel Vital  •  Customer Success Specialist  •  Efcaz SRM",
    5.65, 0.6, 9.0, 0.45, size=13, color=CINZA_MED)

# Linha decorativa teal entre as áreas
rect(s1, 0, 4.55, 13.33, 0.05, TEAL_ESC)

# ═══════════════════════════════════════════════════════════
# SLIDE 2 — RESUMO EXECUTIVO
# ═══════════════════════════════════════════════════════════
print("Slide 2: Resumo Executivo")
s2 = blank()
s2.background.fill.solid(); s2.background.fill.fore_color.rgb = BRANCO
header(s2, "Resumo Executivo — Visão da Carteira (07/07/2026)")

# KPI row 1
box_kpi(s2, "Com Resposta",    "17  (65%)",  1.15, 0.4,  width=2.7, cor_valor=VERDE,   cor_borda=VERDE)
box_kpi(s2, "Sem Resposta",    "8  (31%)",   1.15, 3.25, width=2.7, cor_valor=LARANJA,  cor_borda=LARANJA)
box_kpi(s2, "Sem Acesso",      "1  (4%)",    1.15, 6.1,  width=2.7, cor_valor=VERMELHO, cor_borda=VERMELHO)
box_kpi(s2, "NPS Geral",       "+77",        1.15, 8.95, width=2.7, cor_valor=TEAL,     cor_borda=TEAL)
box_kpi(s2, "Clientes Avaliados", "26",      1.15, 11.8, width=1.3, height=1.85, cor_valor=PRETO, cor_borda=CINZA_MED)

# KPI row 2
box_kpi(s2, "Total Exibições",  "163", 3.2, 0.4,  width=2.7, height=1.55, cor_valor=PRETO, cor_borda=CINZA_MED)
box_kpi(s2, "Respondentes",     "47",  3.2, 3.25, width=2.7, height=1.55, cor_valor=PRETO, cor_borda=CINZA_MED)
box_kpi(s2, "Taxa de Resposta", "28,8%", 3.2, 6.1, width=2.7, height=1.55, cor_valor=TEAL, cor_borda=TEAL)
box_kpi(s2, "Promoters",  "38",  3.2, 8.95, width=1.3, height=1.55, cor_valor=VERDE,   cor_borda=VERDE)
box_kpi(s2, "Passivos",   "7",   3.2, 10.4, width=1.3, height=1.55, cor_valor=LARANJA, cor_borda=LARANJA)
box_kpi(s2, "Detratores", "2",   3.2, 11.85,width=1.25,height=1.55, cor_valor=VERMELHO,cor_borda=VERMELHO)

# Comparativo
rect(s2, 0.4, 4.95, 12.5, 0.04, TEAL_ESC)
rect(s2, 0.4, 5.1, 12.5, 0.9, TEAL_CLARO)
txt(s2, "Comparativo  19/06 → 07/07", 5.12, 0.55, 3.5, 0.35, size=11, color=TEAL, bold=True)
txt(s2, (
    "Clientes: 25→26 (+1)   |   Com resposta: 15→17 (+2)   |   Exibições: 138→163 (+25)   |   "
    "Respondentes: 34→47 (+13)   |   Taxa: 24,6%→28,8% (+4,2pp)   |   NPS: +79→+77 (-2)"
), 5.47, 0.55, 12.3, 0.45, size=11, color=PRETO)
rodape(s2)

# ═══════════════════════════════════════════════════════════
# SLIDE 3 — SEGMENTAÇÃO NPS
# ═══════════════════════════════════════════════════════════
print("Slide 3: Segmentação NPS Completa")
s3 = blank()
s3.background.fill.solid(); s3.background.fill.fore_color.rgb = BRANCO
header(s3, "Segmentação NPS — Respondentes por Categoria")

# Promoters
rect(s3, 0.4, 1.1, 12.5, 0.04, VERDE)
txt(s3, "PROMOTERS (nota 9–10) — 38 respondentes", 1.2, 0.5, 12.0, 0.4, size=13, color=VERDE, bold=True)
bullet(s3, "Clientes 100% promoters: Cielo, Engesp (2), Geistlich (2), Hospital Adventista (4), Paccoby, Ponsse, Tarkett, Unimed Brasil, Unimed Campo Grande (3), Unimed Sou", 1.7)
bullet(s3, "Promoters com passivos/detratores: Dock Brasil 5P | Integral Médica 3P | Unimed Dourados 3P | DOF 2P | Afonso França 5P | Cebrace 3P", 2.15)
bullet(s3, "Destaques: 'SUPER EFICIENTE' (Dock Brasil)  •  'Melhor controle de acesso' (Dock Brasil)  •  'Agilidade' (Afonso França)  •  'Excelente' (Tarkett)", 2.6, color=VERDE, size=12)

# Passivos
rect(s3, 0.4, 3.1, 12.5, 0.04, LARANJA)
txt(s3, "PASSIVOS (nota 7–8) — 7 respondentes", 3.2, 0.5, 12.0, 0.4, size=13, color=LARANJA, bold=True)
bullet(s3, "Bom Futuro — Cairo Freitas (7): 'Ainda preciso de suporte, mas tem evoluído bastante'", 3.7)
bullet(s3, "Cebrace — Leandro Alves (7): 'Filtro não atende a necessidade'  |  Afonso França — Cibele Grechi (8)", 4.15)
bullet(s3, "Dock Brasil — Arthur Sumihara (8): 'Boa plataforma'  |  DOF — Rosilene Rangel (8): 'Nada a declarar'", 4.6)
bullet(s3, "Integral Médica — Ida Laiber (7): 'A conexão cai e tenho que reabrir'  |  Unimed Dourados — Nilton (8): 'Relatórios específicos precisam melhorar'", 5.05)

# Detratores
rect(s3, 0.4, 5.55, 12.5, 0.04, VERMELHO)
txt(s3, "DETRATORES (nota 0–6) — 2 respondentes", 5.65, 0.5, 12.0, 0.4, size=13, color=VERMELHO, bold=True)
bullet(s3, "Afonso França — Thais Barroso (nota 6): 'Navegabilidade, suporte e praticidade operacional'", 6.1, color=VERMELHO)
bullet(s3, "Cebrace — Bruno Machado (nota 1): 'Não faz sentido avaliar o mesmo fornecedor em período tão curto'", 6.55, color=VERMELHO)
rodape(s3)

# ═══════════════════════════════════════════════════════════
# SLIDE 4 — DISTRIBUIÇÃO POR STATUS
# ═══════════════════════════════════════════════════════════
print("Slide 4: Distribuição por Status")
s4 = blank()
s4.background.fill.solid(); s4.background.fill.fore_color.rgb = BRANCO
header(s4, "Distribuição por Status — 26 Clientes Avaliados")

rect(s4, 0.4, 1.1, 12.5, 0.04, VERDE)
txt(s4, "COM RESPOSTA — 17 clientes (65%)", 1.2, 0.5, 12.0, 0.38, size=13, color=VERDE, bold=True)
bullet(s4, "10 clientes NPS +100: Cielo, Engesp, Geistlich, Hospital Adventista, Paccoby, Ponsse, Tarkett, Unimed Brasil, Unimed Campo Grande, Unimed Sou", 1.65)
bullet(s4, "7 clientes entre +40 e +83: Dock Brasil +83 | Integral Médica +75 | Unimed Dourados +75 | DOF +67 | Afonso França +57 | Cebrace +40 | Bom Futuro +0", 2.1)

rect(s4, 0.4, 2.6, 12.5, 0.04, LARANJA)
txt(s4, "ACESSAM MAS NÃO RESPONDERAM — 8 clientes (31%)", 2.7, 0.5, 12.0, 0.38, size=13, color=LARANJA, bold=True)
bullet(s4, "Bunker One (B+): 5 exib — Rodrigo Moura (10v), Antonio Mendes (5v), Roney Gatto (5v)", 3.15, color=PRETO)
bullet(s4, "Federação Paulista (B): 5 exib — Filipe Marques Silva (11v), Marcelo Campos (9v), Fernanda Zanzarini (5v)", 3.6, color=PRETO)
bullet(s4, "ZAB / Zurich (A): 2 exib — SESMT (6v), Débora Coelho (6v)  |  Eucatex (A): 4 exib  |  Alumetaf (B): 2 exib", 4.05, color=PRETO)
bullet(s4, "Amboretto (C): Andressa Gabateli (15v!), Juliana Agostinho (8v)  |  Asso Marítima (C): Ivan Barbosa (10v)  |  Banco Honda: 2 exib", 4.5, color=PRETO)

rect(s4, 0.4, 5.0, 12.5, 0.04, VERMELHO)
txt(s4, "SEM ACESSO — 1 cliente (4%)", 5.1, 0.5, 12.0, 0.38, size=13, color=VERMELHO, bold=True)
bullet(s4, "ADVtec (Tier C) — nenhum usuário acessou o NPS na campanha atual", 5.55, color=VERMELHO)

rect(s4, 0.4, 6.1, 12.5, 0.04, CINZA_MED)
txt(s4, "Taxa de Resposta: 28,8%  (47 de 163 exibições)  |  +4,2pp vs. 19/06/2026",
    6.2, 0.5, 12.0, 0.4, size=12, color=CINZA_MED)
rodape(s4)

# ═══════════════════════════════════════════════════════════
# SLIDE 5 — ANÁLISE POR TIER
# ═══════════════════════════════════════════════════════════
print("Slide 5: Análise por Tier")
s5 = blank()
s5.background.fill.solid(); s5.background.fill.fore_color.rgb = BRANCO
header(s5, "Análise por Tier — Pontos de Atenção")

def bloco_tier(slide, label, top, itens):
    rect(slide, 0.4, top, 12.5, 0.04, TEAL)
    txt(slide, label, top + 0.1, 0.5, 3.0, 0.38, size=13, color=TEAL, bold=True)
    for i, (item, cor) in enumerate(itens):
        bullet(slide, item, top + 0.55 + i * 0.42, color=cor, size=12)

bloco_tier(s5, "Tier A", 1.1, [
    ("ZAB / Zurich (R$8.9k): sem resposta — SESMT e Débora acessam mas não respondem. Prioridade máxima.", VERMELHO),
    ("Bom Futuro (R$6.5k): nova resposta NPS +0 — Cairo passivo (nota 7): 'Ainda preciso de suporte'", LARANJA),
    ("Eucatex (R$5.5k): sem resposta — 4 exibições sem engajamento com pesquisa", LARANJA),
])

bloco_tier(s5, "Tier B+", 3.2, [
    ("Bunker One (R$3.8k): sem resposta — Rodrigo Moura (10v) nunca respondeu ao NPS", LARANJA),
    ("Dock Brasil (R$3.6k): piora -17pp (100→83) — novo passivo Arthur Sumihara (nota 8)", LARANJA),
    ("Cebrace (R$2.3k): NPS +40 — detrator Bruno Machado (nota 1) + passivo Leandro Alves (nota 7)", VERMELHO),
])

bloco_tier(s5, "Tier B", 5.25, [
    ("Federação Paulista (R$2.5k): sem resposta — Filipe Silva (11v) e Marcelo Campos (9v) não engajam", LARANJA),
    ("DOF (R$2.7k): piora -33pp (100→67) — novo passivo Rosilene (8)  |  Integral Médica: passiva Ida Laiber (nota 7)", LARANJA),
])

rodape(s5)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 — RESUMO E PRÓXIMOS PASSOS
# ═══════════════════════════════════════════════════════════
print("Slide 6: Resumo e Próximos Passos")
s6 = blank()
s6.background.fill.solid(); s6.background.fill.fore_color.rgb = BRANCO
header(s6, "Resumo e Próximos Passos")

rect(s6, 0.4, 1.1, 12.5, 0.04, VERDE)
txt(s6, "POSITIVO", 1.2, 0.5, 12.0, 0.38, size=13, color=VERDE, bold=True)
bullet(s6, "NPS Geral +77 com 38 promoters — taxa de resposta subiu 24,6% → 28,8% (+4,2pp | +13 respondentes)", 1.65, color=VERDE)
bullet(s6, "Melhora em Afonso França (+7pp), Cebrace (+15pp) e Integral Médica (+8pp)", 2.1, color=VERDE)

rect(s6, 0.4, 2.6, 12.5, 0.04, VERMELHO)
txt(s6, "ATENÇÃO", 2.7, 0.5, 12.0, 0.38, size=13, color=VERMELHO, bold=True)
bullet(s6, "ZAB (Tier A, R$8.9k): sem resposta — pesquisa não está sendo respondida por nenhum usuário real", 3.15, color=VERMELHO)
bullet(s6, "Detratores persistentes: Thais Barroso (Afonso França, nota 6) e Bruno Machado (Cebrace, nota 1) — mesmos da onda anterior", 3.6, color=VERMELHO)
bullet(s6, "8 heavy users com muitos acessos sem responder: Amboretto (15v), Ivan Barbosa (10v), Rodrigo Moura (10v), Filipe Silva (11v)", 4.05, color=LARANJA)

rect(s6, 0.4, 4.55, 12.5, 0.04, TEAL)
txt(s6, "PRÓXIMOS PASSOS", 4.65, 0.5, 12.0, 0.38, size=13, color=TEAL, bold=True)
bullet(s6, "Acionar ZAB: entender por que pesquisa não é respondida pelos usuários reais da conta", 5.1, color=TEAL)
bullet(s6, "Follow-up com detratores: contato direto com Thais Barroso (Afonso França) e Bruno Machado (Cebrace)", 5.55, color=TEAL)
bullet(s6, "Repassar ao produto: problema de conexão (Integral Médica), filtros (Cebrace), extração de relatórios (Unimed Dourados)", 6.0, color=TEAL)
rodape(s6)

# ═══════════════════════════════════════════════════════════
# SALVAR
# ═══════════════════════════════════════════════════════════
output = Path(r"C:\Users\gabriel.evangelista\Documents\ClaudeGL\clientes\NPS\NPS_Carteira_Efcaz_2026_07_07.pptx")
prs.save(str(output))
print(f"\nSalvo: {output}")
